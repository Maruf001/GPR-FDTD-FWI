"""
Generate a scientific PowerPoint presentation for the GPR-FDTD-FWI project.

Creates ~23 slides covering: FDTD theory, forward simulation results,
inversion methodology and results, GPU acceleration, extensions, and
conclusions. All figures are generated or loaded from existing outputs.

Usage:
    python generate_presentation.py

Output:
    outputs/GPR_FDTD_FWI_Presentation.pptx
"""
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import config as cfg

# ═════════════════════════════════════════════════════════════════════════════
# Design constants
# ═════════════════════════════════════════════════════════════════════════════
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TITLE_COLOR = RGBColor(0x1B, 0x3A, 0x5C)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT_COLOR = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_TITLE = 'Calibri'
FONT_BODY = 'Calibri'

FIG_DIR = 'outputs/figures'
PRES_FIG_DIR = 'outputs/presentation_figures'
os.makedirs(FIG_DIR, exist_ok=True)
os.makedirs(PRES_FIG_DIR, exist_ok=True)


# ═════════════════════════════════════════════════════════════════════════════
# Helper functions
# ═════════════════════════════════════════════════════════════════════════════

def new_slide(prs):
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a dark blue title bar at the top of the slide."""
    # Title background bar
    left, top, width, height = Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.font.name = FONT_BODY


def add_text_box(slide, left, top, width, height, text, size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or BODY_COLOR
    p.font.name = FONT_BODY
    p.alignment = alignment
    return tf


def add_bullets(slide, left, top, width, height, items, size=16):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = BODY_COLOR
        p.font.name = FONT_BODY
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_figure(slide, fig_path, left, top, width=None, height=None):
    """Add a figure image to the slide."""
    if not os.path.exists(fig_path):
        print(f"  WARNING: {fig_path} not found, skipping")
        return
    if width and height:
        slide.shapes.add_picture(fig_path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(fig_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(fig_path, left, top, height=height)
    else:
        slide.shapes.add_picture(fig_path, left, top)


def set_notes(slide, notes_text):
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_key_number(slide, left, top, number_text, label_text):
    """Add a highlighted key number with label."""
    txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label_text
    p2.font.size = Pt(14)
    p2.font.color.rgb = BODY_COLOR
    p2.font.name = FONT_BODY
    p2.alignment = PP_ALIGN.CENTER


# ═════════════════════════════════════════════════════════════════════════════
# Generate new figures
# ═════════════════════════════════════════════════════════════════════════════

def generate_new_figures():
    """Generate the 6 new presentation-specific figures."""
    print("Generating new presentation figures...")

    # 1. Yee Grid Diagram
    fig, ax = plt.subplots(1, 1, figsize=(8, 6))
    ax.set_xlim(-0.5, 4.5)
    ax.set_ylim(-0.5, 4.5)
    ax.set_aspect('equal')
    ax.invert_yaxis()

    for i in range(5):
        ax.axhline(i, color='#CCCCCC', linewidth=0.5)
        ax.axvline(i, color='#CCCCCC', linewidth=0.5)

    for i in range(4):
        for j in range(4):
            ax.plot(j, i, 'o', color='#1B3A5C', markersize=12, zorder=5)
            if i == 1 and j == 1:
                ax.annotate('$E_z[i,j]$', (j, i), xytext=(j + 0.15, i - 0.2),
                            fontsize=13, fontweight='bold', color='#1B3A5C')

    for i in range(4):
        for j in range(4):
            ax.plot(j, i + 0.5, 's', color='#C62828', markersize=10, zorder=5)
            if i == 0 and j == 1:
                ax.annotate('$H_x[i,j]$', (j, i + 0.5), xytext=(j + 0.15, i + 0.3),
                            fontsize=13, fontweight='bold', color='#C62828')

    for i in range(4):
        for j in range(4):
            ax.plot(j + 0.5, i, '^', color='#2E7D32', markersize=10, zorder=5)
            if i == 1 and j == 0:
                ax.annotate('$H_y[i,j]$', (j + 0.5, i), xytext=(j + 0.65, i - 0.2),
                            fontsize=13, fontweight='bold', color='#2E7D32')

    ax.set_xlabel('x index (j)', fontsize=14)
    ax.set_ylabel('z index (i)', fontsize=14)
    ax.set_title('Yee Staggered Grid (2D TMz)', fontsize=16, fontweight='bold')

    legend_elements = [
        plt.Line2D([0], [0], marker='o', color='w', markerfacecolor='#1B3A5C',
                   markersize=12, label='$E_z$ (integer nodes)'),
        plt.Line2D([0], [0], marker='s', color='w', markerfacecolor='#C62828',
                   markersize=10, label='$H_x$ (half-step in z)'),
        plt.Line2D([0], [0], marker='^', color='w', markerfacecolor='#2E7D32',
                   markersize=10, label='$H_y$ (half-step in x)'),
    ]
    ax.legend(handles=legend_elements, loc='lower right', fontsize=11,
              framealpha=0.9)
    plt.tight_layout()
    plt.savefig(f'{PRES_FIG_DIR}/yee_grid_diagram.png', dpi=200, bbox_inches='tight')
    plt.close()
    print("  1/6 Yee grid diagram")

    # 2. Scanning Geometry
    fig, ax = plt.subplots(1, 1, figsize=(10, 5))
    ax.set_xlim(-10, 510)
    ax.set_ylim(-20, 310)
    ax.invert_yaxis()
    ax.set_aspect('equal')

    ax.add_patch(plt.Rectangle((0, 0), 500, 40, facecolor='#E3F2FD', edgecolor='none'))
    ax.text(250, 20, 'Air', ha='center', va='center', fontsize=14, color='#1565C0')
    ax.add_patch(plt.Rectangle((0, 40), 500, 260, facecolor='#BCAAA4', edgecolor='none'))
    ax.text(250, 250, 'Concrete ($\\varepsilon_r = 6.0$)', ha='center', va='center',
            fontsize=14, color='#4E342E')

    for x_c in [150, 250, 350]:
        circle = plt.Circle((x_c, 90), 6, facecolor='#616161', edgecolor='red',
                             linewidth=2, zorder=5)
        ax.add_patch(circle)
    ax.text(410, 90, 'Steel rebars\n(PEC)', fontsize=11, color='#C62828', va='center')

    ax.annotate('', xy=(50, 32), xytext=(450, 32),
                arrowprops=dict(arrowstyle='<->', color='#2E7D32', lw=2))
    ax.text(250, 25, 'Scan path (50-450 mm)', ha='center', fontsize=11,
            color='#2E7D32', fontweight='bold')

    ax.plot(200, 36, 'v', color='#1B3A5C', markersize=12)
    ax.plot(220, 36, 'v', color='#C62828', markersize=12)
    ax.text(200, 44, 'Tx', ha='center', fontsize=10, color='#1B3A5C')
    ax.text(220, 44, 'Rx', ha='center', fontsize=10, color='#C62828')

    ax.annotate('', xy=(420, 40), xytext=(420, 90),
                arrowprops=dict(arrowstyle='<->', color='black', lw=1.5))
    ax.text(435, 65, '50 mm\ncover', fontsize=10, va='center')

    ax.annotate('', xy=(150, 105), xytext=(250, 105),
                arrowprops=dict(arrowstyle='<->', color='black', lw=1.5))
    ax.text(200, 115, '100 mm', fontsize=10, ha='center')

    ax.axhline(40, color='#1565C0', linewidth=2)
    ax.set_xlabel('Lateral position x [mm]', fontsize=13)
    ax.set_ylabel('Depth z [mm]', fontsize=13)
    ax.set_title('GPR Scanning Configuration', fontsize=16, fontweight='bold')
    plt.tight_layout()
    plt.savefig(f'{PRES_FIG_DIR}/scanning_geometry.png', dpi=200, bbox_inches='tight')
    plt.close()
    print("  2/6 Scanning geometry")

    # 3. Inversion Workflow
    fig, ax = plt.subplots(1, 1, figsize=(12, 4))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 4)
    ax.axis('off')

    boxes = [
        (0.5, 1.5, 'Initial\nModel'),
        (2.5, 1.5, 'Forward\nFDTD'),
        (4.5, 1.5, 'Compute\nResidual'),
        (6.5, 1.5, 'Adjoint\nFDTD'),
        (8.5, 1.5, 'Cross-\nCorrelate'),
        (10.5, 1.5, 'Update\nModel'),
    ]
    colors = ['#E3F2FD', '#BBDEFB', '#90CAF9', '#BBDEFB', '#90CAF9', '#C8E6C9']

    for (x, y, text), color in zip(boxes, colors):
        ax.add_patch(FancyBboxPatch((x, y), 1.5, 1.2, boxstyle="round,pad=0.1",
                                      facecolor=color, edgecolor='#1B3A5C', linewidth=2))
        ax.text(x + 0.75, y + 0.6, text, ha='center', va='center',
                fontsize=11, fontweight='bold', color='#1B3A5C')

    for i in range(5):
        x_start = boxes[i][0] + 1.5
        x_end = boxes[i + 1][0]
        ax.annotate('', xy=(x_end, 2.1), xytext=(x_start, 2.1),
                    arrowprops=dict(arrowstyle='->', color='#1B3A5C', lw=2))

    ax.annotate('', xy=(0.5, 1.5), xytext=(10.5 + 0.75, 1.5),
                arrowprops=dict(arrowstyle='->', color='#2E7D32', lw=2,
                                connectionstyle='arc3,rad=0.4'))
    ax.text(6, 0.4, 'Iterate until convergence', ha='center', fontsize=12,
            color='#2E7D32', fontstyle='italic')

    ax.set_title('Adjoint-State Inversion Workflow', fontsize=16, fontweight='bold',
                 color='#1B3A5C', pad=20)
    plt.tight_layout()
    plt.savefig(f'{PRES_FIG_DIR}/inversion_workflow.png', dpi=200, bbox_inches='tight')
    plt.close()
    print("  3/6 Inversion workflow")

    # 4. Dispersive comparison (run actual FDTD)
    from core.materials import MaterialModel
    from core.materials_dispersive import DispersiveMaterialModel
    from core.fdtd import FDTDSimulator
    from core.fdtd_dispersive import FDTDSimulatorDispersive
    from core.source import ricker_wavelet, generate_time_array

    t = generate_time_array(cfg.NT, cfg.DT)
    wav = ricker_wavelet(t, cfg.F_CENTER)
    n = cfg.NPML
    siz = int(np.round(cfg.TX_Z / cfg.DZ)) + n
    six = int(np.round(0.25 / cfg.DX)) + n
    riz = int(np.round(0.20 / cfg.DZ)) + n
    rix = six + 10
    iz_a = int(np.round(cfg.CONCRETE_TOP / cfg.DZ)) + n

    m_nd = MaterialModel(cfg.NZ, cfg.NX)
    m_nd.epsilon_r[:iz_a, :] = cfg.AIR_EPSR
    m_nd.epsilon_r[iz_a:, :] = cfg.CONCRETE_EPSR
    m_nd.sigma[iz_a:, :] = cfg.CONCRETE_SIGMA
    tr_nd = FDTDSimulator(m_nd).run(wav, siz, six, riz, rix)['trace']

    m_d = DispersiveMaterialModel(cfg.NZ, cfg.NX)
    m_d.epsilon_r[:iz_a, :] = cfg.AIR_EPSR
    m_d.eps_inf[:iz_a, :] = cfg.AIR_EPSR
    m_d.eps_static[:iz_a, :] = cfg.AIR_EPSR
    m_d.epsilon_r[iz_a:, :] = 7.0
    m_d.sigma[iz_a:, :] = cfg.CONCRETE_SIGMA
    m_d.eps_inf[iz_a:, :] = 5.0
    m_d.eps_static[iz_a:, :] = 7.0
    m_d.tau[iz_a:, :] = 0.1e-9
    tr_d = FDTDSimulatorDispersive(m_d).run(wav, siz, six, riz, rix)['trace']

    fig, ax = plt.subplots(figsize=(10, 5))
    ax.plot(t * 1e9, tr_nd, 'b-', lw=1.5, label='Non-dispersive ($\\varepsilon_r=6.0$)')
    ax.plot(t * 1e9, tr_d, 'r--', lw=1.5,
            label='Dispersive (Debye: $\\varepsilon_\\infty=5.0$, $\\varepsilon_s=7.0$)')
    ax.set_xlim(0, 6)
    ax.set_xlabel('Time [ns]', fontsize=13)
    ax.set_ylabel('$E_z$ amplitude', fontsize=13)
    ax.set_title('Effect of Debye Dispersion on GPR Trace', fontsize=15, fontweight='bold')
    ax.legend(fontsize=12)
    ax.grid(True, alpha=0.3)
    rms = np.sqrt(np.mean((tr_nd - tr_d)**2)) / np.max(np.abs(tr_nd)) * 100
    ax.text(0.97, 0.95, f'RMS difference: {rms:.1f}%', transform=ax.transAxes,
            ha='right', va='top', fontsize=13, bbox=dict(boxstyle='round', fc='wheat', alpha=0.8))
    plt.tight_layout()
    plt.savefig(f'{PRES_FIG_DIR}/dispersive_comparison.png', dpi=200, bbox_inches='tight')
    plt.close()
    print("  4/6 Dispersive comparison")

    # 5. Antenna comparison (run actual FDTD)
    from core.geometry import build_rebar_model
    from core.source_antenna import run_with_antenna

    model_a = build_rebar_model()
    sa_iz = int(np.round(cfg.TX_Z / cfg.DZ)) + n
    sa_ix = int(np.round(0.25 / cfg.DX)) + n
    ra_iz = int(np.round(cfg.RX_Z / cfg.DZ)) + n
    ra_ix = sa_ix + int(np.round(cfg.TX_RX_OFFSET / cfg.DX))

    tr_pt = FDTDSimulator(model_a).run(wav, sa_iz, sa_ix, ra_iz, ra_ix)['trace']
    tr_dp = run_with_antenna(model_a, wav, sa_iz, sa_ix, ra_iz, ra_ix, 0.020)['trace']

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5))
    ax1.plot(t * 1e9, tr_pt, 'b-', lw=1.5, label='Point source')
    ax1.plot(t * 1e9, tr_dp, 'r--', lw=1.5, label='Dipole (20 mm)')
    ax1.set_xlim(0, 6)
    ax1.set_xlabel('Time [ns]', fontsize=13)
    ax1.set_ylabel('$E_z$ amplitude', fontsize=13)
    ax1.set_title('Raw Amplitude Comparison', fontsize=14, fontweight='bold')
    ax1.legend(fontsize=11)
    ax1.grid(True, alpha=0.3)

    ax2.plot(t * 1e9, tr_pt / np.max(np.abs(tr_pt)), 'b-', lw=1.5, label='Point (norm.)')
    ax2.plot(t * 1e9, tr_dp / np.max(np.abs(tr_dp)), 'r--', lw=1.5, label='Dipole (norm.)')
    ax2.set_xlim(0, 6)
    ax2.set_xlabel('Time [ns]', fontsize=13)
    ax2.set_ylabel('Normalised amplitude', fontsize=13)
    ax2.set_title('Normalised Shape Comparison', fontsize=14, fontweight='bold')
    ax2.legend(fontsize=11)
    ax2.grid(True, alpha=0.3)

    rms_a = np.sqrt(np.mean((tr_pt - tr_dp)**2)) / np.max(np.abs(tr_pt)) * 100
    fig.suptitle(f'Point Source vs Dipole Antenna (RMS diff: {rms_a:.1f}%)',
                 fontsize=15, fontweight='bold', y=1.02)
    plt.tight_layout()
    plt.savefig(f'{PRES_FIG_DIR}/antenna_comparison.png', dpi=200, bbox_inches='tight')
    plt.close()
    print("  5/6 Antenna comparison")

    # 6. Multi-scale bar chart
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 5))
    freqs = ['0.5 GHz', '1.0 GHz', '1.5 GHz']
    x = np.arange(3)
    w = 0.25

    x_errs = {'Rebar 1': [8.4, 2.3, 4.5], 'Rebar 2': [8.2, 0.9, 0.9],
              'Rebar 3': [10.4, 1.4, 0.6]}
    z_errs = {'Rebar 1': [7.1, 2.0, 0.2], 'Rebar 2': [3.4, 1.7, 0.3],
              'Rebar 3': [6.2, 0.7, 3.4]}
    colors_bar = ['#1B3A5C', '#2E7D32', '#C62828']

    for k, (label, vals) in enumerate(x_errs.items()):
        ax1.bar(x + k * w, vals, w, label=label, color=colors_bar[k], alpha=0.85)
    ax1.set_xticks(x + w)
    ax1.set_xticklabels(freqs, fontsize=12)
    ax1.set_ylabel('Lateral error [mm]', fontsize=13)
    ax1.set_title('Position Error by Stage', fontsize=14, fontweight='bold')
    ax1.legend(fontsize=11)
    ax1.grid(True, alpha=0.3, axis='y')

    for k, (label, vals) in enumerate(z_errs.items()):
        ax2.bar(x + k * w, vals, w, label=label, color=colors_bar[k], alpha=0.85)
    ax2.set_xticks(x + w)
    ax2.set_xticklabels(freqs, fontsize=12)
    ax2.set_ylabel('Depth error [mm]', fontsize=13)
    ax2.set_title('Depth Error by Stage', fontsize=14, fontweight='bold')
    ax2.legend(fontsize=11)
    ax2.grid(True, alpha=0.3, axis='y')

    fig.suptitle('Multi-Scale Frequency Continuation (0.5 → 1.0 → 1.5 GHz)',
                 fontsize=15, fontweight='bold', y=1.02)
    plt.tight_layout()
    plt.savefig(f'{PRES_FIG_DIR}/multiscale_bar.png', dpi=200, bbox_inches='tight')
    plt.close()
    print("  6/6 Multi-scale bar chart")

    # Extract a single frame from the GIF for the animation slide
    gif_path = 'outputs/animations/wave_propagation.gif'
    if os.path.exists(gif_path):
        img = PILImage.open(gif_path)
        # Seek to a mid-animation frame
        try:
            img.seek(40)
        except EOFError:
            img.seek(0)
        img.save(f'{PRES_FIG_DIR}/wave_frame.png')
        print("  Extracted wave animation frame")


# ═════════════════════════════════════════════════════════════════════════════
# Build slides
# ═════════════════════════════════════════════════════════════════════════════

def build_presentation():
    """Assemble all slides into a PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── SLIDE 1: Title ──
    slide = new_slide(prs)
    slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    slide.shapes[0].fill.solid()
    slide.shapes[0].fill.fore_color.rgb = TITLE_COLOR

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(2.5),
                 "2D FDTD Forward Modeling & Full-Waveform\nInversion for Ground-Penetrating Radar",
                 size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Rebar Detection in Reinforced Concrete",
                 size=22, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.6),
                 "Abdullah Maruf", size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    set_notes(slide, "Introduce yourself. This is a complete from-scratch implementation "
              "in Python — no external EM solvers used.")

    # ── SLIDE 2: Motivation ──
    slide = new_slide(prs)
    add_title_bar(slide, "Motivation: GPR for Concrete Inspection")
    add_bullets(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(5), [
        "Ground-Penetrating Radar (GPR) is the standard",
        "  non-destructive method for imaging rebars",
        "",
        "Challenge: B-scan hyperbolas give qualitative",
        "  locations — not quantitative geometry",
        "",
        "Goal: Recover rebar positions & sizes from",
        "  waveform data using physics-based inversion",
        "",
        "From scratch: ~2,500 lines of Python",
    ], size=16)
    add_figure(slide, f'{PRES_FIG_DIR}/scanning_geometry.png',
               Inches(6.2), Inches(1.3), width=Inches(6.5))
    set_notes(slide, "GPR used in bridge deck inspection, concrete assessment. "
              "Current practice relies on manual interpretation. FWI automates this.")

    # ── SLIDE 3: Outline ──
    slide = new_slide(prs)
    add_title_bar(slide, "Outline")
    sections = [
        ("1.", "Forward Modeling", "2D TMz FDTD with CPML boundaries"),
        ("2.", "Inversion", "Adjoint-state gradient + geometry-based FWI"),
        ("3.", "GPU Acceleration", "CuPy-based speedup (3-7x)"),
        ("4.", "Extensions", "Dispersion, antenna, multi-scale"),
        ("5.", "Conclusions", "Key results and future work"),
    ]
    for i, (num, title, desc) in enumerate(sections):
        y = Inches(1.6 + i * 1.05)
        add_text_box(slide, Inches(1.5), y, Inches(1), Inches(0.5),
                     num, size=24, bold=True, color=ACCENT_COLOR)
        add_text_box(slide, Inches(2.3), y, Inches(4), Inches(0.5),
                     title, size=22, bold=True, color=TITLE_COLOR)
        add_text_box(slide, Inches(6.5), y, Inches(5.5), Inches(0.5),
                     desc, size=16, color=BODY_COLOR)
    set_notes(slide, "5 sections, ~20 min talk.")

    # ── SLIDE 4: Yee Grid ──
    slide = new_slide(prs)
    add_title_bar(slide, "FDTD Method: The Yee Staggered Grid")
    add_bullets(slide, Inches(0.5), Inches(1.4), Inches(5), Inches(5), [
        "TMz polarisation: Ez, Hx, Hy",
        "E and H staggered by half a grid cell",
        "Leapfrog time-stepping: H → E → H → ...",
        "Second-order accurate, energy-conserving",
        "No matrix inversions (fully explicit)",
    ], size=16)
    add_figure(slide, f'{PRES_FIG_DIR}/yee_grid_diagram.png',
               Inches(5.8), Inches(1.3), width=Inches(6.8))
    set_notes(slide, "Yee's 1966 innovation. Staggering gives 2nd-order accuracy for free "
              "and automatically satisfies div B = 0.")

    # ── SLIDE 5: Parameters ──
    slide = new_slide(prs)
    add_title_bar(slide, "Simulation Parameters")
    params = [
        ["Grid spacing", "2 mm", ">=16 pts/wavelength"],
        ["Time step", "4.24 ps", "90% of CFL limit"],
        ["Domain", "500 x 300 mm", "280 x 180 cells + PML"],
        ["Source", "1.5 GHz Ricker", "Standard for concrete GPR"],
        ["Rebars", "3 x 12 mm dia", "50 mm cover, 100 mm spacing"],
        ["CPML", "15 layers, cubic", "~60 dB absorption"],
        ["Scan", "101 positions", "50-450 mm, 4 mm step"],
    ]
    # Build table
    rows, cols = len(params) + 1, 3
    tbl = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5),
                                  Inches(10), Inches(4.5)).table
    headers = ["Parameter", "Value", "Justification"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR

    for i, row in enumerate(params):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
    set_notes(slide, "Every parameter has a physical justification. "
              "2mm grid = 16 pts/wavelength at max Ricker frequency in concrete.")

    # ── SLIDE 6: Ground Truth ──
    slide = new_slide(prs)
    add_title_bar(slide, "Ground-Truth Model")
    add_figure(slide, f'{FIG_DIR}/ground_truth.png',
               Inches(1.5), Inches(1.3), width=Inches(10))
    set_notes(slide, "The model we aim to recover. Concrete eps_r=6, rebar PEC (sigma=1e7). "
              "3 rebars at 50mm cover, 100mm spacing.")

    # ── SLIDE 7: B-scan ──
    slide = new_slide(prs)
    add_title_bar(slide, "Forward Result: B-Scan Radargram")
    add_figure(slide, f'{FIG_DIR}/bscan.png',
               Inches(0.5), Inches(1.2), width=Inches(8.5))
    add_key_number(slide, Inches(9.5), Inches(2), "3", "rebar hyperbolas")
    add_key_number(slide, Inches(9.5), Inches(4), "6/6", "tests pass")
    set_notes(slide, "Each rebar produces a characteristic hyperbola. "
              "Apex = lateral position, curvature = concrete velocity.")

    # ── SLIDE 8: Animation ──
    slide = new_slide(prs)
    add_title_bar(slide, "Wave Propagation Through Concrete")
    frame_path = f'{PRES_FIG_DIR}/wave_frame.png'
    if os.path.exists(frame_path):
        add_figure(slide, frame_path, Inches(1.5), Inches(1.3), width=Inches(10))
    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "Full animation available: outputs/animations/wave_propagation.gif",
                 size=14, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)
    set_notes(slide, "Show the GIF animation live if possible. "
              "Wavefront refracts at air-concrete, reflects off rebars, CPML absorbs at edges.")

    # ── SLIDE 9: Inversion Problem ──
    slide = new_slide(prs)
    add_title_bar(slide, "The Inversion Problem")
    add_bullets(slide, Inches(0.5), Inches(1.5), Inches(7), Inches(3), [
        "Given: observed B-scan data d_obs",
        "Find: permittivity model m that explains the data",
        "Misfit: J(m) = 0.5 * ||d_obs - d_syn(m)||^2",
        "",
        "Naive approach: perturb each of 50,000 grid cells",
        "  → 50,000 forward simulations per gradient!",
    ], size=18)
    add_key_number(slide, Inches(9), Inches(2.5), "50,000", "parameters")
    add_key_number(slide, Inches(9), Inches(4.5), "~4 hrs", "per naive gradient")
    set_notes(slide, "This frames the core challenge. Direct FD gradient is O(N) simulations.")

    # ── SLIDE 10: Adjoint ──
    slide = new_slide(prs)
    add_title_bar(slide, "Adjoint-State Method: Efficient Gradients")
    add_figure(slide, f'{PRES_FIG_DIR}/inversion_workflow.png',
               Inches(0.5), Inches(1.2), width=Inches(12))
    add_bullets(slide, Inches(0.5), Inches(4.5), Inches(8), Inches(2.5), [
        "Forward: source at Tx → illumination of each grid cell",
        "Adjoint: time-reversed residual at Rx → sensitivity map",
        "Cross-correlate → gradient at every pixel simultaneously",
    ], size=16)
    add_key_number(slide, Inches(9.5), Inches(4.5), "2 sims", "not 50,000")
    set_notes(slide, "The adjoint method computes the full gradient with only 2 simulations "
              "per source: one forward, one adjoint.")

    # ── SLIDE 11: Gradient Validation ──
    slide = new_slide(prs)
    add_title_bar(slide, "Adjoint Gradient Validation")
    add_bullets(slide, Inches(0.5), Inches(1.5), Inches(7), Inches(4), [
        "Gold standard: compare adjoint gradient vs finite differences",
        "Perturbed individual pixels, ran forward simulation",
        "Agreement at max-sensitivity pixel: 0.4% relative error",
        "Random pixels: 0.85-1.09 ratio (expected FD truncation error)",
        "",
        "Conclusion: adjoint implementation is mathematically correct",
        "Any convergence issues are NOT gradient errors",
    ], size=17)
    add_key_number(slide, Inches(9.5), Inches(2.5), "0.4%", "FD agreement")
    set_notes(slide, "This is critical for trusting inversion results.")

    # ── SLIDE 12: Pixel-wise challenge — experimental evidence ──
    slide = new_slide(prs)
    add_title_bar(slide, "Experiment: Pixel-Wise FWI Step-Size Analysis")
    add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5), [
        "Gradient validated (0.4% vs FD) → ran systematic step-size sweep:",
        "",
        "  alpha = 5e-5  :  dJ = -1.75e-5   (descent)",
        "  alpha = 1e-4  :  dJ = -1.70e-5   (descent)",
        "  alpha = 3e-4  :  dJ = -6.7e-6    (descent, marginal)",
        "  alpha = 5e-4  :  dJ = +6.4e-6    (ascent — overshooting!)",
        "  alpha = 1e-3  :  dJ = +4.5e-5    (ascent)",
        "",
        "Only alpha <= 3e-4 produces descent. At best step: -1.7e-5 per iter",
        "from J = 14 — would need ~50,000 iterations (impractical)",
        "",
        "Root cause: 50K pixels changing simultaneously causes nonlinear overshoot",
        "Decision: pivot to geometry-based approach using prior knowledge",
    ], size=15)
    set_notes(slide, "This is the evidence that motivated the pivot. The gradient is "
              "correct but the step size is severely limited by the dimensionality. "
              "We also tested pseudo-Hessian preconditioning (Experiment 06) — "
              "negligible improvement for this well-illuminated domain.")

    # ── SLIDE 13: Geometry insight ──
    slide = new_slide(prs)
    add_title_bar(slide, "Key Insight: Geometry-Based Parameterisation")
    # Two columns
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Pixel-wise FWI", size=22, bold=True, color=RGBColor(0xC6, 0x28, 0x28))
    add_bullets(slide, Inches(0.5), Inches(2.2), Inches(5.5), Inches(4), [
        "50,000 parameters",
        "Adjoint gradient (expensive)",
        "Extremely slow convergence",
        "Needs advanced preconditioning",
    ], size=16)

    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Geometry-based FWI", size=22, bold=True, color=ACCENT_COLOR)
    add_bullets(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4), [
        "9 parameters (x, z, r per rebar)",
        "FD gradient (affordable at 9 params)",
        "Fast convergence (Nelder-Mead)",
        "GPU-accelerated forward sims",
    ], size=16)
    add_key_number(slide, Inches(5), Inches(5), "50K → 9", "parameters")
    set_notes(slide, "We know rebars are circular. Exploit this prior knowledge.")

    # ── SLIDE 14: Results ──
    slide = new_slide(prs)
    add_title_bar(slide, "Inversion Results: Rebar Recovery")
    add_figure(slide, f'{FIG_DIR}/inversion_comparison.png',
               Inches(0.3), Inches(1.2), width=Inches(9))
    add_key_number(slide, Inches(9.5), Inches(1.5), "< 3 mm", "position error")
    add_key_number(slide, Inches(9.5), Inches(3.2), "< 1 mm", "depth error")
    add_key_number(slide, Inches(9.5), Inches(4.9), "3.0%", "NRMS model error")
    set_notes(slide, "All 3 rebars recovered with sub-grid accuracy. "
              "Grid spacing is 2mm, so <3mm error is within 1-2 cells.")

    # ── SLIDE 15: Convergence ──
    slide = new_slide(prs)
    add_title_bar(slide, "Inversion Convergence")
    add_figure(slide, f'{FIG_DIR}/convergence.png',
               Inches(0.5), Inches(1.2), width=Inches(8.5))
    add_bullets(slide, Inches(9.3), Inches(1.5), Inches(3.8), Inches(5), [
        "275 function evaluations",
        "(Nelder-Mead simplex)",
        "",
        "43.5% misfit reduction",
        "",
        "Rapid initial decrease:",
        "optimizer quickly finds",
        "approximate rebar positions",
        "",
        "Gradual refinement:",
        "fine-tuning radii and",
        "sub-grid positioning",
    ], size=14)
    set_notes(slide, "The convergence shows two phases: rapid coarse localisation "
              "in the first ~50 evaluations, then gradual refinement of positions and radii.")

    # ── SLIDE 15b: Signal Comparison ──
    slide = new_slide(prs)
    add_title_bar(slide, "Signal Comparison: Observed vs Synthetic Traces")
    add_figure(slide, f'{FIG_DIR}/signal_comparison.png',
               Inches(0.3), Inches(1.2), width=Inches(9))
    add_bullets(slide, Inches(9.5), Inches(1.5), Inches(3.5), Inches(5), [
        "Blue: observed data",
        "  (from true model)",
        "Red: synthetic data",
        "  (from inverted model)",
        "",
        "Good match at all 5",
        "scan positions, especially",
        "at rebar reflection times",
        "(1.5-4 ns)",
        "",
        "Residual mismatch at late",
        "times from multiple",
        "reflections not captured",
        "by 3-rebar parameterisation",
    ], size=14)
    set_notes(slide, "The trace comparison confirms that the inverted model "
              "reproduces the observed waveforms well. The 5 positions sample "
              "the full scan aperture.")

    # ── SLIDE 16: B-scan fit ──
    slide = new_slide(prs)
    add_title_bar(slide, "B-Scan Data Fit: Observed vs Inverted")
    for i, (fname, lbl) in enumerate([
        ('bscan_observed.png', 'Observed'),
        ('bscan_inverted.png', 'Inverted'),
        ('residual_bscan.png', 'Residual'),
    ]):
        x = Inches(0.3 + i * 4.3)
        add_figure(slide, f'{FIG_DIR}/{fname}', x, Inches(1.5), width=Inches(4.1))
        add_text_box(slide, x, Inches(1.2), Inches(4.1), Inches(0.4),
                     lbl, size=16, bold=True, alignment=PP_ALIGN.CENTER)
    set_notes(slide, "Visual proof: the inverted model reproduces the observed B-scan well.")

    # ── SLIDE 17: GPU — measured benchmarks ──
    slide = new_slide(prs)
    add_title_bar(slide, "GPU Acceleration: Measured Benchmarks")
    add_figure(slide, f'{FIG_DIR}/gpu_scaling.png',
               Inches(0.3), Inches(1.2), width=Inches(8.5))
    add_bullets(slide, Inches(9), Inches(1.5), Inches(4), Inches(5), [
        "Measured on NVIDIA GB10",
        "(DGX Spark, 128 GB)",
        "",
        "5 grid sizes benchmarked:",
        "  50K → 3.2M cells",
        "  500 time steps each",
        "  2 runs averaged",
        "",
        "CuPy: drop-in NumPy",
        "replacement for CUDA",
        "",
        "Speedup increases with",
        "grid size (more parallel",
        "work saturates GPU)",
    ], size=14)
    add_key_number(slide, Inches(9.5), Inches(5.5), "3.2-7.0x", "measured speedup")
    set_notes(slide, "These are MEASURED benchmarks, not estimates. "
              "5 grid sizes from 50K to 3.2M cells, each timed with 2 runs averaged. "
              "Speedup peaks at 7x for 806K cells. For 3D problems with millions of "
              "cells, speedup would be 30-100x.")

    # ── SLIDE 18: GPU CPML ──
    slide = new_slide(prs)
    add_title_bar(slide, "Extension: GPU CPML (Full GPU Time-Stepping)")
    add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5), [
        "Ported all 8 PML boundary corrections to GPU (CuPy vectorised ops)",
        "Replaced per-layer Python loops with 2D array broadcasting",
        "All psi arrays on GPU — complete GPU-resident time-stepping loop",
        "",
        "Verification: bit-identical traces vs CPU (zero error at float64 precision)",
        "",
        "Performance at project grid (180x280):",
        "  CPU with CPML: 1.70 s    |    GPU v2 with CPML: 1.10 s    →   1.5x speedup",
        "",
        "Why modest speedup? CPML operates on thin boundary strips (15 cells wide)",
        "  — low parallelism per kernel launch. At larger grids (e.g., 3D with",
        "  millions of cells), the field-update cost dominates and overall GPU",
        "  speedup approaches the 3-7x measured for field updates alone.",
    ], size=15)
    add_key_number(slide, Inches(9.5), Inches(5.5), "Bit-identical", "to CPU")
    set_notes(slide, "GPU CPML adds correctness (complete absorbing boundaries on GPU) "
              "but the performance gain at this small 2D grid is modest. The thin PML "
              "strips (15x280 or 180x15) don't saturate GPU parallelism. For production "
              "3D grids the CPML overhead is proportionally much smaller.")

    # ── SLIDE 19: Dispersive ──
    slide = new_slide(prs)
    add_title_bar(slide, "Extension: Dispersive Materials (Debye Model)")
    add_figure(slide, f'{PRES_FIG_DIR}/dispersive_comparison.png',
               Inches(0.3), Inches(1.2), width=Inches(8))
    add_bullets(slide, Inches(8.5), Inches(1.5), Inches(4.5), Inches(4), [
        "Single-pole Debye ADE",
        "eps_inf=5.0, eps_s=7.0",
        "tau = 0.1 ns",
        "",
        "Adds polarisation current",
        "to E-field update",
        "",
        "Significant for quantitative",
        "FWI with real data",
    ], size=14)
    add_key_number(slide, Inches(9.5), Inches(5.5), "~20%", "trace effect")
    set_notes(slide, "Debye dispersion causes frequency-dependent velocity and attenuation.")

    # ── SLIDE 20: Antenna ──
    slide = new_slide(prs)
    add_title_bar(slide, "Extension: Distributed Dipole Antenna Model")
    add_figure(slide, f'{PRES_FIG_DIR}/antenna_comparison.png',
               Inches(0.3), Inches(1.2), width=Inches(12.5))
    add_key_number(slide, Inches(10), Inches(5.5), "~40%", "amplitude effect")
    set_notes(slide, "Point source is omnidirectional; dipole has directional pattern. "
              "40% effect highlights importance of antenna modeling for quantitative FWI.")

    # ── SLIDE 21: Multi-scale ──
    slide = new_slide(prs)
    add_title_bar(slide, "Extension: Multi-Scale Frequency Continuation")
    add_figure(slide, f'{PRES_FIG_DIR}/multiscale_bar.png',
               Inches(0.3), Inches(1.2), width=Inches(12.5))
    add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
                 "Progressive refinement: 0.5 GHz (coarse) → 1.0 GHz (medium) → 1.5 GHz (fine)",
                 size=15, alignment=PP_ALIGN.CENTER)
    set_notes(slide, "Low frequencies have smooth misfit landscape — optimizer escapes local minima. "
              "Each stage inherits the previous result as starting point.")

    # ── SLIDE 22: Architecture ──
    slide = new_slide(prs)
    add_title_bar(slide, "Software Architecture")
    modules = [
        ("core/", "FDTD engine, CPML, materials, geometry, source, scan"),
        ("inversion/", "Adjoint, objective, regularisation, optimizer, geometry, multi-scale"),
        ("gpu/", "CuPy FDTD, GPU CPML, benchmarking"),
        ("visualization/", "All plotting modules"),
        ("tests/", "6 verification tests (all pass)"),
    ]
    for i, (mod, desc) in enumerate(modules):
        y = Inches(1.6 + i * 0.85)
        add_text_box(slide, Inches(1), y, Inches(3), Inches(0.5),
                     mod, size=20, bold=True, color=TITLE_COLOR)
        add_text_box(slide, Inches(4), y, Inches(8), Inches(0.5),
                     desc, size=16)
    add_key_number(slide, Inches(9), Inches(5.5), "~2,500", "lines of Python")
    set_notes(slide, "Entire codebase from scratch. No external EM solver libraries.")

    # ── SLIDE 23: Summary ──
    slide = new_slide(prs)
    slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    slide.shapes[0].fill.solid()
    slide.shapes[0].fill.fore_color.rgb = TITLE_COLOR

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8),
                 "Summary & Key Contributions", size=30, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    contributions = [
        "1.  Complete 2D TMz FDTD solver with CPML — 6/6 tests pass",
        "2.  Adjoint gradient validated to 0.4% vs finite differences",
        "3.  Pixel-wise FWI: correct but impractical → evidence-based pivot",
        "4.  Geometry inversion: 9 params, < 3 mm position, < 1 mm depth",
        "5.  GPU acceleration: 3.2-7.0x measured speedup, bit-identical CPML",
        "6.  Extensions: Debye dispersion (~20% effect), dipole antenna",
        "     (~40% effect), multi-scale frequency continuation",
        "",
        "All from scratch — ~2,500 lines of documented Python",
    ]
    for i, item in enumerate(contributions):
        add_text_box(slide, Inches(1.5), Inches(1.6 + i * 0.6), Inches(10), Inches(0.5),
                     item, size=18, color=WHITE)

    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
                 "Thank you — Questions?", size=24, bold=True,
                 color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
    set_notes(slide, "Open for questions. Mention: dual-parameter adjoint and 3D extension "
              "as natural next steps.")

    # Save
    out_path = 'outputs/GPR_FDTD_FWI_Presentation.pptx'
    prs.save(out_path)
    print(f"\nPresentation saved: {out_path}")
    print(f"  Slides: {len(prs.slides)}")
    return out_path


# ═════════════════════════════════════════════════════════════════════════════
# Main
# ═════════════════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    generate_new_figures()
    build_presentation()
