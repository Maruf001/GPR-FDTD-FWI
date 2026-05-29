"""
Generate the second single-rebar deck (post baseline pipeline).

Tells the story of experiments 024-058 as the answer to:
  what assumption, if broken, would actually move radius?
with the five published FWI papers acting as the menu of hypotheses tested.

Output:
    outputs/GPR_FDTD_FWI_SingleRebar_Next.pptx

Custom figures used by the deck are regenerated here so the script is
self-contained. They land in:
    outputs/presentation_figures_single_rebar_next/
"""
import csv
import json
import os
import sys
from collections import defaultdict

import numpy as np
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# -----------------------------------------------------------------------------
# Paths
# -----------------------------------------------------------------------------

PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPUTS_DIR = os.path.join(PROJECT_ROOT, "outputs")
EXP_DIR = os.path.join(OUTPUTS_DIR, "experiments")
FIG_DIR = os.path.join(OUTPUTS_DIR, "presentation_figures_single_rebar_next")
PPTX_PATH = os.path.join(OUTPUTS_DIR, "GPR_FDTD_FWI_SingleRebar_Next.pptx")
os.makedirs(FIG_DIR, exist_ok=True)


def exp_dir(prefix):
    """Resolve experiments/NNN_* by leading three-digit prefix."""
    for name in os.listdir(EXP_DIR):
        if name.startswith(prefix + "_"):
            return os.path.join(EXP_DIR, name)
    raise FileNotFoundError(prefix)


# -----------------------------------------------------------------------------
# Design language (matches the first deck for visual consistency)
# -----------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1A, 0x23, 0x32)
STEEL = RGBColor(0x2C, 0x5F, 0x7C)
TEAL = RGBColor(0x00, 0x7C, 0x7C)
AMBER = RGBColor(0xE8, 0x91, 0x3A)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
CHARCOAL = RGBColor(0x2C, 0x3E, 0x50)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
PALE_AMBER = RGBColor(0xFD, 0xF1, 0xE0)
PALE_GREEN = RGBColor(0xE6, 0xF5, 0xEC)
PALE_RED = RGBColor(0xFA, 0xE5, 0xE2)
DIVIDER_BG = RGBColor(0x1E, 0x3A, 0x5F)

FONT = "Calibri"

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Helvetica", "Arial", "DejaVu Sans"],
    "font.size": 12,
    "axes.labelsize": 13,
    "axes.titlesize": 15,
    "axes.titleweight": "bold",
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "axes.edgecolor": "#CCCCCC",
    "grid.alpha": 0.3,
    "grid.color": "#CCCCCC",
})

SLIDE_NUM = [0]


# -----------------------------------------------------------------------------
# Slide helpers
# -----------------------------------------------------------------------------


def new_slide(prs, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    if bg is not None:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
    return s


def accent_bar(slide, color=STEEL, h=Inches(0.06)):
    sh = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def page_num(slide):
    box = slide.shapes.add_textbox(
        SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.45),
        Inches(0.6), Inches(0.3),
    )
    p = box.text_frame.paragraphs[0]
    p.text = str(SLIDE_NUM[0])
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT


def exp_badge(slide, label):
    """Small amber chip in the top-right that names the experiment reference."""
    if not label:
        return
    chip_w = Inches(2.4)
    chip_l = SLIDE_W - chip_w - Inches(0.7)
    chip_t = Inches(0.32)
    sh = slide.shapes.add_shape(5, chip_l, chip_t, chip_w, Inches(0.40))
    sh.fill.solid()
    sh.fill.fore_color.rgb = PALE_AMBER
    sh.line.color.rgb = AMBER
    sh.line.width = Pt(1.0)
    p = sh.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT


def title_bar(slide, title, subtitle=None, left=Inches(0.7), exp_label=None):
    accent_bar(slide, STEEL)
    # Title gets the full left-to-badge horizontal budget.
    title_w = Inches(9.3) if exp_label else Inches(12.0)
    box = slide.shapes.add_textbox(left, Inches(0.22), title_w, Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    if len(title) > 50:
        p.font.size = Pt(20)
    elif len(title) > 40:
        p.font.size = Pt(22)
    elif len(title) > 32:
        p.font.size = Pt(24)
    else:
        p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT

    bar = slide.shapes.add_shape(1, left, Inches(0.88), Inches(2.4), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()

    if subtitle:
        box2 = slide.shapes.add_textbox(left, Inches(0.95), Inches(12), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        p2.font.name = FONT
    exp_badge(slide, exp_label)
    page_num(slide)


def text(slide, left, top, w, h, body, size=14, bold=False, color=None,
         alignment=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color or CHARCOAL
    p.font.name = FONT
    p.alignment = alignment
    return tf


def bullets(slide, left, top, w, h, items, size=14, color=None, spacing=6):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        level = 0
        s = item
        while s.startswith("  "):
            level += 1
            s = s[2:]
        p.text = s
        p.level = level
        p.font.size = Pt(size)
        p.font.color.rgb = color or CHARCOAL
        p.font.name = FONT
        p.space_after = Pt(spacing)
    return tf


def figure(slide, path, left, top, w=None, h=None):
    if not os.path.exists(path):
        print(f"  WARNING: missing figure {path}")
        return None
    kw = {}
    if w is not None:
        kw["width"] = w
    if h is not None:
        kw["height"] = h
    return slide.shapes.add_picture(path, left, top, **kw)


def caption(slide, left, top, w, body, color=GRAY):
    box = slide.shapes.add_textbox(left, top, w, Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


def panel(slide, left, top, w, h, fill, border=None):
    sh = slide.shapes.add_shape(5, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = border
        sh.line.width = Pt(1.0)
    sh.text_frame.text = ""
    return sh


def set_notes(slide, body):
    slide.notes_slide.notes_text_frame.text = body


# -----------------------------------------------------------------------------
# Custom figures
# -----------------------------------------------------------------------------


def load_spectrum_csv(path):
    by_label = defaultdict(lambda: ([], []))
    with open(path) as f:
        for row in csv.DictReader(f):
            f_hz = float(row["frequency_hz"])
            a = float(row["amplitude"])
            by_label[row["label"]][0].append(f_hz)
            by_label[row["label"]][1].append(a)
    return {k: (np.array(v[0]), np.array(v[1])) for k, v in by_label.items()}


def fig_spectrum_clip(label_filter, title, out_name, fmax_ghz=4.0,
                      color_overrides=None, normalize_each=True):
    """Re-render 043 spectrum entries with x-axis clipped to a useful band."""
    data = load_spectrum_csv(os.path.join(
        exp_dir("043"), "data", "spectrum_records.csv"))
    fig, ax = plt.subplots(1, 1, figsize=(11, 4.5))
    color_overrides = color_overrides or {}
    for label in label_filter:
        if label not in data:
            continue
        f, a = data[label]
        mask = f <= fmax_ghz * 1e9
        f = f[mask]
        a = a[mask]
        if normalize_each:
            peak = a.max() if a.max() > 0 else 1.0
            a = a / peak
        ax.plot(f / 1e9, a, label=label.replace("_", " "),
                color=color_overrides.get(label), linewidth=1.8)
    ax.set_xlabel("Frequency [GHz]")
    ax.set_ylabel("Normalized mean amplitude")
    ax.set_title(title, color="#1A2332")
    ax.set_xlim(0, fmax_ghz)
    ax.grid(True, alpha=0.3)
    ax.legend(loc="upper right", fontsize=10)
    out = os.path.join(FIG_DIR, out_name)
    fig.tight_layout()
    fig.savefig(out, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  wrote {out}")
    return out


def fig_pebdd_stage_progression():
    """Bar chart of recovered radius at each stage of the spectrum-driven PEBDD."""
    stages = [
        ("Initial seed", 6.76),
        ("Stage 1\n0.35–1.10", 6.864),
        ("Stage 2\n0.35–1.50", 6.865),
        ("Stage 3\n0.35–2.00", 6.896),
        ("Stage 4\n0.35–2.50", 6.930),
        ("Final polish\nfull band", 6.000),
    ]
    labels = [s[0] for s in stages]
    radii = [s[1] for s in stages]
    colors = ["#7F8C8D", "#E8913A", "#E8913A", "#E8913A", "#E8913A", "#27AE60"]

    fig, ax = plt.subplots(1, 1, figsize=(12, 4.6))
    bars = ax.bar(labels, radii, color=colors, edgecolor="#1A2332")
    ax.axhline(6.0, color="#27AE60", linestyle="--", linewidth=1.5, label="true r = 6.0 mm")
    ax.set_ylabel("Recovered radius [mm]")
    ax.set_ylim(5.5, 7.2)
    ax.set_title("Radius across PEBDD stages (exp 045)", color="#1A2332")
    for bar, r in zip(bars, radii):
        ax.text(bar.get_x() + bar.get_width() / 2, r + 0.04,
                f"{r:.3f}", ha="center", fontsize=10,
                color="#1A2332", fontweight="bold")
    ax.legend(loc="upper right", fontsize=10)
    out = os.path.join(FIG_DIR, "pebdd_stage_progression.png")
    fig.tight_layout()
    fig.savefig(out, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  wrote {out}")
    return out


def fig_freq_weight_noise_compare():
    """Side-by-side: 5% noise (049) and 10% noise (047) margin curves."""
    paths = [
        (exp_dir("049") + "/figures/frequency_weight_radius_profiles.png", "5% noise (exp 049)"),
        (exp_dir("047") + "/figures/frequency_weight_radius_profiles.png", "10% noise (exp 047)"),
    ]
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.6))
    for ax, (p, title) in zip(axes, paths):
        img = plt.imread(p)
        ax.imshow(img)
        ax.axis("off")
        ax.set_title(title, fontsize=14, color="#1A2332")
    out = os.path.join(FIG_DIR, "freq_weight_noise_compare.png")
    fig.tight_layout()
    fig.savefig(out, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  wrote {out}")
    return out


def fig_wavelet_progression():
    """Four-up: raw mismatch -> amp -> +time -> +freq."""
    paths = [
        (exp_dir("052") + "/figures/wavelet_mismatch_radius_profiles.png",
         "Raw (exp 052)"),
        (exp_dir("053") + "/figures/wavelet_mismatch_radius_profiles.png",
         "+ amplitude scalar (exp 053)"),
        (exp_dir("054") + "/figures/wavelet_mismatch_radius_profiles.png",
         "+ time-shift grid (exp 054)"),
        (exp_dir("055") + "/figures/wavelet_mismatch_radius_profiles.png",
         "+ frequency-scale grid (exp 055)"),
    ]
    fig, axes = plt.subplots(2, 2, figsize=(13, 8.4))
    for ax, (p, title) in zip(axes.flat, paths):
        img = plt.imread(p)
        ax.imshow(img)
        ax.axis("off")
        ax.set_title(title, fontsize=13, color="#1A2332", pad=4)
    out = os.path.join(FIG_DIR, "wavelet_progression_4up.png")
    fig.tight_layout()
    fig.savefig(out, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  wrote {out}")
    return out


def fig_radius_progression_pebdd_first_pass():
    """Stage-by-stage radius for the first-pass PEBDD (exp 029-033)."""
    rows = [
        ("Initial 2 mm seed", 6.76, "#7F8C8D"),
        ("029 stage 1\n0.2–0.8 GHz", 6.573, "#E8913A"),
        ("030 stage 2\n0.2–1.1 GHz", 6.864, "#E8913A"),
        ("031 full-band\nPowell + polish", 6.000, "#27AE60"),
        ("032 full-band\nfrom stage 1, no polish", 6.955, "#C0392B"),
        ("033 stage 1\n→ polish only", 6.000, "#27AE60"),
    ]
    labels = [r[0] for r in rows]
    radii = [r[1] for r in rows]
    colors = [r[2] for r in rows]

    fig, ax = plt.subplots(1, 1, figsize=(12, 4.8))
    bars = ax.bar(labels, radii, color=colors, edgecolor="#1A2332")
    ax.axhline(6.0, color="#27AE60", linestyle="--", linewidth=1.5,
               label="true r = 6.0 mm")
    ax.set_ylabel("Recovered radius [mm]")
    ax.set_ylim(5.5, 7.3)
    ax.set_title("Radius across first-pass PEBDD variants",
                 color="#1A2332")
    for bar, r in zip(bars, radii):
        ax.text(bar.get_x() + bar.get_width() / 2, r + 0.04,
                f"{r:.3f}", ha="center", fontsize=10,
                color="#1A2332", fontweight="bold")
    ax.legend(loc="upper right", fontsize=10)
    out = os.path.join(FIG_DIR, "pebdd_first_pass_radius.png")
    fig.tight_layout()
    fig.savefig(out, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  wrote {out}")
    return out


# -----------------------------------------------------------------------------
# Slide constructors
# -----------------------------------------------------------------------------


def slide_title(prs):
    s = new_slide(prs, bg=NAVY)
    accent_bar(s, AMBER, Inches(0.10))

    chip = s.shapes.add_shape(5, Inches(0.9), Inches(1.4), Inches(3.0), Inches(0.4))
    chip.fill.solid()
    chip.fill.fore_color.rgb = AMBER
    chip.line.fill.background()
    text(s, Inches(0.9), Inches(1.4), Inches(3.0), Inches(0.4),
         "SECOND RESEARCH ARC", size=13, bold=True, color=NAVY,
         alignment=PP_ALIGN.CENTER)

    text(s, Inches(0.85), Inches(1.95), Inches(12), Inches(1.3),
         "Stress-testing the single-rebar pipeline",
         size=42, bold=True, color=WHITE)
    text(s, Inches(0.85), Inches(2.95), Inches(12), Inches(0.6),
         "What assumption, if broken, actually moves radius?",
         size=22, color=RGBColor(0xCC, 0xDD, 0xEE), italic=True)
    text(s, Inches(0.85), Inches(3.5), Inches(12), Inches(0.6),
         "Five published FWI strategies as the menu of hypotheses",
         size=18, color=RGBColor(0xAA, 0xBB, 0xCC))

    tile_y = Inches(4.7)
    tiles = [
        ("Source profiling",
         "Amplitude + time-shift + center-frequency scale"),
        ("Frequency weighting",
         "Lower bands for basin, 1.5 GHz for radius"),
        ("Polish, not Powell",
         "Local grid polish remains the radius selector"),
    ]
    for i, (lbl, body) in enumerate(tiles):
        left = Inches(0.85 + i * 4.0)
        panel(s, left, tile_y, Inches(3.7), Inches(1.5),
              RGBColor(0x24, 0x3B, 0x5E), border=AMBER)
        text(s, left, tile_y + Inches(0.18), Inches(3.7), Inches(0.5),
             lbl, size=18, bold=True, color=AMBER, alignment=PP_ALIGN.CENTER)
        text(s, left, tile_y + Inches(0.75), Inches(3.7), Inches(0.7),
             body, size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    text(s, Inches(0.85), SLIDE_H - Inches(0.6), Inches(12), Inches(0.3),
         "GPR-FDTD-FWI  •  experiments 024–058",
         size=12, color=GRAY)
    page_num(s)
    set_notes(s,
              "The first deck ended at exact-synthetic recovery on the 1 mm grid. "
              "This deck is the systematic stress test that follows, organised "
              "around five published FWI strategies. The deck order mirrors what "
              "we learned: rule out the OT hypothesis (cycle skipping), test "
              "bandwidth scheduling (helpful for seeding, not for radius), test "
              "frequency weighting (use 1.5 GHz-only or carry_low_25), test W2 "
              "(rejected as a radius objective), find the dominant new failure "
              "mode (source wavelet mismatch), and assemble a production "
              "source-profiled polish that recovers radius and source jointly.")


def slide_section_divider(prs, number, title, subtitle):
    s = new_slide(prs, bg=DIVIDER_BG)
    accent_bar(s, AMBER, Inches(0.08))

    text(s, Inches(1.5), Inches(2.0), Inches(2), Inches(1),
         f"{number:02d}", size=72, bold=True, color=AMBER)
    line = s.shapes.add_shape(1, Inches(3.4), Inches(2.2), Inches(0.04), Inches(2.4))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()
    # Auto-shrink title if long so it does not collide with the subtitle.
    title_size = Pt(34) if len(title) <= 24 else Pt(28)
    box = s.shapes.add_textbox(Inches(4.0), Inches(2.3), Inches(8.8), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = title_size
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    text(s, Inches(4.0), Inches(4.2), Inches(8.8), Inches(0.8),
         subtitle, size=16, color=RGBColor(0xAA, 0xBB, 0xCC),
         italic=True)
    page_num(s)


def slide_paper_quadrant(prs, idx, citation, subtitle,
                         problem, method, findings, adoption,
                         exp_label=None, notes_text=None):
    """Dense paper slide with four content cards in a 2x2 grid."""
    s = new_slide(prs)
    # Title is just the method name — the badge already says which paper number.
    title_bar(s, citation, subtitle, exp_label=exp_label)

    quadrants = [
        ("PROBLEM IT ATTACKS", problem, PALE_AMBER, AMBER),
        ("METHOD / ARCHITECTURE", method, PALE_BLUE, STEEL),
        ("KEY FINDINGS REPORTED", findings, LIGHT_BG, TEAL),
        ("OUR ADAPTATION", adoption, PALE_GREEN, GREEN),
    ]
    card_w = Inches(6.20)
    card_h = Inches(2.65)
    positions = [
        (Inches(0.30), Inches(1.40)),
        (Inches(6.83), Inches(1.40)),
        (Inches(0.30), Inches(4.20)),
        (Inches(6.83), Inches(4.20)),
    ]
    for (header, body, fill, border), (left, top) in zip(quadrants, positions):
        panel(s, left, top, card_w, card_h, fill, border=border)
        # Header strip
        text(s, left + Inches(0.15), top + Inches(0.05),
             card_w - Inches(0.3), Inches(0.32),
             header, size=11, bold=True, color=NAVY)
        # Body
        box = s.shapes.add_textbox(left + Inches(0.15), top + Inches(0.42),
                                   card_w - Inches(0.3), card_h - Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body
        p.font.size = Pt(11)
        p.font.color.rgb = CHARCOAL
        p.font.name = FONT
        p.space_after = Pt(0)

    if notes_text:
        set_notes(s, notes_text)


def slide_paper(prs, idx, citation, core, kept, deferred, headline,
                notes_text="", exp_label=None):
    s = new_slide(prs)
    title_bar(s, f"Paper {idx}: {citation}",
              "How we used it for the single-rebar radius problem",
              exp_label=exp_label)

    # Left: core + kept + deferred
    top = Inches(1.5)
    width = Inches(7.5)

    text(s, Inches(0.7), top, width, Inches(0.45),
         "Core idea", size=14, bold=True, color=NAVY)
    text(s, Inches(0.7), top + Inches(0.45), width, Inches(1.5),
         core, size=13, color=CHARCOAL)

    text(s, Inches(0.7), top + Inches(2.1), width, Inches(0.45),
         "What we kept", size=14, bold=True, color=GREEN)
    text(s, Inches(0.7), top + Inches(2.55), width, Inches(1.2),
         kept, size=13, color=CHARCOAL)

    text(s, Inches(0.7), top + Inches(3.95), width, Inches(0.45),
         "What we deferred / rejected", size=14, bold=True, color=RED)
    text(s, Inches(0.7), top + Inches(4.4), width, Inches(1.2),
         deferred, size=13, color=CHARCOAL)

    # Right: highlight card
    card_l = Inches(8.4)
    panel(s, card_l, top, Inches(4.5), Inches(5.6), LIGHT_BG, STEEL)
    text(s, card_l, top + Inches(0.2), Inches(4.5), Inches(0.5),
         "Single-line takeaway", size=13, bold=True, color=NAVY,
         alignment=PP_ALIGN.CENTER)
    text(s, card_l + Inches(0.2), top + Inches(0.8), Inches(4.1), Inches(4.5),
         headline, size=14, italic=True, color=CHARCOAL,
         alignment=PP_ALIGN.CENTER)

    if notes_text:
        set_notes(s, notes_text)


def slide_framing(prs, section_title, subtitle, hypothesis, did, result_lines,
                  exp_label=None, notes_text=None):
    s = new_slide(prs)
    title_bar(s, section_title, subtitle, exp_label=exp_label)

    top = Inches(1.5)
    text(s, Inches(0.7), top, Inches(12), Inches(0.4),
         "Hypothesis", size=14, bold=True, color=NAVY)
    text(s, Inches(0.7), top + Inches(0.4), Inches(12), Inches(0.9),
         hypothesis, size=14, color=CHARCOAL)

    text(s, Inches(0.7), top + Inches(1.5), Inches(12), Inches(0.4),
         "What we did", size=14, bold=True, color=NAVY)
    text(s, Inches(0.7), top + Inches(1.9), Inches(12), Inches(1.4),
         did, size=14, color=CHARCOAL)

    text(s, Inches(0.7), top + Inches(3.5), Inches(12), Inches(0.4),
         "What we'll see", size=14, bold=True, color=AMBER)
    bullets(s, Inches(0.7), top + Inches(3.95), Inches(12), Inches(2.0),
            result_lines, size=14, color=CHARCOAL, spacing=6)
    if notes_text:
        set_notes(s, notes_text)


def slide_figure(prs, title, subtitle, fig_path, caption_body=None,
                 fig_top=Inches(1.4), fig_height=None, fig_width=Inches(12.5),
                 footer=None, exp_label=None, notes_text=None,
                 max_fig_h_in=4.8):
    s = new_slide(prs)
    title_bar(s, title, subtitle, exp_label=exp_label)
    # Auto-size: if width-only sizing would overflow available vertical space,
    # constrain by height instead so the figure fits between title and caption.
    if fig_height:
        h = fig_height
        from PIL import Image as PILImage
        img = PILImage.open(fig_path)
        aspect = img.width / img.height
        w = Inches(h / Inches(1) * aspect)
        left = Inches((13.333 - w / Inches(1)) / 2)
        figure(s, fig_path, left, fig_top, h=h)
    else:
        from PIL import Image as PILImage
        img = PILImage.open(fig_path)
        aspect = img.width / img.height
        w_in = fig_width / Inches(1)
        h_in = w_in / aspect
        if h_in > max_fig_h_in:
            h_in = max_fig_h_in
            w_in = h_in * aspect
        left = Inches((13.333 - w_in) / 2)
        figure(s, fig_path, left, fig_top,
               w=Inches(w_in), h=Inches(h_in))
    if caption_body:
        caption(s, Inches(0.7), Inches(6.2), Inches(12), caption_body)
    if footer:
        text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
             footer, size=12, italic=True, color=AMBER)
    if notes_text:
        set_notes(s, notes_text)


def slide_figure_two(prs, title, subtitle, left_path, right_path,
                     left_caption, right_caption,
                     footer=None, exp_label=None, notes_text=None):
    s = new_slide(prs)
    title_bar(s, title, subtitle, exp_label=exp_label)
    figure(s, left_path, Inches(0.2), Inches(1.4), w=Inches(6.5))
    figure(s, right_path, Inches(6.7), Inches(1.4), w=Inches(6.5))
    caption(s, Inches(0.2), Inches(5.9), Inches(6.5), left_caption)
    caption(s, Inches(6.7), Inches(5.9), Inches(6.5), right_caption)
    if footer:
        text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
             footer, size=13, italic=True, color=AMBER)
    if notes_text:
        set_notes(s, notes_text)


def slide_combined(prs, title, subtitle, fig_path, body_lines,
                   fig_left=Inches(7.0), fig_w=Inches(6.1),
                   fig_top=Inches(1.4), exp_label=None, notes_text=None):
    s = new_slide(prs)
    title_bar(s, title, subtitle, exp_label=exp_label)
    bullets(s, Inches(0.6), Inches(1.5), Inches(6.2), Inches(5),
            body_lines, size=14, color=CHARCOAL, spacing=8)
    figure(s, fig_path, fig_left, fig_top, w=fig_w)
    if notes_text:
        set_notes(s, notes_text)


def slide_text_table(prs, title, subtitle, intro_lines, table_rows,
                     header_row=None, footer=None, exp_label=None,
                     notes_text=None):
    s = new_slide(prs)
    title_bar(s, title, subtitle, exp_label=exp_label)
    bullets(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.6),
            intro_lines, size=14, color=CHARCOAL, spacing=6)

    # Render table-like text in monospace using a single textbox
    box = s.shapes.add_textbox(Inches(0.7), Inches(3.1), Inches(12), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = False
    lines = []
    if header_row:
        lines.append(header_row)
        lines.append("-" * len(header_row))
    lines.extend(table_rows)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(2)
    if footer:
        text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
             footer, size=13, italic=True, color=AMBER)
    if notes_text:
        set_notes(s, notes_text)


def slide_two_column_text(prs, title, subtitle, left_title, left_body,
                          right_title, right_body, exp_label=None,
                          notes_text=None):
    s = new_slide(prs)
    title_bar(s, title, subtitle, exp_label=exp_label)
    # Left card
    panel(s, Inches(0.6), Inches(1.5), Inches(5.9), Inches(5.0),
          PALE_AMBER, AMBER)
    text(s, Inches(0.6), Inches(1.65), Inches(5.9), Inches(0.6),
         left_title, size=18, bold=True, color=NAVY,
         alignment=PP_ALIGN.CENTER)
    bullets(s, Inches(0.9), Inches(2.3), Inches(5.3), Inches(4.0),
            left_body, size=13, color=CHARCOAL, spacing=8)
    # Right card
    panel(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.0),
          PALE_BLUE, STEEL)
    text(s, Inches(6.8), Inches(1.65), Inches(5.9), Inches(0.6),
         right_title, size=18, bold=True, color=NAVY,
         alignment=PP_ALIGN.CENTER)
    bullets(s, Inches(7.1), Inches(2.3), Inches(5.3), Inches(4.0),
            right_body, size=13, color=CHARCOAL, spacing=8)
    if notes_text:
        set_notes(s, notes_text)


# -----------------------------------------------------------------------------
# Deck content
# -----------------------------------------------------------------------------


def slide_definitions_two_col(prs, title, subtitle, left_entries, right_entries,
                              exp_label="reference", notes_text=None):
    """A two-column reference slide: each entry is term + dash + definition."""
    s = new_slide(prs)
    title_bar(s, title, subtitle, exp_label=exp_label)

    def write_column(left_in, top_in_start, entries, col_w_in=6.05):
        top_in = top_in_start
        for term, definition in entries:
            box = s.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                       Inches(col_w_in), Inches(0.50))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run1 = p.add_run()
            run1.text = term
            run1.font.size = Pt(11)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run1.font.name = FONT
            run2 = p.add_run()
            run2.text = "  —  " + definition
            run2.font.size = Pt(10)
            run2.font.color.rgb = CHARCOAL
            run2.font.name = FONT
            # Approx height per entry; varies a little with line wrap.
            top_in += 0.45 if len(definition) < 80 else 0.62

    write_column(0.4, 1.40, left_entries)
    write_column(6.85, 1.40, right_entries)
    if notes_text:
        set_notes(s, notes_text)


def build_abbreviations(prs):
    """Acronyms used through the deck."""
    general = [
        ("FWI", "Full-Waveform Inversion — fit a wave-equation forward model to observed traces."),
        ("GPR", "Ground-Penetrating Radar."),
        ("FDTD", "Finite-Difference Time-Domain — the wave-equation solver this project uses."),
        ("CPML", "Convolutional Perfectly Matched Layer — the absorbing boundary used by the FDTD solver."),
        ("B-scan", "GPR data array: lateral scan position × two-way time."),
        ("LS / L2", "Least squares — the standard waveform misfit (sum of squared point-wise residuals)."),
        ("NRMS", "Normalised Root-Mean-Square misfit between two B-scans (or between two models)."),
        ("εr", "Relative permittivity (dielectric constant)."),
        ("σ", "Electrical conductivity, in S/m."),
        ("PEC", "Perfect Electric Conductor — the steel-rebar idealisation in our forward model."),
        ("Tx / Rx", "Transmitter / receiver antenna positions on the scan line."),
        ("fc", "Source-wavelet centre frequency."),
        ("Ricker wavelet", "The time-domain pulse used as our source (the negative second derivative of a Gaussian)."),
        ("ps / mm", "Picoseconds (10⁻¹² s) and millimetres — units we report time-shift and geometry in."),
    ]
    papers = [
        ("WRI", "Wavefield Reconstruction Inversion — paper 1 (Feng et al. 2022)."),
        ("PEBDD", "Progressively Expanded Bandwidths of the Data — paper 2 (Zhou et al. 2021)."),
        ("OT", "Optimal Transport — a family of distances between distributions."),
        ("OT-LS", "OT-then-LS hybrid objective — paper 3 (Hunziker et al. 2025)."),
        ("W2", "Quadratic Wasserstein distance — a specific OT-based objective (paper 4, Lu et al. 2024)."),
        ("IFWI", "Implicit FWI — paper 5 (Sun et al. 2025)."),
        ("Sinkhorn", "Entropy-regularised algorithm for cheap W2 computation."),
        ("Softplus", "log(1+e^(βx)) — turns a signed oscillating trace into a non-negative mass distribution."),
        ("β (W2)", "Softplus scale — larger β preserves more signed-shape information."),
        ("ε (W2)", "Sinkhorn entropy regularisation — smaller is more accurate, less numerically stable."),
        ("ds (W2)", "Sinkhorn downsample factor for trace length (cost control)."),
        ("RCCC", "Relative Cross-Correlation Criterion = |trace shift| / dominant period (per trace pair)."),
        ("NRCCC", "Normalised RCCC = fraction of traces with RCCC < 0.5 (half-period gate from the OT-LS paper)."),
        ("carry_low_25", "Weighting scheme: 0.25 × 1.0 GHz + 1.0 × 1.5 GHz (radius-friendly multi-frequency)."),
    ]
    slide_definitions_two_col(
        prs,
        "Abbreviations and shorthand",
        "Used throughout the deck — refer back as needed",
        general, papers,
        exp_label="reference",
        notes_text=(
            "This is a reference slide; we will not dwell on it. The rest "
            "of the deck uses these acronyms freely, so jump back here if "
            "anything is unclear. The left column has general FWI, GPR, "
            "and FDTD terms. The right column has paper-specific and "
            "numerical terms, including the five paper methods and the "
            "two W2 hyperparameters. The name carry_low_25 is one we "
            "coined for the frequency-weighting scheme that puts 25 "
            "percent weight on the 1 GHz residual and full weight on "
            "the 1.5 GHz residual; it shows up in the frequency-weighting "
            "section later."
        ))


def build_jargon(prs):
    """Key terms used through the deck."""
    left = [
        ("Powell", "Derivative-free local optimiser available in SciPy; explores along directions, not gradients. Our default for the continuous local refinement stage."),
        ("Basin", "The local valley in the objective landscape where an optimiser settles."),
        ("High-radius basin", "The specific wrong local minimum at r ≈ 6.95 mm that continuous Powell consistently finds in our setup."),
        ("Depth-radius coupling", "Near the true minimum, a slightly deeper rebar (z+) trades against a slightly larger radius (r+) and produces an almost identical waveform — this is what causes the high-radius basin."),
        ("Grid polish", "A deterministic small local (z, r) grid search after Powell. Corrects radius bias because rasterised geometry is piecewise-constant in physical coordinates, so smooth-tolerance optimisers cannot find the exact rasterised cell."),
        ("Top-k candidates", "Instead of one best (x, z, r), the polish reports the K best candidates and their misfits — the next-best alternative is auditable."),
        ("Twin experiment", "Both the 'observed' and the candidate B-scans come from the same FDTD solver, so any failure is in the inverse pipeline, not in the physics model."),
    ]
    right = [
        ("Cycle skipping", "When a modelled trace is more than half a period out of phase with the observed one, least-squares can settle into the wrong cycle. This is the classical failure mode the OT family of methods targets."),
        ("Variable projection (WRI)", "Solve a fast subproblem for the wavefield first (with the model fixed), then reduce the optimisation to just the material model."),
        ("Matched filter (PEBDD)", "Apply the same band-pass to both observed and modelled traces inside the objective — the distinguishing detail vs older bandwidth methods."),
        ("Source-wavelet mismatch", "The observed source wavelet differs from the modelled one (amplitude / time-zero / centre-frequency)."),
        ("Source profile", "A small set of nuisance parameters (amplitude scalar + time-shift grid + fc-scale grid) fit per candidate to absorb wavelet mismatch."),
        ("Distinct-radius margin", "Objective gap between the best radius and the best next-distinct-radius candidate — our standard confidence measure."),
        ("Residual spectrum", "FFT amplitude of (synthetic − observed) for a specific candidate; shows which frequencies that candidate fails to fit."),
    ]
    slide_definitions_two_col(
        prs,
        "Key terms and jargon",
        "Plain-language gloss for the rest of the deck",
        left, right,
        exp_label="reference",
        notes_text=(
            "Two entries on the left matter more than the others. "
            "'High-radius basin' is the wrong local minimum at about "
            "6.95 mm radius that Powell consistently finds — it will come "
            "up on most of the later slides. 'Depth-radius coupling' is "
            "the physical reason that basin exists: a deeper rebar trades "
            "against a larger radius and produces an almost identical "
            "waveform. The grid-polish entry is a carry-over from the "
            "first deck but central enough to repeat here. The right "
            "column lists features of the production source-profiled "
            "polish later in the deck — when 'source profile' and "
            "'distinct-radius margin' appear, the definitions are on "
            "this slide."
        ))


def build_papers_intro(prs):
    """Why we used five published FWI papers as a menu of hypotheses."""
    s = new_slide(prs)
    title_bar(s,
              "Five papers as the menu of hypotheses",
              "What this section is and how it lands in the deck",
              exp_label="setup")

    # Left: framing text
    bullets(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(5.4),
            [
                "The first deck ended at exact-synthetic recovery on the 1 mm grid. That left a sharper question:",
                "  what assumption, if broken, actually moves radius?",
                "",
                "We did not invent a new objective. Instead, the recent FWI literature already proposes five distinct strategies for the failure modes we might run into. Each paper is a hypothesis about which assumption matters:",
                "  Paper 1 (WRI): the wave equation is too tight a constraint",
                "  Paper 2 (PEBDD): the data bandwidth is too wide too soon",
                "  Paper 3 (OT-LS): the objective traps on cycle skips",
                "  Paper 4 (W2): least squares is the wrong distance",
                "  Paper 5 (IFWI): the model representation is too rigid",
                "",
                "Each paper was tested through a controlled landscape gate before any optimizer integration. Some were kept, some were rejected, some were deferred — with reasons.",
            ], size=14, color=CHARCOAL, spacing=4)

    # Right: paper stack with title + single-sentence position
    card_l = Inches(7.0)
    card_w = Inches(5.9)
    card_h = Inches(1.00)
    gap = Inches(0.08)
    papers = [
        ("Paper 1 · WRI 2022 (Feng et al.)",
         "Frequency-domain wavefield reconstruction + cumulative frequency strategy.",
         PALE_AMBER, AMBER),
        ("Paper 2 · PEBDD 2021 (Zhou et al.)",
         "Progressively expanded bandwidth, matched filter on both d_obs and d_syn.",
         PALE_BLUE, STEEL),
        ("Paper 3 · OT-LS 2025 (Hunziker et al.)",
         "Hybrid OT-then-LS with an NRCCC switching criterion.",
         LIGHT_BG, TEAL),
        ("Paper 4 · W2 GPR 2024 (Lu et al.)",
         "Quadratic Wasserstein objective via Softplus normalisation + Sinkhorn.",
         PALE_RED, RED),
        ("Paper 5 · IFWI 2025 (Sun et al.)",
         "Implicit neural representation (x, z) → (εr, σ) with frequency-principle prior.",
         PALE_GREEN, GREEN),
    ]
    top = Inches(1.5)
    for title_p, body_p, fill, border in papers:
        panel(s, card_l, top, card_w, card_h, fill, border=border)
        text(s, card_l + Inches(0.18), top + Inches(0.10),
             card_w - Inches(0.3), Inches(0.4),
             title_p, size=13, bold=True, color=NAVY)
        text(s, card_l + Inches(0.18), top + Inches(0.50),
             card_w - Inches(0.3), Inches(0.45),
             body_p, size=11, color=CHARCOAL)
        top += card_h + gap

    set_notes(s,
              "The framing of this whole arc is deliberate. We did not pick "
              "a single paper and try to apply it; we collected five "
              "distinct strategies for FWI failure modes and tested each "
              "as a hypothesis. The pattern was the same every time: "
              "build the smallest reusable piece of the paper's machinery, "
              "put it through a controlled landscape gate, and only "
              "decide whether to adopt it based on what the landscape "
              "says — before any optimiser integration. That order "
              "matters. If we had wired W2 into Powell first, we would "
              "have burned compute to learn the same lesson the landscape "
              "test delivered cheaply.")


def build_papers_bridge(prs):
    """Map each paper's hypothesis to the experiment that tested it."""
    s = new_slide(prs)
    title_bar(s,
              "From papers to experiments — the mapping",
              "Which experiments tested each paper's claim, and the verdict",
              exp_label="roadmap")

    headers = ["Source", "Hypothesis tested", "Experiments", "Gate result",
               "Verdict"]
    rows = [
        ("OT-LS",
         "Cycle skipping drives radius bias",
         "024–027",
         "NRCCC = 1.0 for every candidate, including the wrong-radius one",
         "ruled out for this problem; diagnostic kept",
         RED),
        ("PEBDD",
         "Bandwidth scheduling fixes the radius bias",
         "028–037, 043–045",
         "Powell stops in the high-radius basin at every stage; polish stays the radius selector",
         "kept as seed builder; spectrum-design tool added",
         AMBER),
        ("WRI",
         "Cumulative-frequency / weighted multi-frequency helps radius",
         "038–041, 046–049",
         "1.0 GHz weak; unweighted averaging dilutes margin; carry_low_25 best compromise",
         "kept as time-domain weighted LS; WRI solver deferred",
         AMBER),
        ("W2",
         "Softplus / Sinkhorn beats LS for radius selection",
         "048, 050, 051",
         "Shift convexity ✓ (gate 1); rebar margin collapses ×10⁴ (gate 2)",
         "rejected as radius objective; module kept as diagnostic",
         RED),
        ("PEBDD risk",
         "Source-wavelet mismatch (field-data flag) shifts radius",
         "052–055",
         "Raw mismatch pegs r to grid bound; structured source profile recovers truth",
         "amplitude + time + freq profile fixes all tested cases",
         GREEN),
        ("(geometry)",
         "Material parameters can hide a wrong radius",
         "056",
         "εr identified; σ saturates above 1e5; distinct-radius margin preserved",
         "no material params in the radius optimizer yet",
         AMBER),
        ("all 5",
         "Production runner combines surviving ideas",
         "057, 058",
         "Radius + source profile co-recovered under injected mismatch",
         "recommended pipeline going forward",
         GREEN),
        ("IFWI",
         "Neural implicit field needed for radius",
         "doc 29 only",
         "Flexible residual would absorb radius bias",
         "deferred to multi-rebar / field-data stage",
         RED),
    ]

    cols = len(headers)
    rcount = len(rows) + 1
    table_shape = s.shapes.add_table(
        rcount, cols,
        Inches(0.4), Inches(1.4),
        Inches(12.5), Inches(5.3),
    )
    tbl = table_shape.table

    # Column widths
    widths_in = [1.0, 3.0, 1.5, 4.3, 2.7]
    for i, w in enumerate(widths_in):
        tbl.columns[i].width = Inches(w)

    # Row heights
    tbl.rows[0].height = Inches(0.45)
    for r in range(1, rcount):
        tbl.rows[r].height = Inches(0.58)

    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.text = h
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT
        cell.margin_left = Inches(0.10)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

    # Data rows
    for r_idx, row in enumerate(rows, start=1):
        source, hyp, exps, gate, verdict, verdict_color = row
        contents = [source, hyp, exps, gate, verdict]
        bg = LIGHT_BG if r_idx % 2 == 1 else WHITE
        for c_idx, body in enumerate(contents):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = body
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.size = Pt(10)
            run.font.name = FONT
            if c_idx == 0:
                run.font.bold = True
                run.font.color.rgb = NAVY
            elif c_idx == 4:
                run.font.bold = True
                run.font.color.rgb = verdict_color
            else:
                run.font.color.rgb = CHARCOAL
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

    text(s, Inches(0.5), Inches(6.78), Inches(12.4), Inches(0.6),
         "Acronym key: OT-LS = optimal-transport / least-squares hybrid · PEBDD = progressively expanded "
         "bandwidths of the data · WRI = wavefield reconstruction inversion · W2 = quadratic Wasserstein · "
         "IFWI = implicit full-waveform inversion · NRCCC = fraction of trace pairs within half a period · "
         "LS = least squares.",
         size=10, italic=True, color=GRAY)

    set_notes(s,
              "This is the deck roadmap. The order of the rest of the "
              "presentation follows the rows of this table top to bottom. "
              "First we rule out OT as a radius solver and keep its "
              "diagnostic. Then we test bandwidth scheduling and keep it "
              "only as a seed builder. Then frequency weighting, where we "
              "discover unweighted multi-frequency dilutes radius margin "
              "and carry_low_25 is the best compromise. Then W2, where "
              "we run two gates — the shift-convexity test passes, the "
              "rebar landscape fails, and W2 gets rejected as a final "
              "radius objective. Then material parameters, which turn out "
              "not to explain the radius bias. Then the wavelet-mismatch "
              "finding, which is the new contribution from this arc. And "
              "finally the production runner that synthesises the "
              "surviving ideas. The OT-LS paper actually appears twice — "
              "once directly for the cycle-skipping hypothesis, and once "
              "indirectly because it shares PEBDD's concern about source-"
              "wavelet handling on field data.")


def build_papers(prs):
    slide_paper_quadrant(prs, 1,
        "Wavefield Reconstruction Inversion (WRI)",
        subtitle="Feng et al., Remote Sensing 2022 — frequency-domain GPR-FWI for underground pipeline imaging",
        problem=(
            "Standard FWI for GPR pipeline imaging is brittle: enforcing the wave "
            "equation exactly each iteration makes the objective highly nonlinear, "
            "so an inaccurate starting model creates cycle skipping and trapped "
            "local minima. Multiparameter inversion (permittivity and conductivity "
            "together) amplifies this; high-frequency content amplifies it further."
        ),
        method=(
            "Frequency-domain wavefield reconstruction inversion (WRI). The wave "
            "equation is no longer enforced exactly — it becomes a penalty term, "
            "so wavefield and model are co-updated. Variable projection: solve a "
            "least-squares problem for the reconstructed wavefield given the "
            "model, substitute back into a reduced objective over the model, "
            "optimize with L-BFGS over relative permittivity and conductivity. "
            "Pair this with a multi-scale cumulative frequency strategy: stages "
            "B1..B5 add 1, 2, 3, 4, or 6 new frequencies per step while keeping "
            "previous ones; inverse-frequency-squared weighting prevents the high "
            "band from dominating the objective."
        ),
        findings=(
            "Two synthetic pipeline cases and one field-style case. Cumulative "
            "schedules beat the simultaneous baseline in every test. Reported "
            "error reductions: 7–11% (permittivity) and 7–15% (conductivity) on "
            "case 1; 20–40% and 17–24% on case 2. Smaller cumulative batches were "
            "more accurate but cost more PDE solves. Permittivity reconstructed "
            "better than conductivity. Field case recovered plausible pipeline "
            "locations even with imperfect starts."
        ),
        adoption=(
            "Kept: the cumulative-frequency philosophy, adapted as time-domain "
            "per-frequency objective contributions and a --frequency-weights CLI. "
            "This is what enables the carry_low_25 scheme we use for radius.\n"
            "Deferred: full WRI itself. Faithful implementation needs a new "
            "frequency-domain Helmholtz/Maxwell operator, complex sparse solves, "
            "frequency-domain adjoint, and penalty continuation — its own "
            "project. The current radius bias is already well explained by "
            "source handling and weighted LS, so the cost is not justified. "
            "Documented in docs/experiments/28_wri_feasibility.md."
        ),
        exp_label="paper 1 / WRI",
        notes_text=(
            "Three points to land from this slide. First, the WRI "
            "relaxation is a structural change to the optimisation — the "
            "wave equation is no longer a hard constraint enforced "
            "exactly each iteration, it is a penalty term, so the "
            "wavefield can be slightly inconsistent with the model in "
            "exchange for better basin behaviour. Second, the cumulative-"
            "frequency idea is independent of WRI itself, and that is "
            "the part we actually adopted — we added per-frequency "
            "objective contributions and a frequency-weight CLI flag, "
            "which is what makes the carry_low_25 scheme possible later "
            "in the deck. Third, our deferral of full WRI is conditional, "
            "not categorical. If a future setup creates an initial-model "
            "failure that source handling cannot fix, WRI moves back "
            "onto the table; the docs/experiments/28 tracker has the "
            "design notes."
        ))

    slide_paper_quadrant(prs, 2,
        "Progressively Expanded Bandwidth FWI (PEBDD)",
        subtitle="Zhou, Klotzsche, Vereecken, NSG 2021 — crosshole GPR FWI with bandwidth scheduling + source-wavelet update",
        problem=(
            "Crosshole GPR FWI requires the modeled and observed traces to be "
            "within roughly half a wavelength or half a period at every position "
            "for least-squares to be safe. Ray-based starting models often violate "
            "this in high-contrast media. The visible failure mode is the "
            "optimizer settling into a depth/contrast tradeoff that fits the "
            "B-scan well but is physically wrong."
        ),
        method=(
            "PEBDD = Progressively Expanded Bandwidths of the Data. Apply the "
            "SAME tapered bandpass to both observed and modeled traces — that is "
            "the differentiator from older bandwidth-expansion methods, which "
            "filtered only the modeled side through the source wavelet. "
            "Algorithm: estimate an effective source wavelet; build a tapered "
            "filter; fix a low cut; start with a low high cut; run a small "
            "number of FWI iterations; expand the high cut by a fixed increment; "
            "repeat. After the staged build-up, use the recovered model as the "
            "starting model for full-band FWI, then update the source wavelet "
            "by deconvolution and run full-band FWI again."
        ),
        findings=(
            "Synthetic case I with a good ray-based start: permittivity MAE 2.0 → "
            "1.7 with PEBDD; data RMS roughly halved. Synthetic case II with a "
            "deliberately bad start outside the half-wavelength criterion: "
            "standard FWI MAE 3.6 → PEBDD MAE 2.1; standard FWI was trapped, "
            "PEBDD was not. Field-data RMS improved in all four crosshole "
            "sections tested; CPT comparisons mostly better, not uniformly."
        ),
        adoption=(
            "Kept: the matched-filter design (FFT bandpass applied to both "
            "d_obs and d_syn) as inversion/trace_filters.py, plus a new "
            "spectrum-design tool so the band edges come from real residual "
            "spectra instead of guesses. Used as the seed-building stage of "
            "the staged inversion.\n"
            "Adapted: the source-wavelet update is replaced in our pipeline "
            "by the explicit source profiling described later in the deck "
            "(amplitude / time-shift / center-frequency-scale)."
        ),
        exp_label="paper 2 / PEBDD",
        notes_text=(
            "The single most important detail in PEBDD is that the same "
            "band-pass filter is applied to both observed and modelled "
            "traces inside the objective. Older bandwidth methods "
            "filtered only the modelled side through the source wavelet, "
            "and they did not survive on field data. PEBDD did. Our "
            "adoption is faithful to that detail — the matched-filter "
            "band-pass lives in inversion/trace_filters.py and is "
            "applied identically to both d_obs and d_syn before the "
            "residual is computed. The schedule itself we re-derived "
            "from actual residual spectra after seeing that our first-"
            "pass guessed bands did not carry radius information; that "
            "story takes up the next several slides."
        ))

    slide_paper_quadrant(prs, 3,
        "Optimal-Transport / Least-Squares Hybrid (OT-LS)",
        subtitle="Hunziker, Meles, Linde, JAG 2025 — get into the basin with OT, refine with LS, switch via NRCCC",
        problem=(
            "Crosshole GPR FWI is cycle-skipping-prone with L2 because the "
            "objective has multiple local minima along trace shifts. When the "
            "wavelet's main lobe moves past truth, pointwise amplitude "
            "cancellation creates a misleading local decrease and the optimizer "
            "slides into the wrong basin. The proposed fix is a hybrid objective, "
            "not just a different optimizer."
        ),
        method=(
            "Two-stage objective. OT first: convert each trace into a 'fingerprint' "
            "— a 2D pseudo probability density over (time, amplitude), take its "
            "two 1D marginals, compute Wasserstein distances on those marginals. "
            "The fingerprint OT objective has a broad global basin that is "
            "shift-monotonic. Then switch to LS for sharpness once trace pairs "
            "are aligned. Switching is governed by the cross-correlation "
            "criterion RCCC = |trace shift| / dominant period; NRCCC is the "
            "fraction of traces with RCCC < 0.5; switch when NRCCC > Cs ≈ 0.7. "
            "Gradients computed sparsely at random master points away from "
            "antennas, then interpolated — implicit regularization plus runtime."
        ),
        findings=(
            "Synthetic crosshole with a large circular anomaly and a deliberately "
            "poor homogeneous starting model. Pure LS converged to bad local "
            "minima. Pure OT moved toward the right structure but did not resolve "
            "details. Hybrid OT → LS recovered good models for switching "
            "thresholds Cs ≈ 0.7–0.75. Runs that ended with NRCCC = 1 and "
            "weighted RMS close to 1 were reliable; runs that did not were not."
        ),
        adoption=(
            "Kept: the NRCCC trace-shift diagnostic, implemented in "
            "inversion/trace_distances.py and saved in every run summary. We use "
            "it as a safety report — does the model live in a least-squares-"
            "safe basin?\n"
            "Deferred: the OT objective itself. Across exact / 5% / 10% noise, "
            "NRCCC = 1.0 for our wrong-radius candidates — they are already in "
            "the LS-safe basin. The OT machinery would be solving a problem we "
            "do not have."
        ),
        exp_label="paper 3 / OT-LS",
        notes_text=(
            "Two design points on this slide deserve unpacking. First, "
            "the trace fingerprint — converting each trace to a 2D "
            "pseudo probability density over time and amplitude and "
            "then taking marginals — is what lets optimal transport "
            "work on signed oscillatory data without a heavy "
            "normalisation. That is a different approach from the W2 / "
            "Softplus method in the next paper, which transforms the "
            "trace itself into a non-negative mass distribution. "
            "Second, the sparse-gradient idea using random master "
            "points away from the antennas is a regularisation choice "
            "with a side benefit — gradients become independent of the "
            "specific distance measure, so swapping in a new objective "
            "does not require deriving a new adjoint. That property "
            "makes the OT-LS framework unusually flexible for trying "
            "objectives quickly."
        ))

    slide_paper_quadrant(prs, 4,
        "Quadratic Wasserstein GPR-FWI (W2)",
        subtitle="Lu et al., Remote Sensing 2024 — two-parameter GPR-FWI with W2 as a direct L2 replacement",
        problem=(
            "Two-parameter GPR FWI (relative permittivity and conductivity) is "
            "unstable under standard L2: high parameter sensitivity, weak "
            "conductivity recovery, noise vulnerability, strong initial-model "
            "dependence. The paper proposes a single objective change rather "
            "than a new optimization framework, so the rest of the FWI stack "
            "stays the same."
        ),
        method=(
            "Replace L2 with the quadratic Wasserstein distance W2. Three "
            "implementation pieces matter. (1) Softplus normalisation: raw "
            "signed GPR traces oscillate around zero, but optimal transport "
            "needs non-negative distributions with equal mass; softplus(b·x) / "
            "Σsoftplus turns the trace into a usable distribution. The scale b "
            "controls how aggressive the transform is. (2) Sinkhorn entropy "
            "regularisation: exact OT is too expensive for repeated FWI "
            "evaluations; ε trades transport accuracy for numerical stability. "
            "(3) Wrap in multi-scale frequency-domain FWI with L-BFGS over εr "
            "and σ, frequencies added in batches (e.g. 15 batches of 4 "
            "frequencies, 10 iterations each)."
        ),
        findings=(
            "Three numerical examples. W2-FWI depends less on the initial model "
            "than L2-FWI. W2-FWI is more robust to additive noise. W2-FWI "
            "improves conductivity recovery more strongly than permittivity "
            "recovery (the opposite ordering from L2-based work). L2-FWI can "
            "generate high-frequency artifacts and obscure targets when the "
            "starting model is poor. Adding higher frequencies did not always "
            "improve W2 results."
        ),
        adoption=(
            "Built: a standalone Softplus / Sinkhorn W2 module with unit tests "
            "(inversion/trace_wasserstein.py). Ran it through two gates before "
            "any optimizer integration.\n"
            "Gate 1 (shifted Ricker, paper claim): passed.\n"
            "Gate 2 (rebar landscape): failed. Softplus normalisation removes "
            "the amplitude content radius depends on; W2 margin collapses by "
            "four orders of magnitude vs LS.\n"
            "Decision: reject W2 as the final radius objective for this "
            "problem. Kept available as a basin diagnostic for future poor-"
            "initial-model / field-data scenarios."
        ),
        exp_label="paper 4 / W2",
        notes_text=(
            "The two-gate structure of our W2 evaluation is important — "
            "we did not dismiss W2 from theoretical arguments alone. We "
            "replicated the paper's positive claim about shift convexity "
            "on shifted Ricker pulses, and only then ran the same "
            "machinery on the actual rebar landscape. The negative "
            "result on the rebar landscape is specific to this single-"
            "rebar radius problem; it is not a general judgement on "
            "W2 for FWI. Two more things to keep in mind: β, the "
            "Softplus scale, and ε, the Sinkhorn entropy regularisation, "
            "are hyperparameters that must be chosen, and Sinkhorn cost "
            "grows with trace length. Both of these are reasons to test "
            "in a controlled landscape gate before wiring W2 into an "
            "optimiser."
        ))

    slide_paper_quadrant(prs, 5,
        "Implicit Multiparameter GPR-FWI (IFWI)",
        subtitle="Sun et al., GJI 2025 — neural implicit representation of the subsurface material model",
        problem=(
            "Standard FWI stores material parameters on a dense grid. That "
            "makes multiparameter inversion brittle: huge model dimensionality, "
            "manual parameter weighting between permittivity and conductivity, "
            "strong initial-model dependence, ad-hoc multi-scale schedules. "
            "The paper proposes a fundamentally different model representation "
            "rather than a new objective."
        ),
        method=(
            "Implicit FWI. Replace the grid with a small neural network: "
            "coordinates (x, z) → N_θ(x, z) → (εr, σ). The forward solver is "
            "still Maxwell / FDTD, the loss is still waveform mismatch — only "
            "the model parameterization changes. Optimize the network weights "
            "with Adam. Architecture details matter: SIREN / sinusoidal "
            "activations with a frequency scale ω₀ that controls detail "
            "sensitivity; too small ω₀ misses fine structure, too large "
            "ω₀ produces high-frequency noise. Dropout helps in some "
            "configurations. The paper also describes an RNN-style "
            "differentiable FDTD framing so AD gives gradients consistent "
            "with adjoint-state methods."
        ),
        findings=(
            "The 'frequency principle' of neural networks acts as an automatic "
            "multiscale prior: broad/low-frequency structure is learned before "
            "fine detail. IFWI reconstructed useful subsurface structure from "
            "poorer initial models than standard FWI. Permittivity and "
            "conductivity inverted simultaneously without the manual parameter "
            "weighting that classical multiparameter FWI relies on. Network "
            "architecture (depth, width, ω₀, dropout) materially affects "
            "results — these are tuning choices, not free parameters."
        ),
        adoption=(
            "Deferred. Two concrete reasons: (1) our parameter vector is "
            "[x, z, radius] — three numbers; an implicit field is overkill. "
            "(2) a flexible neural residual would risk absorbing radius bias "
            "into a 'background correction' instead of revealing it, which "
            "would make the radius estimate less meaningful, not more. The "
            "frequency-principle lesson is already implicit in our staged "
            "pipeline (basin first, detail later). Reserved for multi-rebar "
            "or field-data stages where explicit nuisance parameters become "
            "insufficient. See docs/experiments/29_ifwi_feasibility.md for "
            "the design of a constrained future prototype."
        ),
        exp_label="paper 5 / IFWI",
        notes_text=(
            "The IFWI deferral is the most subtle one in the deck. The "
            "risk is not that IFWI does not work — it does, in the "
            "paper's own results. The risk is what it can be made to "
            "fit. A free neural residual field around the rebar could "
            "match observed data perfectly while leaving the radius "
            "estimate wrong, because the radius error would be absorbed "
            "into the residual field as a small material perturbation. "
            "That would make our radius confidence number meaningless. "
            "The future-prototype design in tracker 29 is intentionally "
            "constrained for exactly this reason: explicit geometry "
            "stays the primary unknown, the neural residual is small, "
            "smooth, and regularised, and held-out scan positions "
            "validate that the residual generalises rather than absorbs "
            "errors. That constrained version is what we would build "
            "first if we did revisit IFWI."
        ))


def build_ot_diagnostic(prs):
    slide_section_divider(prs, 1, "Was it cycle skipping?",
                          "OT diagnostic on the existing high-radius candidate")

    slide_framing(prs,
                  "OT-style trace-shift diagnostic",
                  "Rule out the OT branch before committing to a new objective",
                  hypothesis=("If the high-radius Powell basin is a cycle-skipping "
                              "failure, the OT-LS paper's NRCCC criterion should "
                              "flag it as unsafe for least-squares."),
                  did=("We implemented RCCC = |trace shift| / dominant period and "
                       "NRCCC = fraction with RCCC < 0.5. Ran it post-hoc on the "
                       "known high-radius Powell candidate, the polished true "
                       "candidate, and the top-k polish candidates under 5% and "
                       "10% noise."),
                  result_lines=[
                      "If NRCCC < 1 for the wrong-radius case, the OT branch is worth pursuing.",
                      "If NRCCC = 1 for every candidate, this problem is not a transport problem.",
                  ],
                  exp_label="exp 024–027",
                  notes_text=(
                      "Module: inversion/trace_distances.py. Runner: "
                      "run_single_rebar_trace_diagnostics.py. We added per-frequency "
                      "trace_shift_by_frequency to every new single_rebar_summary. The "
                      "design point of mentioning this on the slide is *not* that OT "
                      "is bad — it is that the diagnostic itself ruled out the branch "
                      "for the current problem before any expensive optimizer "
                      "integration. Even at 10% noise the wrong-radius candidate "
                      "still passes the half-period gate."
                  ))

    # Custom rich layout for the NRCCC slide — not just a table.
    s = new_slide(prs)
    title_bar(s,
              "NRCCC results across the candidate set",
              "Does the half-period gate flag the wrong-radius candidates as unsafe?",
              exp_label="exp 024–026")

    # Definitions strip
    bullets(s, Inches(0.5), Inches(1.35), Inches(12.4), Inches(1.20),
            [
                "RCCC (Relative Cross-Correlation Criterion) = |trace shift| / dominant period, computed for each (observed, synthetic) trace pair.",
                "NRCCC (Normalised RCCC) = fraction of traces with RCCC < 0.5 — the half-period gate the OT-LS paper uses to decide a trace is least-squares-safe.",
                "Paper's switching rule: when NRCCC > ~0.7, switch from the OT objective to LS. If a candidate is already at NRCCC = 1.0, OT brings no extra basin-finding benefit.",
            ], size=11, color=CHARCOAL, spacing=2)

    # Table built with python-pptx for clean alignment
    headers = ["Candidate", "Misfit J", "NRCCC", "med RCCC", "max RCCC",
               "Radius", "Note"]
    data_rows = [
        ("polished true (014 / 015)", "0", "1.000", "0.000", "0.000",
         "6.000 mm", "exact recovery", GREEN),
        ("Powell high-radius", "2.08e-3", "1.000", "0.003", "0.003",
         "6.955 mm", "wrong radius, but phase-safe", RED),
        ("noise 10% top 1 (023)", "1.99e-1", "1.000", "0.000", "0.003",
         "6.000 mm", "recovered", GREEN),
        ("noise 10% top 2", "1.99e-1", "1.000", "0.000", "0.003",
         "6.000 mm", "alternative cell", GREEN),
        ("noise 10% high-radius", "2.01e-1", "1.000", "0.003", "0.006",
         "6.991 mm", "wrong radius, phase-safe", RED),
        ("noise 5% recovered", "5.86e-2", "1.000", "0.000", "0.000",
         "6.000 mm", "recovered", GREEN),
        ("noise 5% optimizer-end", "6.09e-2", "1.000", "0.003", "0.006",
         "6.973 mm", "wrong radius, phase-safe", RED),
    ]
    n_rows = len(data_rows) + 1
    table_shape = s.shapes.add_table(
        n_rows, len(headers),
        Inches(0.4), Inches(2.95),
        Inches(12.55), Inches(2.85),
    )
    tbl = table_shape.table
    widths = [3.0, 1.3, 1.0, 1.3, 1.3, 1.2, 3.45]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    tbl.rows[0].height = Inches(0.36)
    for r in range(1, n_rows):
        tbl.rows[r].height = Inches(0.34)

    for c_idx, h in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.text = h
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)

    for r_idx, row in enumerate(data_rows, start=1):
        contents = list(row[:7])
        note_color = row[7]
        bg = LIGHT_BG if r_idx % 2 == 1 else WHITE
        for c_idx, body in enumerate(contents):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = body
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.size = Pt(10)
            run.font.name = "Consolas" if c_idx in (1, 2, 3, 4, 5) else FONT
            if c_idx == 6:
                run.font.bold = True
                run.font.color.rgb = note_color
            elif c_idx == 0:
                run.font.color.rgb = NAVY
            else:
                run.font.color.rgb = CHARCOAL
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)

    # Interpretation panel at the bottom
    panel(s, Inches(0.4), Inches(5.95), Inches(12.55), Inches(1.10),
          PALE_AMBER, AMBER)
    text(s, Inches(0.55), Inches(5.99), Inches(12.2), Inches(0.4),
         "What this means", size=12, bold=True, color=NAVY)
    text(s, Inches(0.55), Inches(6.32), Inches(12.2), Inches(0.7),
         "Every candidate — including the wrong-radius Powell solutions at all noise levels — already sits at NRCCC = 1.0. "
         "The wrong-radius rows have a tiny non-zero median RCCC (≈ 0.003, i.e. 0.3% of a dominant period), but that is far below the 0.5 half-period gate. "
         "So the radius failure is not a phase / cycle-skipping problem; an OT objective designed for cycle skips cannot help here. "
         "Decision: rule the OT branch out for this problem; keep the NRCCC diagnostic for field-data use.",
         size=11, color=CHARCOAL)

    set_notes(s,
              "The result on this slide is what rules out the OT branch "
              "for our specific problem. Three observations to land. "
              "First, the OT-LS paper's whole motivation is cycle "
              "skipping at scale — the trigger condition for using OT "
              "is NRCCC less than 1, meaning some traces are out of "
              "phase by more than half a period. Second, every one of "
              "our wrong-radius candidates already sits at NRCCC equal "
              "to 1, including under 5% and 10% noise; the paper itself "
              "would say 'use least squares here'. Third, the wrong-"
              "radius rows do have very slightly higher RCCC than the "
              "polished true rows, but the absolute value — about "
              "0.003 of a period — is two orders of magnitude below the "
              "0.5 gate. So the radius issue is an amplitude and detail "
              "issue, not a transport issue. That observation drives "
              "the next section on PEBDD bandwidth and eventually "
              "motivates the source-profiling work that closes the "
              "deck. We keep the NRCCC diagnostic available because "
              "field data may still have basin-finding problems that "
              "this synthetic case does not.")


def build_pebdd_first_pass(prs, fig_first_pass_radius):
    slide_section_divider(prs, 2, "Bandwidth scheduling",
                          "PEBDD first pass, then a spectrum-driven re-design")

    slide_framing(prs,
                  "First PEBDD attempt with guessed bands",
                  "Same band-pass filter applied to observed and synthetic traces (matched-filter design)",
                  hypothesis=("A staged objective that starts in a low frequency band and "
                              "expands the high cut should reduce the high-radius Powell "
                              "bias (the wrong local minimum at r ≈ 6.95 mm that the "
                              "derivative-free SciPy Powell optimiser keeps finding) "
                              "before the polish stage."),
                  did=("Built a trace-filter helper and a staged runner. First "
                       "schedule used 0.2–0.8 then 0.2–1.1 then full band, with a "
                       "final coarse polish. Also ran controls without polish and a "
                       "polish-only path from the low-band seed."),
                  result_lines=[
                      "Exact data: low-band Powell did move radius closer to the truth, but full-band Powell pulled it back.",
                      "Noise: the low-band stage stayed in the high-radius basin.",
                      "Polish remained the reliable radius selector in every variant.",
                  ],
                  exp_label="exp 028–033",
                  notes_text=(
                      "At this stage the bands were chosen by inspection: "
                      "0.2 to 0.8 GHz, then 0.2 to 1.1, then 0.2 to 1.5, "
                      "then the full band. The matched band-pass filter "
                      "lives in inversion/trace_filters.py and is "
                      "exposed through the --objective-bandpass-ghz CLI "
                      "flag with a taper width. The fact that this "
                      "first-pass schedule did not move radius is what "
                      "motivated the spectrum-driven design on the next "
                      "few slides; we will see that the residual energy "
                      "that actually distinguishes nearby radii does "
                      "not live in the 0.2–0.8 GHz band at all."
                  ))

    slide_figure(prs,
                 "Radius across first-pass PEBDD variants",
                 "Single-frequency 1.5 GHz, 5 scan positions (PEBDD = progressively expanded bandwidths of the data)",
                 fig_first_pass_radius,
                 caption_body=("Each bar is the recovered radius at the end of one run. "
                               "'Powell' is SciPy's derivative-free local optimiser; 'polish' is a "
                               "deterministic small (z, r) grid search after Powell. Polish is the only "
                               "step that lands on 6.0 mm; full-band Powell without polish returns to the "
                               "high-radius basin (r ≈ 6.95 mm) — the wrong local minimum it consistently finds."),
                 exp_label="exp 029–033",
                 notes_text=(
                     "The bars go in chronological order. The grey bar on "
                     "the left is the initial seed at radius 6.76 mm. "
                     "Stage 1 with the low band brings radius down to "
                     "6.57 — that is a real improvement. Stage 2, which "
                     "expands the band to 1.1 GHz, drifts back up to "
                     "6.86. The red bar in the middle is the control: "
                     "stage 1 seed followed by a full-band Powell run "
                     "without polish — it returns to the same 6.955 mm "
                     "wrong basin as the original first-arc Powell "
                     "result, which tells us the regression is caused "
                     "by the full-band Powell objective itself, not by "
                     "any particular staging choice. The two green bars "
                     "are the only ones that land at the true 6.0 mm, "
                     "and both of them used the polish stage."
                 ))

    slide_figure(prs,
                 "Stage 1: low-band Powell",
                 "0.2–0.8 GHz objective; full unfiltered traces still used for NRMS",
                 os.path.join(exp_dir("029"), "figures",
                              "single_rebar_model_comparison.png"),
                 caption_body=("Stage-1 recovered r = 6.573 mm at x ≈ 249.8, z ≈ 90.6 mm. "
                               "The low band moved radius toward the truth but did not reach it."),
                 exp_label="exp 029",
                 notes_text=(
                     "The filtered objective at the end of this stage is "
                     "about 1e-4 and the full-data NRMS is about 1.07 "
                     "percent — both small numbers. But the recovered "
                     "model still has the rebar slightly off in depth "
                     "and at the wrong radius. The PEBDD-faithful detail "
                     "here is that the observed and synthetic traces are "
                     "filtered through the same band-pass before the "
                     "residual is computed; older bandwidth-expansion "
                     "methods filtered only the modelled side through "
                     "the source wavelet, and that is the difference "
                     "that made PEBDD survive on field data."
                 ))

    slide_figure(prs,
                 "Stage 2: expanded to 0.2–1.1 GHz",
                 "Bandwidth expansion from the stage-1 seed (same band-pass on observed and synthetic)",
                 os.path.join(exp_dir("030"), "figures",
                              "single_rebar_model_comparison.png"),
                 caption_body=("Recovered r drifts back up to 6.864 mm. Expanding the band re-introduces "
                               "the depth–radius coupling (a slightly deeper rebar trades against a "
                               "slightly larger radius to produce a near-identical waveform)."),
                 exp_label="exp 030",
                 notes_text=(
                     "Full-data NRMS (normalised RMS over the unfiltered traces) "
                     "actually improves to about 0.79%, but the recovered radius "
                     "worsens from 6.57 mm back toward 6.86 mm. This is the "
                     "single most important lesson from the PEBDD-first-pass "
                     "section: data fit and radius fidelity can move in opposite "
                     "directions when the objective expansion reactivates the "
                     "depth-radius tradeoff. The expanded band lets the "
                     "optimiser explore amplitude detail the low band could not "
                     "see, but the dominant valley in that detail is the high-"
                     "radius basin, not the true r = 6.0 mm minimum. This "
                     "behaviour was not predictable from the paper's claims "
                     "alone — it is specific to our rasterised-geometry setting "
                     "and is what motivates the polish stage that follows."
                 ))

    slide_figure(prs,
                 "Stage 3: full band + polish",
                 "Full-band Powell from stage 2, then a deterministic coarse grid polish",
                 os.path.join(exp_dir("031"), "figures",
                              "single_rebar_model_comparison.png"),
                 caption_body=("Polish recovers r = 6.000 mm (NRMS model = 0). Powell alone reverts "
                               "to r = 6.955 mm before polish takes over — same wrong basin as the "
                               "no-PEBDD path."),
                 exp_label="exp 031",
                 notes_text=(
                     "The top polish candidates here are radius 6.0 mm "
                     "with misfit zero, radius 6.2 with misfit 1.04e-3, "
                     "and radius 6.8 with misfit 2.08e-3. That is "
                     "exactly the same ranking the simpler single-"
                     "frequency polish produced in the first deck. The "
                     "implication is unambiguous: even with the multi-"
                     "stage bandwidth schedule wrapped around it, the "
                     "polish stage is what does the radius work. PEBDD "
                     "as we have implemented it has not added to radius "
                     "accuracy; it has only added compute cost. This is "
                     "the first hard evidence that PEBDD on its own "
                     "cannot solve radius, and it sets up the spectrum-"
                     "design section that asks whether the band choices "
                     "themselves were wrong. One caveat: this is the "
                     "exact-synthetic case. Under noise the same "
                     "conclusion holds, but the polish stage starts to "
                     "need careful early-stop thresholds to avoid "
                     "following noise instead of geometry."
                 ))

    slide_figure(prs,
                 "Control: full-band Powell, no polish",
                 "Same stage-1 seed, run full-band Powell without polish — isolates the bandwidth role",
                 os.path.join(exp_dir("032"), "figures",
                              "single_rebar_model_comparison.png"),
                 caption_body=("Final r = 6.955 mm. Confirms the regression we saw in stage 2 is the "
                               "full-band Powell objective itself, not a specific staging choice."),
                 exp_label="exp 032",
                 notes_text=(
                     "This run isolates blame between two suspects. "
                     "Either the staged path itself is unstable — maybe "
                     "stage 2 destabilised something — or the full-band "
                     "Powell objective is the cause because of the "
                     "depth-radius valley we discussed earlier. If the "
                     "staged path were at fault, removing the polish "
                     "step would still leave us at a moderate radius. "
                     "If the full-band objective were at fault, we would "
                     "land back at the original wrong basin. The result "
                     "is unambiguous: r = 6.955 mm, exactly the same as "
                     "the original first-arc Powell result. The full-"
                     "band Powell objective is the cause, not the "
                     "staging. That conclusion lets us promote the "
                     "'stage 1 seed → polish only' variant on the next "
                     "slide, which removes the full-band Powell step "
                     "entirely."
                 ))

    slide_figure(prs,
                 "Pragmatic path: stage 1 seed → polish only",
                 "Skip full-band Powell; use low-band seed to enter the polish window",
                 os.path.join(exp_dir("033"), "figures",
                              "single_rebar_model_comparison.png"),
                 caption_body=("Recovered r = 6.000 mm. The low-band seed plus polish "
                               "is the cleanest staged path on exact data."),
                 exp_label="exp 033",
                 notes_text=(
                     "This is the variant we promote on exact data: low-band "
                     "Powell to land in the right basin, then full-band "
                     "polish to select radius. It is faster than the full "
                     "Powell-plus-polish path because we skip the most "
                     "expensive Powell stage. The caveat is that this only "
                     "works on exact data — under noise the low-band stage "
                     "does not move radius much, and the next slide shows "
                     "the noise variant explicitly."
                 ))

    slide_text_table(prs,
                     "Noise variant: PEBDD seed + polish under noise",
                     "Experiments 034–037",
                     exp_label="exp 034–037",
                     notes_text=(
                         "Two additive-noise seeds, 5 percent and 10 "
                         "percent RMS. The low-band Powell stage stays "
                         "near r = 7.0 mm in both cases — it does not "
                         "move radius under noise. The polish stage is "
                         "what recovers the true 6.0 mm. So under noise "
                         "the low-band stage is a seed-builder, not a "
                         "radius corrector, and the headline is the same "
                         "as the exact-data case: bandwidth scheduling "
                         "moves rough seeds into the right window, but "
                         "the polish stage is the radius selector."
                     ),
                     intro_lines=[
                         "Under noise, the low-band Powell stage stops in the high-radius basin. The polish stage still recovers r = 6.0 mm.",
                     ],
                     header_row=(
                         f"{'experiment':46s}  {'noise':>6s}   "
                         f"{'powell r':>9s}   {'polish r':>9s}   {'NRMS data':>9s}"
                     ),
                     table_rows=[
                         f"{'034 stage 1, 5% noise':46s}  {'5%':>6s}   "
                         f"{'6.999 mm':>9s}   {'-':>9s}   {'5.09%':>9s}",
                         f"{'035 polish from stage 1, 5% noise':46s}  {'5%':>6s}   "
                         f"{'-':>9s}   {'6.000 mm':>9s}   {'5.02%':>9s}",
                         f"{'036 stage 1, 10% noise':46s}  {'10%':>6s}   "
                         f"{'6.927 mm':>9s}   {'-':>9s}   {'10.05%':>9s}",
                         f"{'037 polish from stage 1, 10% noise':46s}  {'10%':>6s}   "
                         f"{'-':>9s}   {'6.000 mm':>9s}   {'10.01%':>9s}",
                     ],
                     footer="Bandwidth scheduling helps move rough seeds; the polish is the radius selector.")


def build_cumulative_freq(prs):
    slide_framing(prs,
                  "Cumulative frequency — first try",
                  "Add a 1.0 GHz source carrier alongside the 1.5 GHz one, average the two",
                  hypothesis=("Adding a complementary source carrier frequency should "
                              "help radius, in the spirit of WRI's cumulative-frequency "
                              "strategy (low frequencies kept while higher ones are added)."),
                  did=("Added per-frequency objective contributions (the residual misfit "
                       "split by source centre frequency, not just the average) to the "
                       "summary, and ran a coarse grid polish at the known high-radius "
                       "seed using two source frequencies (1.0 + 1.5 GHz)."),
                  result_lines=[
                      "Polish still recovers r = 6.0 mm in this exact case.",
                      "But the per-frequency split shows that 1.0 GHz is a much weaker radius discriminator than 1.5 GHz.",
                      "Equal averaging dilutes the useful radius evidence.",
                  ],
                  exp_label="exp 038–041",
                  notes_text=(
                      "Two new pieces of plumbing were added here. First, "
                      "per-frequency reporting in the summary JSON, so we "
                      "can see which source band a candidate's misfit "
                      "actually came from. Second, a --frequency-weights "
                      "CLI flag for weighted averaging across frequencies. "
                      "These two together set up the frequency-weighting "
                      "matrix that we will see in the next section."
                  ))

    slide_figure(prs,
                 "Two-frequency polish",
                 "1.0 + 1.5 GHz objective, coarse polish from the high-radius seed",
                 os.path.join(exp_dir("040"), "figures",
                              "single_rebar_model_comparison.png"),
                 caption_body=("Polished r = 6.000 mm. Per-frequency margins: 1.0 GHz "
                               "= 3.56e-5, 1.5 GHz = 1.04e-3. The average dilutes the "
                               "1.5 GHz evidence by roughly half."),
                 exp_label="exp 040",
                 notes_text=(
                     "The recovered model panel looks the same as the "
                     "earlier polish results because the true model has "
                     "zero misfit at both frequencies. The interesting "
                     "information is in the margin numbers. At r = 6.2 "
                     "mm, the 1.0 GHz objective gives a misfit of "
                     "3.56e-5, while the 1.5 GHz objective gives "
                     "1.04e-3 — roughly a thirty-times ratio. When we "
                     "average the two equally, the effective margin "
                     "drops to about 5.4e-4, roughly halfway between. "
                     "The conclusion to take away when designing multi-"
                     "frequency objectives: equal weighting is not "
                     "neutral. It dilutes the most informative frequency."
                 ))


def build_spectrum_design(prs, fig_src, fig_obs, fig_res):
    slide_framing(prs,
                  "Designing PEBDD bands from real spectra",
                  "Choose band edges from what the data actually contains",
                  hypothesis=("The guessed low band (0.2–0.8 GHz) may have contained "
                              "almost none of the energy that distinguishes nearby "
                              "radii. Choose bands from real signal/residual spectra "
                              "instead."),
                  did=("Built a spectrum-design runner. Saved source, observed and "
                       "candidate-residual spectra for the true model, a near-radius "
                       "(r = 6.2) candidate, and the known high-radius Powell "
                       "candidate, with and without 10% noise."),
                  result_lines=[
                      "Radius-discriminating residual energy lives mostly above 1 GHz, peaking near 2 GHz.",
                      "The earlier 0.2–0.8 GHz band contained ≈ 2% of that energy.",
                      "Noise residual is broadband — useful as a stress test, not as a faithful field-noise model.",
                  ],
                  exp_label="exp 043–044",
                  notes_text=(
                      "The implementation lives in two files: "
                      "inversion/spectrum_analysis.py for the core "
                      "logic, and run_single_rebar_spectrum_design.py "
                      "for the experiment runner. The three spectra — "
                      "source, observed, and candidate residual — are "
                      "all saved as a single CSV so we can re-render at "
                      "any zoom level. The figures on the next three "
                      "slides are re-rendered with the x-axis cropped to "
                      "0–4 GHz, because the original saved figures "
                      "extend to about 250 GHz Nyquist and the relevant "
                      "content is invisible at that scale."
                  ))

    slide_figure(prs,
                 "Source wavelet spectrum",
                 "1.5 GHz Ricker pulse, FFT amplitude (Ricker = negative second derivative of a Gaussian)",
                 fig_src,
                 caption_body=("Useful source band is roughly 0.75–2.5 GHz with peak at 1.5 GHz. "
                               "x-axis cropped to 0–4 GHz for readability; the FDTD sampling rate "
                               "(~250 GHz Nyquist) sets the figure's untrimmed extent."),
                 exp_label="exp 043",
                 notes_text=(
                     "The runner reports the 5-to-95 percent source-"
                     "energy band as 0.75 to 2.5 GHz, with peak at "
                     "1.5 GHz, which matches the nominal centre frequency "
                     "of the Ricker pulse. This sets a hard upper limit "
                     "for any band design that comes next: anything we "
                     "ask the inversion objective to fit above about "
                     "2.5 GHz is fitting noise or numerical content, "
                     "not real source content. It also explains why the "
                     "W2 / Sinkhorn experiments later in the deck need "
                     "to downsample along the time axis — the source "
                     "band is narrow compared to the trace sampling "
                     "rate, so most of the trace samples carry no "
                     "source information and are pure simulator "
                     "overhead from Sinkhorn's perspective."
                 ))

    slide_figure(prs,
                 "Observed muted B-scan spectrum",
                 "Spectrum of the synthetic observed traces",
                 fig_obs,
                 caption_body=("Observed energy band 0.375–2.75 GHz, peak at 0.75 GHz. "
                               "The medium and antenna geometry have shifted energy downward "
                               "relative to the source spectrum."),
                 exp_label="exp 043",
                 notes_text=(
                     "Notice that the observed spectrum peak is at "
                     "0.75 GHz, not 1.5 GHz. That can surprise the "
                     "audience: the source carrier is 1.5 GHz, but by "
                     "the time the wave has propagated through air, "
                     "crossed the air-concrete interface, scattered off "
                     "the rebar, and returned through the Tx–Rx "
                     "geometry, the energy distribution has shifted "
                     "downward. The important caveat is that this does "
                     "not mean the radius information lives at 0.75 "
                     "GHz. The radius-discriminating part is in the "
                     "residual, not in the observed trace itself — and "
                     "the next slide shows the residual spectra "
                     "explicitly."
                 ))

    slide_figure(prs,
                 "Candidate residual spectra — the headline plot",
                 "What frequencies actually distinguish nearby radii (residual = synthetic − observed, then FFT amplitude)",
                 fig_res,
                 caption_body=("Residual energy from a near-radius candidate (r = 6.2 mm) and the "
                               "known high-radius Powell candidate (r ≈ 6.95 mm — Powell is SciPy's "
                               "derivative-free local optimiser that consistently lands here) both peak "
                               "near 2 GHz. Bands below 1 GHz carry almost none of the radius evidence."),
                 footer="This is why the first-pass 0.2–0.8 GHz schedule could not fix radius.",
                 exp_label="exp 043",
                 notes_text=(
                     "The numerical context behind this chart: for the "
                     "r = 6.2 mm near-radius candidate, the residual "
                     "energy captured by progressively wider bands is "
                     "2.1 percent in 0.2–0.8 GHz, 5.5 percent in 0.35–"
                     "1.10, 22.5 percent in 0.35–1.50, 60.4 percent in "
                     "0.35–2.00, and 89.7 percent in 0.35–2.50. So the "
                     "difference between r = 6.0 and r = 6.2 lives "
                     "almost entirely above 1 GHz, with the peak near "
                     "2 GHz. This single chart is what justified the "
                     "spectrum-derived schedule on the next slides."
                 ))


def build_faithful_pebdd(prs, fig_stage_progression):
    slide_framing(prs,
                  "Faithful spectrum-derived PEBDD schedule",
                  "Run the new bands as a staged Powell + polish",
                  hypothesis=("If the band edges were the problem, the spectrum-"
                              "derived schedule should land Powell closer to the "
                              "true radius before polish."),
                  did=("Built a staged runner using the bands 0.35–1.10 → 0.35–1.50 "
                       "→ 0.35–2.00 → 0.35–2.50 GHz, then a final full-band coarse "
                       "polish. Exact data, 5 sources, 1 mm grid."),
                  result_lines=[
                      "x and z stay on target throughout.",
                      "Each Powell stage stops in the high-radius basin (r ≈ 6.86–6.93 mm).",
                      "Final full-band polish recovers r = 6.000 mm exactly.",
                  ],
                  exp_label="exp 045",
                  notes_text=(
                      "The runner is run_single_rebar_bandwidth_schedule.py. "
                      "Each stage's outputs are saved as a nested "
                      "directory inside the experiment directory — "
                      "stage01_0.35_1.10GHz, stage02_0.35_1.50GHz, and "
                      "so on, ending with a final_fullband_coarse_polish "
                      "stage. This run is the controlled comparison to "
                      "the first-pass PEBDD: the same staged shape, but "
                      "with bands chosen by the spectrum tool rather "
                      "than by guesswork."
                  ))

    slide_figure(prs,
                 "Stage-by-stage recovered radius",
                 "Spectrum-derived PEBDD schedule",
                 fig_stage_progression,
                 caption_body=("Powell stages stay in the high-radius basin even as the "
                               "band widens to 2.5 GHz. The final polish is the step that "
                               "lands radius on truth."),
                 exp_label="exp 045",
                 notes_text=(
                     "The Powell stages move slightly away from truth as "
                     "the band widens — from 6.864 in stage 1 up to "
                     "6.930 in stage 4. The right interpretation is that "
                     "radius is already as well determined as Powell can "
                     "make it in the lowest band, and the wider bands "
                     "reactivate the depth-radius valley rather than "
                     "adding new constraints. Polish is what removes the "
                     "bias. The headline of this section is the same as "
                     "the first-pass section: bandwidth scheduling helps "
                     "with seed building, not with radius selection."
                 ))

    slide_figure(prs,
                 "Final full-band coarse polish — model comparison",
                 "Final stage of the faithful PEBDD schedule (NRMS = normalised root-mean-square misfit)",
                 os.path.join(exp_dir("045"), "final_fullband_coarse_polish",
                              "figures", "single_rebar_model_comparison.png"),
                 caption_body=("Recovered geometry matches truth on the 1 mm grid (NRMS model = 0, NRMS data = 0). "
                               "Top polish candidates (J = least-squares waveform misfit; lower is better): "
                               "r = 6.0 mm (J = 0), r = 6.2 mm (J = 1.04 × 10⁻³), r = 6.8 mm (J = 2.08 × 10⁻³)."),
                 exp_label="exp 045",
                 notes_text=(
                     "Normalised root-mean-square misfit is zero on both "
                     "the model and the data. The polish margin — the "
                     "objective gap between the chosen r = 6.0 mm and "
                     "the next-best r = 6.2 mm — is about 1.04e-3, "
                     "identical to the simpler polish from the first "
                     "arc. So the bandwidth schedule did not change the "
                     "radius evidence the polish stage operates on; it "
                     "only ensured the polish started inside the right "
                     "(x, z) window."
                 ))


def build_frequency_weighting(prs, fig_freq_compare):
    slide_section_divider(prs, 3, "Frequency weighting and W2",
                          "Where to put the spectral evidence weight, and when not to")

    slide_framing(prs,
                  "Frequency-weighted radius margins",
                  "Five weight schemes evaluated on the same local x / z / r candidate grid",
                  hypothesis=("If multi-frequency information is genuinely useful for "
                              "radius, some weighted combination should beat single-"
                              "frequency LS (least squares)."),
                  did=("Built a frequency-weight matrix runner. Swept radius 5.4–7.8 mm "
                       "at four depths under five weight schemes: low_only (1.0 GHz "
                       "alone), onepointfive_only (1.5 GHz alone), unweighted (equal "
                       "average), carry_low_25 (0.25 × 1.0 GHz + 1.0 × 1.5 GHz), "
                       "carry_low_50 (0.5 × 1.0 GHz + 1.0 × 1.5 GHz). "
                       "Ran exact, 5% noise, and 10% noise versions."),
                  result_lines=[
                      "low_only has the weakest radius margin (30x smaller than 1.5 GHz alone).",
                      "Unweighted averaging cuts the radius margin nearly in half.",
                      "carry_low_25 preserves most of the 1.5 GHz margin while keeping low-frequency carry.",
                  ],
                  exp_label="exp 046–049",
                  notes_text=(
                      "The runner is run_single_rebar_frequency_weight_"
                      "matrix.py. Weight schemes are specified as comma-"
                      "separated lists — for example, '0.25,1' means "
                      "25 percent weight on the 1.0 GHz residual and "
                      "full weight on the 1.5 GHz residual. The runner "
                      "reports the best-radius-versus-next-radius "
                      "margin per scheme, so the schemes can be ranked "
                      "at a glance from a single output table."
                  ))

    slide_figure(prs,
                 "Exact-data radius margin under five weight schemes",
                 "Common local x / z / r candidate grid evaluated for each scheme",
                 os.path.join(exp_dir("046"), "figures",
                              "frequency_weight_radius_profiles.png"),
                 caption_body=("Scheme legend — low_only: 1.0 GHz alone; onepointfive_only: 1.5 GHz alone; "
                               "unweighted: equal 1.0 + 1.5 GHz average; carry_low_25 / carry_low_50: 0.25 "
                               "(or 0.50) weight on 1.0 GHz plus full weight on 1.5 GHz. The deepest valley "
                               "at r = 6.0 mm wins on radius discrimination; flatter curves are weaker."),
                 exp_label="exp 046",
                 notes_text=(
                     "The distinct-radius margins between r = 6.0 and "
                     "r = 6.2 are: 3.56e-5 for low_only, 1.04e-3 for "
                     "onepointfive_only, 5.36e-4 for unweighted, 8.37e-4 "
                     "for carry_low_25, and 7.03e-4 for carry_low_50. "
                     "The pattern is intuitive once we have seen the "
                     "spectrum-design slide: 1.0 GHz is below the "
                     "residual-energy band that distinguishes nearby "
                     "radii, so it carries little radius information, "
                     "and equal averaging just dilutes the useful "
                     "evidence from 1.5 GHz."
                 ))

    slide_figure(prs,
                 "Same conclusion under noise — 5% and 10%",
                 "Side-by-side: 5% additive trace-RMS noise on the left, 10% on the right",
                 fig_freq_compare,
                 caption_body=("Margin = objective value at r = 6.2 mm minus value at r = 6.0 mm "
                               "(larger margin → stronger radius discrimination). The headline numbers "
                               "at 10% noise: 1.04 × 10⁻³ (1.5 GHz only), 8.4 × 10⁻⁴ (carry_low_25 = "
                               "0.25 × 1.0 GHz + 1.0 × 1.5 GHz), 5.4 × 10⁻⁴ (unweighted average), "
                               "3.1 × 10⁻⁵ (1.0 GHz only). Ranking is identical at 5%."),
                 exp_label="exp 047 + 049",
                 notes_text=(
                     "The two panels are nearly identical, which is the "
                     "point: the ranking is robust to the noise level. "
                     "The practical recommendation from this whole "
                     "matrix is: use 1.5 GHz-only least squares for the "
                     "final radius decision; use carry_low_25 if a "
                     "low-frequency contribution is needed for basin "
                     "continuity — for example, when the initial seed "
                     "is poorer and we still need some sensitivity to "
                     "broad-scale features. Never use the unweighted "
                     "equal average for radius selection; the margin "
                     "loss is consistent across noise levels."
                 ))


def build_w2(prs):
    slide_framing(prs,
                  "W2 / Optimal Transport — gate 1",
                  "Replicate the paper's shift-convexity claim on shifted Ricker pulses",
                  hypothesis=("The W2 paper claims that quadratic Wasserstein with "
                              "Softplus normalisation (log(1+e^(βx)) — turns a signed "
                              "trace into a non-negative mass distribution) and "
                              "Sinkhorn iterations (entropy-regularised cheap OT) is "
                              "monotonic over time shifts where standard L2 has "
                              "spurious local minima from oscillatory cancellation."),
                  did=("Built a standalone Softplus / Sinkhorn W2 module. Computed "
                       "both L2 and W2 over a −28 … +28 sample shift range for a "
                       "Ricker pulse, with β ∈ {4, 8, 12} (Softplus scale) and "
                       "ε = 0.02 (Sinkhorn entropy)."),
                  result_lines=[
                      "L2 has 24 monotonicity violations across the tested shift range.",
                      "W2 is smoothly monotonic for every tested β.",
                      "Paper claim reproduced — but on shifted Ricker traces, not on the rebar problem.",
                  ],
                  exp_label="exp 048",
                  notes_text=(
                      "The W2 module lives in inversion/trace_wasserstein.py "
                      "as a standalone object with unit tests. The "
                      "Softplus scale β controls how strongly the signed "
                      "waveform is transformed before being normalised "
                      "to a mass distribution; the Sinkhorn entropy "
                      "regularisation ε trades transport accuracy for "
                      "numerical stability. We test three β values — 4, "
                      "8, and 12 — as a small sensitivity sweep. We "
                      "call this gate 1 because passing the trace-shift "
                      "test is necessary but not sufficient for W2 to "
                      "be useful on the actual radius problem; gate 2 "
                      "is the rebar landscape three slides from now."
                  ))

    slide_figure(prs,
                 "W2 trace convexity — paper claim holds",
                 "Softplus / Sinkhorn quadratic Wasserstein (W2) vs least-squares (L2) over shifted Ricker pulses",
                 os.path.join(exp_dir("048"), "figures",
                              "trace_w2_convexity.png"),
                 caption_body=("L2 (left) shows the classical cycle-skip pattern: zero only at shift = 0, "
                               "with a second 'mountain' that traps optimisers. W2 (right) is smooth and "
                               "monotonic on both sides of zero for all three Softplus scales "
                               "(β = 4 / 8 / 12) — the larger β, the sharper the curvature near zero."),
                 exp_label="exp 048",
                 notes_text=(
                     "The L2 panel on the left hits zero only at the "
                     "shift-equals-zero point, but it rises and then "
                     "falls again on each side — that second mountain "
                     "is the cycle-skip trap that misleads optimisers. "
                     "The W2 panel on the right is U-shaped and "
                     "monotonic on both sides of zero for all three β "
                     "values. β = 12 has the sharpest curvature near "
                     "the minimum, which would matter for an optimiser. "
                     "We use β = 8 on the next slide for the rebar "
                     "landscape as a balance between sharpness and "
                     "numerical stability."
                 ))

    slide_framing(prs,
                  "W2 — gate 2 on the actual rebar landscape",
                  "Move quadratic Wasserstein (W2) from the shift smoke test to the radius problem",
                  hypothesis=("If W2's better convexity transfers to the rebar "
                              "landscape, its margin (objective gap) between r = 6.0 mm "
                              "and r = 6.2 mm should be comparable to or better than "
                              "the standard least-squares (LS) margin."),
                  did=("Evaluated the same Softplus / Sinkhorn W2 on the same local "
                       "x / z / r candidate grid as the LS landscape (β = 8 = Softplus "
                       "scale, ε = 0.02 = Sinkhorn entropy regularisation, ds = 8 / 16 "
                       "= downsample factor for trace length; ds reduces the Sinkhorn "
                       "cost by taking every nth sample of the trace)."),
                  result_lines=[
                      "LS margin (r=6.0 vs r=6.2): 1.04e-3.",
                      "W2 margin: 1.06e-7 (downsample 16); 1.03e-7 (downsample 8).",
                      "Softplus normalisation removes the amplitude content that radius depends on.",
                  ],
                  exp_label="exp 050–051",
                  notes_text=(
                      "The runner is run_single_rebar_w2_landscape.py. "
                      "The downsample parameter is needed because "
                      "Sinkhorn on the full trace length is expensive — "
                      "ds = 16 means we take every sixteenth sample of "
                      "the muted trace before computing W2. The key "
                      "story to land here is that the W2 paper's claim "
                      "about shifts genuinely holds — we just saw that "
                      "in gate 1 — but the radius problem in our setup "
                      "is not a shift problem, so W2's advantage does "
                      "not transfer."
                  ))

    slide_figure(prs,
                 "W2 radius margin collapses on the rebar problem",
                 "Note the W2 y-axis is in units of 1 × 10⁻⁵",
                 os.path.join(exp_dir("050"), "figures",
                              "w2_radius_profiles.png"),
                 caption_body=("L2 = least squares (left): the familiar deep V at the true radius "
                               "r = 6.0 mm. W2 = quadratic Wasserstein with Softplus normalisation "
                               "(right): nearly flat — radius margin is 1.06 × 10⁻⁷ vs 1.04 × 10⁻³ "
                               "for L2 (≈ 10 000× smaller). Softplus removed the amplitude content "
                               "that radius depends on."),
                 footer="Decision: reject W2 as the final radius objective. Keep it as a basin diagnostic for field data.",
                 exp_label="exp 050",
                 notes_text=(
                     "Both panels reach zero at r = 6.0 mm, but the L2 "
                     "plot on the left has wings on the order of 1e-3, "
                     "while the W2 plot on the right has wings on the "
                     "order of 1e-7. The downsample-equals-8 version "
                     "in experiment 051 is qualitatively identical, so "
                     "this is not a numerical artefact of how we set up "
                     "Sinkhorn. The decision is to reject W2 as the "
                     "final radius objective for this problem. We keep "
                     "the module available because it may still be "
                     "useful as a basin diagnostic for future poor "
                     "initial models or field data, but for radius "
                     "selection we stay with least squares."
                 ))


def build_material(prs):
    slide_combined(prs,
                   "Material tradeoff — does material explain radius?",
                   "Fixed x = 250, z = 90 mm; sweep radius × εr × σ (εr = relative permittivity / dielectric constant; σ = electrical conductivity, S/m)",
                   os.path.join(exp_dir("056"), "figures",
                                "material_profiled_radius.png"),
                   body_lines=[
                       "Hypothesis: at the correct (x, z) location, free material parameters could let a wrong radius still fit the data.",
                       "Sweep: r = 5.4–7.8 mm × εr ∈ {5.5, 6.0, 6.5} × log10 σ ∈ {5, 6, 7} S/m.",
                       "",
                       "Result:",
                       "  • Best material set is the true one (εr = 6.0).",
                       "  • Rebar σ saturates above ≈ 10⁵ S/m (1e5, 1e6, 1e7 all behave the same — the rebar already acts as a perfect conductor).",
                       "  • Best radius is 6.0 mm; distinct-radius margin against r = 6.2 mm is J = 1.04 × 10⁻³ (J = least-squares misfit).",
                       "",
                       "Decision: do not add material parameters to the radius optimiser at this stage.",
                   ],
                   exp_label="exp 056",
                   notes_text=(
                       "The runner is run_single_rebar_material_tradeoff.py. "
                       "The top four candidates in the saved summary "
                       "tell the whole story. Rank 1 is the true model: "
                       "r = 6.0, εr = 6.0, σ = 1e7, misfit zero. Ranks "
                       "2 and 3 are also r = 6.0 with εr = 6.0, but with "
                       "σ values of 1e6 and 1e5 — and the misfit values "
                       "are tiny, 3.18e-10 and 3.85e-8. That is the "
                       "saturation: once σ is high enough for the rebar "
                       "to behave as a perfect conductor, the exact "
                       "value stops mattering. Rank 4 is r = 6.2 with "
                       "the true εr, and the misfit is 1.04e-3 — the "
                       "familiar radius margin. The εr value is "
                       "genuinely identified; the σ value is saturated "
                       "above 1e5 S/m. That tells us we do not need to "
                       "introduce free material parameters into the "
                       "radius optimiser at this stage."
                   ))


def build_wavelet_mismatch(prs):
    slide_section_divider(prs, 4, "Wavelet mismatch — the key new finding",
                          "What breaks the pipeline first when we leave the synthetic comfort zone")

    slide_framing(prs,
                  "Source-wavelet mismatch — does radius survive?",
                  "Perturb only the observed source pulse; modelled source stays nominal",
                  hypothesis=("Real data don't have a perfectly known source pulse. If "
                              "the observed source differs from the modelled one in "
                              "amplitude, time-zero, or centre frequency (fc), does the "
                              "radius selection still hold?"),
                  did=("Built a wavelet-mismatch runner. Held the modelled source at "
                       "nominal 1.5 GHz Ricker; perturbed only the observed source "
                       "across seven cases. Case-code legend: nominal = no change; "
                       "fc_low10 / fc_high10 = source centre frequency ×0.9 or ×1.1; "
                       "delay_±50 ps = observed source shifted later (+50 ps) or "
                       "earlier (−50 ps); amp_low10 / amp_high10 = observed source "
                       "scaled by 0.9 or 1.1. Same local radius grid as the "
                       "frequency-weighting matrix."),
                  result_lines=[
                      "Pure amplitude mismatch shifts the best radius by one grid step.",
                      "Pure centre-frequency or time-delay mismatch can peg the best radius to the grid bound (5.4 or 7.8 mm).",
                      "Adding source profiling one nuisance parameter at a time (amplitude → + time → + frequency-scale) recovers the true radius.",
                  ],
                  exp_label="exp 052–055",
                  notes_text=(
                      "The runner is run_single_rebar_wavelet_mismatch.py. "
                      "The modelled source wavelet stays at the nominal "
                      "1.5 GHz Ricker; only the observed truth changes "
                      "from case to case. This is the first experiment "
                      "series in the deck where we allow any aspect of "
                      "the source to be wrong, and it is where the most "
                      "important new finding of this arc comes from. "
                      "The fix unfolds across four slides — each slide "
                      "adds one nuisance parameter to the objective and "
                      "shows which mismatch cases recover and which do "
                      "not."
                  ))

    slide_figure(prs,
                 "Raw mismatch — radius pegged to the grid bound",
                 "Fixed-source least-squares (LS) with no nuisance parameters fit per candidate",
                 os.path.join(exp_dir("052"), "figures",
                              "wavelet_mismatch_radius_profiles.png"),
                 caption_body=("Case codes: fc_low10 / fc_high10 = observed centre frequency ×0.9 / ×1.1; "
                               "delay_±50 ps = observed source shifted later / earlier by 50 ps; "
                               "amp_low10 / amp_high10 = observed amplitude ×0.9 / ×1.1. Only the nominal "
                               "and amp_low10 curves keep their minimum at r = 6.0 mm — the four fc / delay "
                               "cases push the minimum to the lower (5.4 mm) or upper (7.8 mm) radius bound."),
                 exp_label="exp 052",
                 notes_text=(
                     "The best-radius per case tells the story. Nominal "
                     "recovers 6.0 mm, fc_low10 pegs to the lower grid "
                     "bound at 5.4, fc_high10 pegs to the upper bound at "
                     "7.8, delay_+50 ps goes to 5.4, delay_−50 ps goes "
                     "to 7.8, amp_low10 stays at 6.0, and amp_high10 "
                     "moves to 7.0. The asymmetry is real: pegging to "
                     "the lower bound corresponds to a slower observed "
                     "wavelet — later arrival, longer period — while "
                     "pegging to the upper bound corresponds to a faster "
                     "one. This is the kind of failure that would only "
                     "show up after moving to real data. It is the "
                     "motivation for the three nuisance-parameter slides "
                     "that follow, which fix the cases one at a time."
                 ))

    slide_figure(prs,
                 "+ amplitude scalar fit",
                 "Fix amplitude mismatch first — fit one scalar amplitude per candidate (closed form, essentially free)",
                 os.path.join(exp_dir("053"), "figures",
                              "wavelet_mismatch_radius_profiles.png"),
                 caption_body=("Amplitude mismatch is cleanly separated from radius. "
                               "amp_high10 (observed amplitude ×1.1) returns to r = 6.0 mm; the four "
                               "fc / delay cases (centre-frequency and time-shift perturbations) are unchanged."),
                 exp_label="exp 053",
                 notes_text=(
                     "The optimal scalar amplitude for each candidate "
                     "has a closed form — it is the standard inner-"
                     "product over norm-squared ratio, the dot product "
                     "of d_obs and d_syn divided by the dot product of "
                     "d_syn with itself. That makes it essentially free "
                     "compared to the cost of a forward solve. The "
                     "amp_high10 line, where the observed source was "
                     "scaled by 1.1, now sits on top of the nominal "
                     "line. The fc_low10, fc_high10, delay_+50, and "
                     "delay_-50 lines are unchanged because amplitude "
                     "fitting cannot absorb phase or bandwidth "
                     "differences. The takeaway is that amplitude is "
                     "the smallest, cheapest nuisance parameter that "
                     "meaningfully reduces source-mismatch bias — if "
                     "the budget only allows one, this is the one to "
                     "add."
                 ))

    slide_figure(prs,
                 "+ global time-shift grid",
                 "Add a small ±80 ps time-shift profile to the nuisance fit",
                 os.path.join(exp_dir("054"), "figures",
                              "wavelet_mismatch_radius_profiles.png"),
                 caption_body=("Time-shift grid: {−80, −50, −25, 0, +25, +50, +80} ps applied globally "
                               "per scan position (not per candidate). delay_±50 ps cases land back at "
                               "r = 6.0 mm. The remaining failure is the fc_low10 / fc_high10 pair (centre-"
                               "frequency ×0.9 / ×1.1) — bandwidth mismatch is not yet absorbed."),
                 exp_label="exp 054",
                 notes_text=(
                     "The time-shift grid is small on purpose — minus "
                     "80, minus 50, minus 25, zero, plus 25, plus 50, "
                     "plus 80 picoseconds. The shift is global per "
                     "scan position, meaning one value per (scan, "
                     "candidate) pair, not a free per-trace parameter; "
                     "that keeps the nuisance space small enough to "
                     "enumerate exhaustively. The delay-plus-50 and "
                     "delay-minus-50 cases now land back at r = 6.0 mm "
                     "because a global plus 50 picosecond shift in the "
                     "modelled trace fully absorbs a minus 50 "
                     "picosecond observed mismatch. The fc_low10 case "
                     "partially recovers — the best radius moves from "
                     "the lower grid bound at 5.4 mm up to 6.2 mm — "
                     "because a time shift can partially compensate "
                     "for a longer-period observed pulse, but "
                     "fc_high10 still fails. Bandwidth mismatch is "
                     "fundamentally different from a pure time shift "
                     "and needs its own nuisance parameter, which the "
                     "next slide adds."
                 ))

    slide_figure(prs,
                 "+ frequency-scale profile — all cases recover",
                 "Add a small centre-frequency scale grid (fc_scale ∈ {0.9, 1.0, 1.1}) to the source profile",
                 os.path.join(exp_dir("055"), "figures",
                              "wavelet_mismatch_radius_profiles.png"),
                 caption_body=("Per-candidate nuisance fit is now: amplitude scalar + global time-shift + "
                               "centre-frequency scale. All seven mismatch cases collapse onto a clean V "
                               "at the true r = 6.0 mm. Distinct-radius margin (objective gap to r = 6.2 mm) "
                               "ranges from 4.55 × 10⁻⁴ (fc_low10) to 1.20 × 10⁻³ (fc_high10)."),
                 footer="This is the most important new contribution of this arc.",
                 exp_label="exp 055",
                 notes_text=(
                     "The frequency-scale grid is small on purpose — "
                     "just three values, 0.9, 1.0, and 1.1 — because "
                     "each value needs its own FDTD forward solve, so "
                     "the grid is the expensive part of the nuisance "
                     "fit. All seven mismatch cases now land at r = "
                     "6.0 mm. This slide is the diagnostic matrix that "
                     "tells us the three-piece source profile — "
                     "amplitude scalar, time-shift grid, and frequency-"
                     "scale grid — is sufficient to absorb the kinds "
                     "of mismatch we tested. The next section packages "
                     "the same three pieces into a reusable production "
                     "polish runner."
                 ))


def build_synthesis(prs):
    slide_section_divider(prs, 5, "Source-profiled polish — synthesis",
                          "Every surviving idea assembled into one runner")

    slide_framing(prs,
                  "Production source-profiled radius polish",
                  "What the recommended pipeline actually runs at the final stage",
                  hypothesis=("If the surviving ideas are right, a single runner "
                              "that does (a) local x/z/r grid, (b) per-candidate "
                              "amplitude + time-shift + frequency-scale source "
                              "profile, and (c) top-k reporting should recover the "
                              "true radius in both nominal and mismatched cases."),
                  did=("Built run_single_rebar_source_profiled_polish.py. "
                       "Validated it on a nominal smoke run and a combined-mismatch "
                       "run with injected fc = 1.1, shift = -50 ps, amplitude = 1.1."),
                  result_lines=[
                      "Nominal case: best r = 6.0 mm, margin 9.815e-4. Source profile recovers truth (fc=1.0, shift=0, amp=1.0).",
                      "Combined mismatch: best r = 6.0 mm, margin 1.146e-3. Source profile recovers fc=1.1, shift=-50 ps, amp=1.1.",
                      "Confidence is reported as top-k candidates and distinct-radius margin, not a single point estimate.",
                  ],
                  exp_label="exp 057–058",
                  notes_text=(
                      "The source-profile module lives in "
                      "inversion/source_profile.py and the runner is "
                      "Stage 1 of the post-summary plan in tracker 31. "
                      "It exists so that the next stage, which "
                      "replicates results across noise and source-"
                      "mismatch seeds, has a clean validated entry "
                      "point. One design choice on this slide worth "
                      "noting: the confidence reporting — top-k "
                      "candidates plus distinct-radius margin — is "
                      "part of the runner interface, not a separate "
                      "analysis step. The pipeline does not produce a "
                      "single point estimate; it produces a list."
                  ))

    slide_figure(prs,
                 "Nominal source-profiled polish",
                 "No injected source mismatch — confirms the runner reproduces the exact-synthetic case",
                 os.path.join(exp_dir("057"), "figures",
                              "source_profiled_radius_profile.png"),
                 caption_body=("Best r = 6.0 mm at (x = 250.0 mm, z = 90.0 mm). Recovered source profile "
                               "lands at (fc_scale = 1.0 = unchanged centre frequency, shift = 0 ps, "
                               "amp = 1.000 = unchanged amplitude). Next-distinct radius is r = 6.2 mm "
                               "with objective gap 9.815 × 10⁻⁴."),
                 exp_label="exp 057",
                 notes_text=(
                     "This is the nominal smoke run: 12 geometry "
                     "candidates evaluated with no injected source "
                     "mismatch. The source-profile parameters land at "
                     "exactly the nominal values — frequency scale 1.0, "
                     "shift 0 picoseconds, amplitude 1.000 — which is "
                     "the right behaviour. The runner should not "
                     "introduce drift when there is no mismatch to fit. "
                     "The distinct-radius margin — the objective gap "
                     "between the best radius and the best next-distinct "
                     "radius — matches the earlier polish result at "
                     "about 9.815e-4. The runner is conservative: it "
                     "does not inflate margins or pull the source "
                     "profile away from truth when no mismatch is "
                     "present. That property is important because we "
                     "want this runner to be safe to use as the default "
                     "final stage, not just as an opt-in audit step."
                 ))

    slide_figure(prs,
                 "Combined mismatch — joint recovery",
                 "Injected fc_scale = 1.1, time shift = −50 ps, amp scale = 1.1 in the observed source",
                 os.path.join(exp_dir("058"), "figures",
                              "source_profiled_radius_profile.png"),
                 caption_body=("Best r = 6.0 mm at (x = 250.0 mm, z = 90.0 mm). Recovered source profile "
                               "matches the injected one exactly: (fc_scale = 1.1, shift = −50 ps, "
                               "amp = 1.100). Distinct-radius margin (J = least-squares misfit gap "
                               "to next radius r = 6.2 mm) is 1.146 × 10⁻³."),
                 exp_label="exp 058",
                 notes_text=(
                     "This run evaluates 52 geometry candidates with "
                     "three modelled source-frequency scales each. The "
                     "observed source was deliberately mismatched in "
                     "three ways at once — centre frequency scaled by "
                     "1.1, time shifted by minus 50 picoseconds, and "
                     "amplitude scaled by 1.1. The recovered source "
                     "profile matches the injected one exactly: fc "
                     "scale 1.1, shift minus 50 picoseconds, amplitude "
                     "1.100. This is the headline result of the whole "
                     "arc — the production runner recovers both the "
                     "geometry and the injected nuisance source "
                     "parameters at the same time. Looking at the "
                     "sub-ranks, a few r = 6.2 mm candidates also "
                     "prefer a minus 25 picosecond shift, which is "
                     "exactly the kind of competing explanation top-k "
                     "reporting is meant to surface."
                 ))


def build_verdict(prs):
    """Concluding verdict: what each paper's method gave us, and what we rejected."""
    s = new_slide(prs)
    title_bar(s,
              "Verdict by paper — what worked and what did not",
              "Concrete adoption / rejection per paper, with the experiment that justifies it",
              exp_label="conclusion")

    headers = ["Paper / method", "What worked  ✓",
               "What did not work / deferred  ✗"]
    rows = [
        ("Paper 1 — WRI\n(Feng et al. 2022)",
         "Cumulative-frequency philosophy adopted as time-domain weighted LS.  "
         "carry_low_25 (25% × 1.0 GHz + 1.0 × 1.5 GHz) preserves ~80% of the "
         "1.5 GHz-only radius margin while keeping a low-frequency contribution.\n"
         "[exp 046–049]",
         "Full WRI solver itself — would require a new frequency-domain "
         "Helmholtz/Maxwell operator, complex sparse solves, and a new adjoint. "
         "Deferred; the current radius bias is already explained by source "
         "handling.\n[doc 28]"),

        ("Paper 2 — PEBDD\n(Zhou et al. 2021)",
         "Matched-filter trace band-pass (same filter on observed and synthetic) "
         "as the seed-builder stage. Spectrum-design tool to choose band edges "
         "from real residual spectra rather than guesses.\n[exp 028–045]",
         "PEBDD did not fix radius on its own. Every Powell stage stopped in the "
         "high-radius basin (r ≈ 6.86–6.93 mm). The polish stage was still the "
         "radius selector. Use PEBDD as a seed-builder only.\n[exp 031, 045]"),

        ("Paper 3 — OT-LS\n(Hunziker et al. 2025)",
         "NRCCC trace-shift diagnostic implemented and saved in every run "
         "summary. Useful as a safety check for future field-data and "
         "poor-initial-model scenarios.\n[exp 024–027]",
         "OT objective inside the optimiser was ruled out. Every wrong-radius "
         "candidate, including under 10% noise, already had NRCCC = 1.0 — the "
         "radius problem is not cycle skipping. OT brings no benefit here.\n"
         "[exp 024–026]"),

        ("Paper 4 — W2\n(Lu et al. 2024)",
         "Standalone Softplus / Sinkhorn W2 module with tests, used as a "
         "landscape gate. Paper's shift-convexity claim reproduced on shifted "
         "Ricker pulses (gate 1 passed).\n[exp 048]",
         "W2 as the final radius objective. On the rebar landscape the margin "
         "collapsed from 1.04 × 10⁻³ (LS) to 1.06 × 10⁻⁷ (W2) — Softplus "
         "removes the amplitude content radius depends on. Rejected.\n"
         "[exp 050, 051]"),

        ("Paper 5 — IFWI\n(Sun et al. 2025)",
         "Frequency-principle thinking — smooth structure before detail — "
         "already implicit in our staged pipeline (coarse → fine → polish).\n"
         "[design]",
         "Full neural implicit field — flexible neural residual would risk "
         "absorbing radius bias into a 'background correction', making radius "
         "less meaningful, not more. Deferred to multi-rebar / field data.\n"
         "[doc 29]"),

        ("Additional finding\n(beyond the 5 papers)",
         "Source-wavelet profiling — per-candidate amplitude scalar + global "
         "time-shift + centre-frequency scale grid. Recovers radius even when "
         "the modelled source differs from the observed one (all seven mismatch "
         "cases land at r = 6.0 mm).\n[exp 052–055, productionised in 057, 058]",
         "Fixed-source LS alone — bias to the radius grid bound under modest fc "
         "or time-zero mismatch. Cannot be relied on for field data without "
         "the source profile.\n[exp 052]"),
    ]

    n_rows = len(rows) + 1
    table_shape = s.shapes.add_table(
        n_rows, 3,
        Inches(0.40), Inches(1.40),
        Inches(12.55), Inches(5.70),
    )
    tbl = table_shape.table

    widths = [2.40, 5.10, 5.05]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)

    tbl.rows[0].height = Inches(0.40)
    for r in range(1, n_rows):
        tbl.rows[r].height = Inches(0.88)

    for c_idx, h in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.text = h
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)

    for r_idx, (label, worked, didnt) in enumerate(rows, start=1):
        # Label cell
        cell = tbl.cell(r_idx, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BG
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = label
        run = para.runs[0]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = NAVY
        run.font.name = FONT
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

        # "Worked" cell — green-tinted
        cell = tbl.cell(r_idx, 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PALE_GREEN
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = worked
        run = para.runs[0]
        run.font.size = Pt(10)
        run.font.color.rgb = CHARCOAL
        run.font.name = FONT
        cell.margin_left = Inches(0.10)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

        # "Did not work" cell — red-tinted
        cell = tbl.cell(r_idx, 2)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PALE_RED
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = didnt
        run = para.runs[0]
        run.font.size = Pt(10)
        run.font.color.rgb = CHARCOAL
        run.font.name = FONT
        cell.margin_left = Inches(0.10)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

    set_notes(s,
              "This is the single-slide summary of the whole arc, "
              "organised by paper. The headline sits in the bottom row "
              "— the additional finding that was not in any of the "
              "five papers: source-wavelet profiling with a small "
              "structured nuisance fit per candidate. That is the most "
              "important new contribution of this arc, and it is the "
              "recipe the recommended pipeline is built around. Above "
              "that, the five papers each gave us something we kept — "
              "frequency weighting from WRI, the matched-filter band-"
              "pass from PEBDD, the trace-shift diagnostic from OT-LS, "
              "the W2 module as a landscape gate, and the frequency-"
              "principle staging idea from IFWI. Each paper also had a "
              "piece we ruled out for this specific problem, and each "
              "rejection is justified by a specific experiment number "
              "in the table. Two of the rejections are conditional, "
              "not categorical. WRI moves back onto the table if a "
              "future setup creates an initial-model failure source "
              "handling cannot fix. IFWI moves back onto the table "
              "when explicit nuisance parameters become insufficient "
              "for multi-rebar or field data.")


def build_deferred_and_next(prs):
    slide_section_divider(prs, 6, "Deferred branches and what's next",
                          "What we are explicitly not building next, and what we are doing instead")

    slide_two_column_text(prs,
                          "What we deferred",
                          "Documented in 28 / 29 trackers; reserved for later stages",
                          left_title="WRI (full)",
                          left_body=[
                              "Needs a new frequency-domain Helmholtz / Maxwell operator",
                              "Complex sparse solves at large grid",
                              "Frequency-domain adjoint and penalty continuation",
                              "Defer reason: radius is already well explained by source handling and weighted LS",
                          ],
                          right_title="IFWI (neural implicit field)",
                          right_body=[
                              "Flexible neural residual can absorb radius bias",
                              "Needs differentiable FDTD, held-out source validation, regularisation guardrails",
                              "Defer reason: not the right tool while explicit nuisance parameters still suffice",
                              "Reserved for multi-rebar / field-data stages",
                          ],
                          exp_label="docs 28 + 29",
                          notes_text=(
                              "Both deferrals are conditional, not "
                              "categorical. WRI becomes worth "
                              "revisiting if a future use case creates "
                              "a poor initial model where time-domain "
                              "FDTD optimisation reliably fails. IFWI "
                              "becomes worth revisiting when explicit "
                              "nuisance parameters can no longer "
                              "absorb material or background "
                              "residuals — most likely with multi-"
                              "rebar geometry or field data. Both "
                              "tracker documents lay out a constrained "
                              "design for a future prototype if "
                              "either of those conditions is met."
                          ))

    slide_two_column_text(prs,
                          "Recommended pipeline",
                          "For controlled single-rebar synthetic data; transferable to field data with the source profile",
                          left_title="What the pipeline does",
                          left_body=[
                              "1. Bring x/z into the basin (staged or PEBDD path).",
                              "2. Final radius from 1.5 GHz-only LS or carry_low_25.",
                              "3. Local radius profiling / grid polish, not continuous Powell.",
                              "4. Profile source amplitude + time-shift + center-frequency scale at the final stage.",
                              "5. Report top-k candidates and distinct-radius margin.",
                          ],
                          right_title="What we are explicitly not doing next",
                          right_body=[
                              "Another broad global search.",
                              "W2 inside the optimizer for radius.",
                              "Free rebar conductivity in the radius optimizer.",
                              "Full WRI or full IFWI at this stage.",
                          ],
                          exp_label="tracker 30 summary",
                          notes_text=(
                              "The five-step pipeline is the production "
                              "recipe distilled from this whole arc. "
                              "Walking through the steps: first, use "
                              "the existing staged or PEBDD path only "
                              "to bring x and z into the right basin — "
                              "we should not expect it to fix radius. "
                              "Second, evaluate the final radius with "
                              "1.5 GHz-only or carry_low_25 — never "
                              "use unweighted multi-frequency, because "
                              "it dilutes margin. Third, the radius "
                              "decision is a local grid polish — a "
                              "small enumeration over (z, r) — not "
                              "Powell; continuous Powell sits in the "
                              "depth-radius valley. Fourth, a per-"
                              "candidate source profile combining "
                              "amplitude, time-shift, and frequency-"
                              "scale is mandatory before trusting any "
                              "radius decision on data where the "
                              "source may differ from nominal. Fifth, "
                              "confidence is the distinct-radius "
                              "margin and the top-k candidate list, "
                              "not a single point estimate. The right "
                              "column matters just as much — each "
                              "'not doing' item is a specific branch "
                              "we tested and chose not to spend the "
                              "next iteration on, with reasons in the "
                              "relevant tracker documents."
                          ))

    slide_two_column_text(prs,
                          "Next experiments",
                          "Direct follow-ons that reuse the source-profiled polish runner",
                          left_title="Replication matrix",
                          left_body=[
                              "Run source-profiled polish across several noise seeds.",
                              "Run across several source-mismatch combinations (fc, shift, amp).",
                              "Stress with offset x/z/r seeds to map the basin where the polish still works.",
                          ],
                          right_title="Scaling out",
                          right_body=[
                              "Extend to two then three rebars: same forward model, same polish stage; the global stage becomes the load-bearing piece.",
                              "Introduce a hyperbola or migration baseline for a non-FWI cross-check.",
                              "Move toward field-style data with the source profile in place.",
                          ],
                          exp_label="post-058 plan",
                          notes_text=(
                              "This is the closing slide. The choice "
                              "here is deliberate: stay narrow, with a "
                              "focused replication matrix and a "
                              "scaling-out path, rather than opening "
                              "another menu of new research branches. "
                              "The hyperbola baseline on the right is "
                              "intentionally non-FWI — having an "
                              "independent method available makes the "
                              "FWI result credible when we get to "
                              "field-data validation."
                          ))


# -----------------------------------------------------------------------------
# Driver
# -----------------------------------------------------------------------------


def generate_figures():
    print("Custom figures...")

    fig_src = fig_spectrum_clip(
        ["source_wavelet"],
        "Source wavelet spectrum — cropped to 0–4 GHz",
        "spectrum_source_cropped.png",
        fmax_ghz=4.0,
        color_overrides={"source_wavelet": "#2C5F7C"},
    )
    fig_obs = fig_spectrum_clip(
        ["observed"],
        "Observed muted B-scan spectrum — cropped to 0–4 GHz",
        "spectrum_observed_cropped.png",
        fmax_ghz=4.0,
        color_overrides={"observed": "#E8913A"},
    )
    fig_res = fig_spectrum_clip(
        ["residual_near_radius_6p2", "residual_high_radius_powell",
         "residual_high_radius_grid"],
        "Candidate residual spectra — what frequencies discriminate radius",
        "spectrum_residual_cropped.png",
        fmax_ghz=4.0,
        color_overrides={
            "residual_near_radius_6p2": "#27AE60",
            "residual_high_radius_powell": "#C0392B",
            "residual_high_radius_grid": "#E8913A",
        },
    )

    fig_first_pass_radius = fig_radius_progression_pebdd_first_pass()
    fig_stage_progression = fig_pebdd_stage_progression()
    fig_freq_compare = fig_freq_weight_noise_compare()

    return {
        "src": fig_src,
        "obs": fig_obs,
        "res": fig_res,
        "first_pass_radius": fig_first_pass_radius,
        "stage_progression": fig_stage_progression,
        "freq_compare": fig_freq_compare,
    }


def build():
    figs = generate_figures()

    print("Assembling deck...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    build_abbreviations(prs)
    build_jargon(prs)
    build_papers_intro(prs)
    build_papers(prs)
    build_papers_bridge(prs)
    build_ot_diagnostic(prs)
    build_pebdd_first_pass(prs, figs["first_pass_radius"])
    build_cumulative_freq(prs)
    build_spectrum_design(prs, figs["src"], figs["obs"], figs["res"])
    build_faithful_pebdd(prs, figs["stage_progression"])
    build_frequency_weighting(prs, figs["freq_compare"])
    build_w2(prs)
    build_material(prs)
    build_wavelet_mismatch(prs)
    build_synthesis(prs)
    build_verdict(prs)
    build_deferred_and_next(prs)

    prs.save(PPTX_PATH)
    print(f"\nSaved: {PPTX_PATH}")


if __name__ == "__main__":
    build()
