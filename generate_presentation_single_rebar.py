"""
Generate a focused PowerPoint deck for the single-rebar GPR-FDTD inversion pipeline.

The deck is short and result-oriented:
  1. Problem framing and identifiability hierarchy.
  2. Forward model and twin experiment.
  3. Staged inversion architecture and the reason each stage exists.
  4. Final synthetic-recovery result, noise robustness, and ambiguity margin.

Output:
    outputs/GPR_FDTD_FWI_SingleRebar_Pipeline.pptx

It also writes a few custom figures it needs into:
    outputs/presentation_figures_single_rebar/
"""
import json
import os
import sys

import numpy as np
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# -----------------------------------------------------------------------------
# Project paths
# -----------------------------------------------------------------------------

PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPUTS_DIR = os.path.join(PROJECT_ROOT, "outputs")
EXPERIMENTS_DIR = os.path.join(OUTPUTS_DIR, "experiments")
FIG_DIR = os.path.join(OUTPUTS_DIR, "presentation_figures_single_rebar")
PPTX_PATH = os.path.join(OUTPUTS_DIR, "GPR_FDTD_FWI_SingleRebar_Pipeline.pptx")

os.makedirs(FIG_DIR, exist_ok=True)

# -----------------------------------------------------------------------------
# Design language reused from generate_presentation_v2 for visual consistency
# -----------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(0x1A, 0x23, 0x32)
STEEL = RGBColor(0x2C, 0x5F, 0x7C)
TEAL = RGBColor(0x00, 0x7C, 0x7C)
AMBER = RGBColor(0xE8, 0x91, 0x3A)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
CHARCOAL = RGBColor(0x2C, 0x3E, 0x50)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
PALE_AMBER = RGBColor(0xFD, 0xF1, 0xE0)
PALE_GREEN = RGBColor(0xE6, 0xF5, 0xEC)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

plt.rcParams.update(
    {
        "font.family": "sans-serif",
        "font.sans-serif": ["Helvetica", "Arial", "DejaVu Sans"],
        "font.size": 12,
        "axes.labelsize": 13,
        "axes.titlesize": 15,
        "axes.titleweight": "bold",
        "figure.facecolor": "white",
        "axes.facecolor": "white",
        "axes.edgecolor": "#CCCCCC",
        "grid.alpha": 0.3,
        "grid.color": "#CCCCCC",
    }
)

SLIDE_NUM = [0]


# -----------------------------------------------------------------------------
# Slide helpers
# -----------------------------------------------------------------------------


def new_slide(prs, bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    if bg_color is not None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
    return slide


def add_accent_bar(slide, color=STEEL, height=Inches(0.06)):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_slide_number(slide):
    box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.45),
        Inches(0.6), Inches(0.3),
    )
    p = box.text_frame.paragraphs[0]
    p.text = str(SLIDE_NUM[0])
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title_text, subtitle_text=None, left=Inches(0.7)):
    add_accent_bar(slide, STEEL)
    box = slide.shapes.add_textbox(left, Inches(0.22), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_TITLE

    bar = slide.shapes.add_shape(1, left, Inches(0.88), Inches(2.4), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()

    if subtitle_text:
        box2 = slide.shapes.add_textbox(left, Inches(0.95), Inches(12), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.font.name = FONT_BODY

    add_slide_number(slide)


def add_text(slide, left, top, width, height, text, size=16, bold=False,
             color=None, alignment=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color or CHARCOAL
    p.font.name = FONT_BODY
    p.alignment = alignment
    return tf


def add_bullets(slide, left, top, width, height, items, size=16, color=None, spacing=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        indent = 0
        clean = item
        while clean.startswith("  "):
            indent += 1
            clean = clean[2:]
        p.text = clean
        p.level = indent
        p.font.size = Pt(size)
        p.font.color.rgb = color or CHARCOAL
        p.font.name = FONT_BODY
        p.space_after = Pt(spacing)
    return tf


def add_figure(slide, fig_path, left, top, width=None, height=None):
    if not os.path.exists(fig_path):
        print(f"  WARNING: figure not found: {fig_path}")
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(fig_path, left, top, **kwargs)


def add_caption(slide, left, top, width, text):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = GRAY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER


def add_panel(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(5, left, top, width, height)  # rounded rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.25)
    shape.text_frame.text = ""
    return shape


def add_metric_pill(slide, left, top, label, value, color=AMBER, width=Inches(2.6)):
    """A two-line metric: big number/label on top, descriptor below."""
    add_panel(slide, left, top, width, Inches(1.2), PALE_BLUE, STEEL)
    add_text(slide, left, top + Inches(0.10), width, Inches(0.5), value,
             size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.65), width, Inches(0.5), label,
             size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


def set_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


# -----------------------------------------------------------------------------
# Custom figures
# -----------------------------------------------------------------------------


def fig_pipeline_architecture():
    """Three-stage staged solver, with role of each stage. Cleaner layout."""
    fig, ax = plt.subplots(1, 1, figsize=(13, 4.8))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 5.0)
    ax.axis("off")

    # Forward-model bar across the top — two lines of text fit comfortably
    fwd_box = FancyBboxPatch(
        (0.4, 3.50), 12.2, 1.20,
        boxstyle="round,pad=0.05",
        facecolor="#E3F2FD", edgecolor="#2C5F7C", linewidth=2.0,
    )
    ax.add_patch(fwd_box)
    ax.text(6.5, 4.30,
            "GPU FDTD forward model",
            ha="center", va="center", fontsize=15, fontweight="bold",
            color="#1A2332")
    ax.text(6.5, 3.80,
            "Yee 2D TMz Maxwell solver  ·  CPML boundaries  ·  batched B-scan"
            "  ·  shared by all three stages",
            ha="center", va="center", fontsize=11, color="#2C3E50",
            style="italic")

    # Three pipeline boxes underneath
    stages = [
        ("STAGE 1", "Coarse global search",
         "2 mm grid  ·  Differential evolution\nFinds the (x, z, r) basin",
         "#FDF1E0", "#E8913A"),
        ("STAGE 2", "Fine continuous refinement",
         "1 mm grid  ·  Powell local search\nLocks x and z; biases radius high",
         "#E6F5EC", "#27AE60"),
        ("STAGE 3", "Deterministic grid polish",
         "1 mm grid  ·  Local (z, r) enumeration\nRemoves the radius bias",
         "#EBF5FB", "#2C5F7C"),
    ]
    x_centers = [2.30, 6.50, 10.70]
    box_w, box_h = 3.70, 2.40
    bottom = 0.30
    for cx, (label, subtitle, body, fc, ec) in zip(x_centers, stages):
        left = cx - box_w / 2
        rect = FancyBboxPatch(
            (left, bottom), box_w, box_h,
            boxstyle="round,pad=0.08",
            facecolor=fc, edgecolor=ec, linewidth=2.0,
        )
        ax.add_patch(rect)
        # Section label (small, coloured to match border)
        ax.text(cx, bottom + box_h - 0.40, label,
                ha="center", va="center", fontsize=11,
                fontweight="bold", color=ec)
        # Role
        ax.text(cx, bottom + box_h - 0.95, subtitle,
                ha="center", va="center", fontsize=13,
                fontweight="bold", color="#1A2332")
        # Body
        ax.text(cx, bottom + 0.55, body,
                ha="center", va="center", fontsize=11, color="#2C3E50")

    # Forward arrows between stages (clean, single arrow each)
    for i in range(2):
        x_start = x_centers[i] + box_w / 2 + 0.06
        x_end = x_centers[i + 1] - box_w / 2 - 0.06
        ax.annotate(
            "", xy=(x_end, bottom + box_h / 2),
            xytext=(x_start, bottom + box_h / 2),
            arrowprops=dict(arrowstyle="-|>", color="#2C5F7C", lw=2.5),
        )
        ax.text((x_start + x_end) / 2, bottom + box_h / 2 + 0.20, "seed",
                ha="center", fontsize=10, color="#7F8C8D", style="italic")

    out = os.path.join(FIG_DIR, "pipeline_architecture.png")
    plt.tight_layout()
    plt.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close()
    print(f"  wrote {out}")
    return out


def fig_noise_robustness():
    """Recovered radius and data NRMS across noise levels."""
    rows = [
        ("0%, exact", 0.00, 6.0, 0.0),
        ("1%, seed 13", 0.01, 6.0, 0.0101),
        ("5%, seed 13", 0.05, 6.0, 0.0502),
        ("5%, seed 21", 0.05, 6.0, 0.0498),
        ("10%, seed 13", 0.10, 6.0, 0.1001),
    ]
    labels = [r[0] for r in rows]
    radii = [r[2] for r in rows]
    nrms = [r[3] for r in rows]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.4))

    bars = ax1.bar(labels, radii, color="#2C5F7C", edgecolor="#1A2332")
    ax1.axhline(6.0, color="#27AE60", linestyle="--", linewidth=1.5, label="true r = 6.0 mm")
    ax1.set_ylim(0, 8)
    ax1.set_ylabel("Recovered radius [mm]")
    ax1.set_title("Recovered radius vs noise level", color="#1A2332")
    ax1.legend(loc="lower right", fontsize=10)
    for bar, val in zip(bars, radii):
        ax1.text(bar.get_x() + bar.get_width() / 2, val + 0.15,
                 f"{val:.1f}", ha="center", fontsize=11, color="#1A2332",
                 fontweight="bold")
    ax1.tick_params(axis="x", labelrotation=20)

    ax2.bar(labels, nrms, color="#E8913A", edgecolor="#7A4A1A")
    ax2.set_ylabel("Data NRMS")
    ax2.set_title("Data NRMS tracks injected noise floor", color="#1A2332")
    for x, val, n in zip(range(len(rows)), nrms, [r[1] for r in rows]):
        ax2.text(x, val + 0.004, f"{val*100:.1f}%", ha="center",
                 fontsize=10, color="#1A2332", fontweight="bold")
    ax2.set_ylim(0, max(nrms) * 1.4 + 0.005)
    ax2.tick_params(axis="x", labelrotation=20)

    out = os.path.join(FIG_DIR, "noise_robustness.png")
    plt.tight_layout()
    plt.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close()
    print(f"  wrote {out}")
    return out


def fig_top_k_margin():
    """Top polish candidates at 10% noise from run 023."""
    with open(os.path.join(EXPERIMENTS_DIR,
                           "023_single_rebar_grid1mm_noise10_coarsepolish_topk",
                           "data", "single_rebar_summary.json")) as f:
        summary = json.load(f)
    top = summary["grid_polish"]["top_candidates"]
    labels = [f"r={c['params']['radius_mm']:.1f}\nz={c['params']['z_mm']:.1f}" for c in top]
    misfits = [c["misfit"] for c in top]
    best = min(misfits)
    worst = max(misfits)
    floor = best - (worst - best) * 0.15

    fig, ax = plt.subplots(1, 1, figsize=(11, 4.6))
    colors = ["#27AE60" if c["params"]["radius_mm"] == 6.0 else "#C0392B" for c in top]
    bars = ax.bar(labels, misfits, color=colors, edgecolor="#1A2332")
    for bar, m in zip(bars, misfits):
        ax.text(bar.get_x() + bar.get_width() / 2,
                m + (worst - best) * 0.04,
                f"J={m:.4f}", ha="center", fontsize=10,
                color="#1A2332")
    ax.set_ylim(floor, worst + (worst - best) * 0.35)
    ax.axhline(best, color="#27AE60", linestyle="--", linewidth=1.0, alpha=0.6)
    ax.text(len(labels) - 0.5, best + (worst - best) * 0.01,
            f"winning J = {best:.4f}",
            color="#27AE60", fontsize=10, ha="right", va="bottom")
    ax.set_ylabel("Objective J")
    ax.set_title(
        "Top polish candidates at 10% additive noise (true radius wins by ≈ 5.6e-4)",
        color="#1A2332",
    )
    ax.grid(True, axis="y", alpha=0.3)

    from matplotlib.patches import Patch
    legend = [
        Patch(facecolor="#27AE60", label="true radius (6.0 mm)"),
        Patch(facecolor="#C0392B", label="competing radius"),
    ]
    ax.legend(handles=legend, loc="upper left", fontsize=10)

    out = os.path.join(FIG_DIR, "top_k_margin.png")
    plt.tight_layout()
    plt.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close()
    print(f"  wrote {out}")
    return out


def fig_radius_sweep_pair():
    """Side-by-side radius sweeps for 2 mm vs 1 mm grid (re-plotted from CSVs)."""
    import csv

    def load(path):
        r, j = [], []
        with open(path) as f:
            reader = csv.DictReader(f)
            for row in reader:
                r.append(float(row["radius_mm"]))
                j.append(float(row["objective"]))
        idx = np.argsort(r)
        return np.array(r)[idx], np.array(j)[idx]

    r2, j2 = load(os.path.join(EXPERIMENTS_DIR,
                                "002_radius_grid2mm_probe",
                                "data", "radius_sweep.csv"))
    r1, j1 = load(os.path.join(EXPERIMENTS_DIR,
                                "003_radius_grid1mm_probe",
                                "data", "radius_sweep.csv"))

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 4.6))

    floor = max(min(j2[j2 > 0]), min(j1[j1 > 0])) * 0.1
    j2p = np.maximum(j2, floor)
    j1p = np.maximum(j1, floor)

    ax1.semilogy(r2, j2p, marker="o", color="#C0392B", linewidth=2.0)
    ax1.axvline(6.0, color="#1A2332", linestyle="--", label="true r = 6.0 mm")
    ax1.set_xlabel("Radius [mm]")
    ax1.set_ylabel("Objective J (log scale)")
    ax1.set_title("2 mm grid: radius plateaus from rasterization", color="#1A2332")
    ax1.legend(loc="upper right", fontsize=10)
    ax1.grid(True, alpha=0.3)

    ax2.semilogy(r1, j1p, marker="o", color="#27AE60", linewidth=2.0)
    ax2.axvline(6.0, color="#1A2332", linestyle="--", label="true r = 6.0 mm")
    ax2.set_xlabel("Radius [mm]")
    ax2.set_ylabel("Objective J (log scale)")
    ax2.set_title("1 mm grid: clean basin around the true radius", color="#1A2332")
    ax2.legend(loc="upper right", fontsize=10)
    ax2.grid(True, alpha=0.3)

    out = os.path.join(FIG_DIR, "radius_sweep_pair.png")
    plt.tight_layout()
    plt.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close()
    print(f"  wrote {out}")
    return out


# -----------------------------------------------------------------------------
# Slides
# -----------------------------------------------------------------------------


def slide_title(prs):
    slide = new_slide(prs, bg_color=NAVY)

    # Accent strip
    add_accent_bar(slide, AMBER, Inches(0.10))

    # Pre-title chip
    chip = slide.shapes.add_shape(5, Inches(0.9), Inches(1.6), Inches(2.6), Inches(0.35))
    chip.fill.solid()
    chip.fill.fore_color.rgb = AMBER
    chip.line.fill.background()
    add_text(slide, Inches(0.9), Inches(1.6), Inches(2.6), Inches(0.35),
             "SINGLE REBAR  •  TWIN EXPERIMENT", size=12, bold=True,
             color=NAVY, alignment=PP_ALIGN.CENTER)

    # Title
    add_text(slide, Inches(0.85), Inches(2.1), Inches(12), Inches(1.4),
             "Single-Rebar Geometry Inversion",
             size=46, bold=True, color=WHITE)
    add_text(slide, Inches(0.85), Inches(3.0), Inches(12), Inches(1.0),
             "Recovering position, depth, and radius from synthetic GPR B-scans",
             size=22, color=RGBColor(0xCC, 0xDD, 0xEE), italic=True)

    # Three small highlight tiles
    tile_y = Inches(4.5)
    tiles = [
        ("GPU FDTD", "2D TMz + CPML, batched"),
        ("Staged solver", "Coarse → fine → polish"),
        ("Verified", "Exact recovery, robust to 10% noise"),
    ]
    for i, (label, body) in enumerate(tiles):
        left = Inches(0.85 + i * 4.0)
        add_panel(slide, left, tile_y, Inches(3.7), Inches(1.4),
                  RGBColor(0x24, 0x3B, 0x5E), border_color=AMBER)
        add_text(slide, left, tile_y + Inches(0.18), Inches(3.7), Inches(0.5),
                 label, size=18, bold=True, color=AMBER, alignment=PP_ALIGN.CENTER)
        add_text(slide, left, tile_y + Inches(0.75), Inches(3.7), Inches(0.5),
                 body, size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.85), SLIDE_HEIGHT - Inches(0.6), Inches(12), Inches(0.3),
             "GPR-FDTD-FWI  •  DGX Spark  •  2026",
             size=12, color=GRAY)
    add_slide_number(slide)
    set_notes(slide,
              "This deck is about recovering the position, depth, and radius of a "
              "single rebar in concrete from a ground-penetrating radar B-scan. "
              "Everything is synthetic — observed data and candidate data come "
              "from the same wave-equation solver, so any failure we see is in "
              "the inverse pipeline, not the physics. The story has three pieces. "
              "First, the forward model is a GPU-accelerated 2D finite-difference "
              "time-domain solver with absorbing boundaries, batched across scan "
              "positions. Second, the inversion is a staged solver — coarse "
              "global search, then fine continuous refinement, then a "
              "deterministic local grid polish. Third, the result is exact "
              "recovery on clean data, and the polished radius stays correct "
              "under additive noise up to ten percent. Where the radius starts "
              "becoming ambiguous, the pipeline reports a top-k list rather "
              "than a single point estimate.")


def slide_problem(prs):
    slide = new_slide(prs)
    add_title(slide,
              "The problem",
              "What we are estimating and what is actually hard")

    add_bullets(slide, Inches(0.7), Inches(1.4), Inches(7.5), Inches(4.5),
                [
                    "Unknown parameter vector: [x, z, radius] for a single circular rebar in concrete.",
                    "Observed data: a 2D B-scan (depth-time × lateral position) from common-offset GPR.",
                    "Twin experiment: synthetic-to-synthetic. The 'observed' B-scan and the candidate B-scans share the same forward operator, so any failure is in the inverse pipeline, not the physics.",
                    "Identifiability hierarchy (literature-consistent):",
                    "  Lateral position: easy, set by hyperbola apex.",
                    "  Cover depth: easy once velocity is known.",
                    "  Radius / diameter: weakly identified; nuisance-coupled and grid-sensitive.",
                ], size=15)

    # Right side: a parameter card
    card_left = Inches(8.6)
    card_top = Inches(1.4)
    add_panel(slide, card_left, card_top, Inches(4.1), Inches(4.6),
              LIGHT_BG, STEEL)
    add_text(slide, card_left, card_top + Inches(0.15), Inches(4.1), Inches(0.5),
             "Synthetic truth", size=16, bold=True, color=NAVY,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, card_left, card_top + Inches(0.65), Inches(4.1), Inches(0.4),
             "x = 250.0 mm   z = 90.0 mm   r = 6.0 mm",
             size=14, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
    add_text(slide, card_left, card_top + Inches(1.25), Inches(4.1), Inches(0.5),
             "Initial guess (intentionally wrong)",
             size=14, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text(slide, card_left, card_top + Inches(1.70), Inches(4.1), Inches(0.4),
             "x = 235.0 mm   z = 80.0 mm   r = 8.0 mm",
             size=14, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
    add_text(slide, card_left, card_top + Inches(2.35), Inches(4.1), Inches(0.5),
             "Acquisition", size=14, bold=True, color=NAVY,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, card_left, card_top + Inches(2.85), Inches(4.1), Inches(1.6),
             "1.5 GHz Ricker source\n5–9 scan positions, 20 mm Tx–Rx offset\n"
             "Concrete εr = 6.0, σ = 0.01 S/m\nAir layer above slab",
             size=13, color=CHARCOAL, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "The unknowns are intentionally tiny. We are not estimating a full "
              "subsurface image — only three numbers from a B-scan: lateral "
              "position, depth, and radius of one circular rebar. That keeps the "
              "inverse problem clean enough to diagnose identifiability. The "
              "twin-experiment design — observed and candidate B-scans come from "
              "the same forward operator — means any wrong answer must come from "
              "the inverse pipeline. The identifiability hierarchy on the left "
              "is consistent with the GPR-FWI literature: lateral position is "
              "set by the hyperbola apex, depth is set by the moveout once "
              "velocity is known, and radius is the weakly identified quantity "
              "that couples to nuisance parameters and to the grid. The true "
              "model on the right is what the observed B-scan is synthesised "
              "from; the initial guess is deliberately offset so the optimiser "
              "has real work to do.")


def slide_pipeline_arch(prs, fig_path):
    slide = new_slide(prs)
    add_title(slide,
              "Pipeline architecture",
              "Same forward model, three optimizers staged from coarse to fine")
    add_figure(slide, fig_path, Inches(0.4), Inches(1.4), width=Inches(12.5))
    add_text(slide, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
             "Each stage hands its best estimate to the next. The forward model is identical; "
             "what changes is the resolution and the search strategy.",
             size=13, italic=True, color=GRAY)
    set_notes(slide,
              "The same forward model sits at the top, used by all three "
              "stages — what changes from stage to stage is the search "
              "strategy and the grid resolution. Stage one is a coarse global "
              "search on a 2 mm grid using differential evolution. Its job is "
              "to find the right basin in (x, z, r); it does not need to be "
              "accurate. Stage two takes that seed and runs a fine continuous "
              "local search on the production 1 mm grid using SciPy's "
              "derivative-free Powell optimiser. It locks in lateral position "
              "and depth, but it consistently biases the radius slightly high "
              "because of a depth–radius coupling we will see in a moment. "
              "Stage three is a deterministic small grid search over depth "
              "and radius on the 1 mm grid. It is the step that closes the "
              "radius gap by accepting that the rasterised geometry is "
              "piecewise constant in physical coordinates.")


def slide_forward(prs):
    slide = new_slide(prs)
    add_title(slide,
              "Forward model",
              "GPU-batched 2D TMz FDTD with CPML; parity-tested against the CPU reference")

    add_bullets(slide, Inches(0.7), Inches(1.4), Inches(6.5), Inches(5),
                [
                    "Solver: 2D transverse-magnetic (TMz) Yee-grid FDTD with CPML boundaries.",
                    "Backend: CuPy on the DGX Spark GB10. Batched B-scan support advances every scan position in one GPU launch.",
                    "Parity tests cover the single-trace and the batched paths against the CPU CPML solver, so changes in CuPy or kernels can be caught.",
                    "Numbered-run convention: every diagnostic and inversion writes a run_manifest.json with the command, git commit, backend, and summary path.",
                ], size=15)

    bscan = os.path.join(EXPERIMENTS_DIR,
                        "015_single_rebar_grid1mm_powell_gridpolish_from_2mm",
                        "figures", "single_rebar_observed_bscan.png")
    add_figure(slide, bscan, Inches(7.7), Inches(1.3), width=Inches(5.4))
    add_caption(slide, Inches(7.7), Inches(5.9), Inches(5.4),
                "Synthetic 'observed' B-scan: rebar hyperbola near 2 ns, 5 scan positions, 1.5 GHz.")
    set_notes(slide,
              "The forward model is a 2D transverse-magnetic finite-difference "
              "time-domain solver on a Yee staggered grid, with convolutional "
              "perfectly matched layer boundaries to absorb outgoing waves. It "
              "runs on the GPU through CuPy. The important property for "
              "inversion runtime is that the batched B-scan path advances "
              "every scan position in one GPU launch, so a single objective "
              "evaluation costs roughly the same as a single forward solve. "
              "Every change to the solver is parity-tested against the CPU "
              "reference for both the single-trace and the batched paths, so "
              "GPU regressions are caught early. Every inversion or "
              "diagnostic run writes a manifest with the exact command, git "
              "commit, backend, and summary path, so the experiment log is "
              "reproducible. On the right is one observed B-scan with five "
              "scan positions; the rebar hyperbola sits near 2 nanoseconds.")


def slide_radius_hard(prs, sweep_pair_path):
    slide = new_slide(prs)
    add_title(slide,
              "Why radius is hard",
              "Hard-grid rasterization, not lack of signal")
    add_figure(slide, sweep_pair_path, Inches(0.4), Inches(1.3), width=Inches(12.6))

    add_text(slide, Inches(0.7), Inches(5.6), Inches(12), Inches(1.0),
             "Same waveform objective, same scan, same true radius — only the grid step changed. "
             "The 2 mm grid produces flat plateaus where nearby radii rasterize to the same "
             "material mask. The 1 mm grid breaks the plateaus and reveals a clean basin at "
             "r = 6.0 mm.",
             size=14, color=CHARCOAL)
    add_text(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
             "Significance: radius identifiability in this code is dominated by geometry "
             "discretization. Multi-frequency stacking did not fix it; a finer grid did.",
             size=14, color=AMBER, bold=True)
    set_notes(slide,
              "Both panels show the inversion objective as a function of "
              "radius, on a log scale. The only thing that changes between "
              "the left and right panels is the grid step of the forward "
              "model. On the 2 mm grid, the radius values 4.5, 5.0, and 5.5 "
              "all sit on a flat plateau because the circular rebar "
              "rasterises to the same set of cells for all three — the "
              "optimiser literally cannot distinguish them. The 1 mm grid "
              "breaks those plateaus and gives us a clean V-shape with the "
              "true 6.0 mm radius at the bottom. The implication is "
              "important: radius identifiability in this code is dominated "
              "by geometry discretisation. We tried multi-frequency "
              "stacking first; it did not fix the plateau. A finer grid "
              "did. That is why the production pipeline runs at 1 mm even "
              "though the seed stage uses 2 mm for speed.")


def slide_coupling(prs):
    slide = new_slide(prs)
    add_title(slide,
              "Depth–radius coupling",
              "Why the continuous optimizer still drifts after the grid is fine enough")

    # Three info cards
    cards = [
        ("At Powell's stop (run 010)",
         "x = 249.5 mm   z = 90.65 mm\nbest radius ≈ 6.80 mm\nJ ≈ 2.1e-3",
         PALE_AMBER, AMBER),
        ("Same x, true depth (run 011)",
         "x = 249.5 mm   z = 90.00 mm\nbest radius = 6.00 mm\nJ ≈ 0",
         PALE_GREEN, GREEN),
        ("More sources, same z bias (run 013)",
         "9 sources, x = 249.5 mm, z = 90.65 mm\nbest radius ≈ 6.80 mm\n9-src didn't move the radius",
         PALE_BLUE, STEEL),
    ]
    for i, (title, body, fill, border) in enumerate(cards):
        left = Inches(0.7 + i * 4.2)
        top = Inches(1.4)
        add_panel(slide, left, top, Inches(4.0), Inches(2.2), fill, border)
        add_text(slide, left, top + Inches(0.15), Inches(4.0), Inches(0.5),
                 title, size=14, bold=True, color=NAVY,
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(0.75), Inches(4.0), Inches(1.3),
                 body, size=13, color=CHARCOAL, alignment=PP_ALIGN.CENTER)

    add_bullets(slide, Inches(0.7), Inches(4.0), Inches(12), Inches(3),
                [
                    "Continuous Powell stops in a low-misfit valley where a slightly deeper rebar trades against a slightly larger radius.",
                    "Profiling at the true depth recovers r = 6 mm exactly — so the signal carries radius information.",
                    "Adding scan positions narrows the valley but does not move its radius minimum. The fix has to come from the search strategy, not from more data.",
                ], size=15)
    set_notes(slide,
              "This is the central technical insight of the deck. The continuous "
              "Powell optimiser stops at radius about 6.8 mm and depth about "
              "0.65 mm deeper than truth, with a small but non-zero misfit. "
              "When we fix x and z and profile radius alone, the right answer "
              "is 6.0 mm exactly. When we fix x and the wrong depth and "
              "profile radius, we get the same 6.8 mm Powell found. So a "
              "deeper rebar trades against a larger radius and produces a "
              "near-identical waveform — the objective near the optimum is "
              "an elongated valley, not a clean basin. Adding more scan "
              "positions narrows the valley but does not move its radius "
              "minimum, because the trade-off is intrinsic to the physics at "
              "the wavelengths we are using. The remedy cannot come from "
              "more data; it has to come from the search strategy. That is "
              "what motivates stage 3.")


def slide_grid_polish(prs):
    slide = new_slide(prs)
    add_title(slide,
              "Stage 3 — the grid polish",
              "Accepting that geometry is piecewise constant on the grid")

    add_bullets(slide, Inches(0.7), Inches(1.4), Inches(6.5), Inches(5),
                [
                    "Powell's continuous step doesn't help past a point: rasterized geometry is piecewise constant in physical coordinates.",
                    "Polish takes the Powell seed and evaluates a small set of (z, r) candidates on an absolute millimeter grid.",
                    "Defaults: 40-candidate coarse grid (z step 0.5 mm, r step 0.2 mm). A 160-candidate fine preset is kept for audit runs.",
                    "Polish exposes the next-best candidates, so the margin between the chosen rebar and its competitors is auditable.",
                ], size=15)

    # Right side: result of polish
    card_left = Inches(7.7)
    card_top = Inches(1.4)
    add_panel(slide, card_left, card_top, Inches(5.0), Inches(2.0),
              PALE_AMBER, AMBER)
    add_text(slide, card_left, card_top + Inches(0.15), Inches(5.0), Inches(0.4),
             "Powell hand-off", size=14, bold=True, color=NAVY,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, card_left, card_top + Inches(0.7), Inches(5.0), Inches(1.1),
             "x = 249.5 mm   z = 90.65 mm\nr = 6.96 mm   J = 2.08e-3",
             size=14, color=CHARCOAL, alignment=PP_ALIGN.CENTER)

    add_panel(slide, card_left, card_top + Inches(2.3), Inches(5.0), Inches(2.0),
              PALE_GREEN, GREEN)
    add_text(slide, card_left, card_top + Inches(2.45), Inches(5.0), Inches(0.4),
             "After polish (run 015)", size=14, bold=True, color=NAVY,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, card_left, card_top + Inches(3.0), Inches(5.0), Inches(1.4),
             "x = 250.0 mm   z = 89.75 mm\nr = 6.00 mm   J = 0",
             size=14, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
    set_notes(slide,
              "The grid polish is the single most important step for radius "
              "in this pipeline. It is conceptually simple: take Powell's "
              "approximate answer, build a small list of nearby (depth, "
              "radius) candidates on an absolute millimetre grid, run the "
              "forward solver on each, and pick the best. The default is a "
              "40-candidate coarse grid with depth step 0.5 mm and radius "
              "step 0.2 mm; an 160-candidate fine preset exists for audit "
              "runs. The polish also keeps the top few candidates, not just "
              "the winner, so the margin between the chosen geometry and "
              "its closest competitor is auditable. On the right, Powell "
              "hands off radius 6.96 mm with a small but non-zero misfit. "
              "After polish, we land at radius 6.00 mm with misfit zero. "
              "The depth reads 89.75 mm rather than 90.00 mm because at "
              "this 1 mm grid resolution the rasterised circle is identical "
              "for several sub-millimetre depth values — that is a "
              "representation detail, not a result detail.")


def slide_final_model(prs):
    slide = new_slide(prs)
    add_title(slide,
              "Final inversion (exact synthetic)",
              "Truth = recovered on the 1 mm grid; model and data misfit both zero")

    model_fig = os.path.join(EXPERIMENTS_DIR,
                             "015_single_rebar_grid1mm_powell_gridpolish_from_2mm",
                             "figures", "single_rebar_model_comparison.png")
    add_figure(slide, model_fig, Inches(0.3), Inches(1.3), width=Inches(13.0))

    # Metric pills below the figure
    metric_y = Inches(5.3)
    add_metric_pill(slide, Inches(0.6), metric_y, "x recovery",
                    "250.0 mm", color=GREEN, width=Inches(2.6))
    add_metric_pill(slide, Inches(3.4), metric_y, "z recovery",
                    "89.75 mm", color=GREEN, width=Inches(2.6))
    add_metric_pill(slide, Inches(6.2), metric_y, "radius",
                    "6.00 mm", color=GREEN, width=Inches(2.6))
    add_metric_pill(slide, Inches(9.0), metric_y, "NRMS (model / data)",
                    "0  /  0", color=GREEN, width=Inches(3.8))

    set_notes(slide,
              "This is the exact-synthetic final result. The three model "
              "panels are the initial guess on the left, the recovered "
              "model in the middle, and the ground truth on the right. The "
              "rebar appears as a small high-permittivity inclusion in the "
              "concrete; the recovered and ground-truth panels are visually "
              "identical because, on the 1 mm grid, they are the same "
              "model. The metric tiles below confirm: lateral position 250 "
              "mm exactly, depth 89.75 mm (which is the same rasterised "
              "cell as 90 mm), radius 6.0 mm exactly, normalised root mean "
              "square misfit zero for both the model and the data. The "
              "depth reads 89.75 rather than 90.00 because at 1 mm grid "
              "resolution several sub-millimetre depth values produce the "
              "same rasterised circle. That is a representation detail, "
              "not a science result.")


def slide_bscan_compare(prs):
    slide = new_slide(prs)
    add_title(slide,
              "Observed vs recovered B-scans",
              "The recovered geometry reproduces the data trace-for-trace")
    obs = os.path.join(EXPERIMENTS_DIR,
                      "015_single_rebar_grid1mm_powell_gridpolish_from_2mm",
                      "figures", "single_rebar_observed_bscan.png")
    rec = os.path.join(EXPERIMENTS_DIR,
                      "015_single_rebar_grid1mm_powell_gridpolish_from_2mm",
                      "figures", "single_rebar_recovered_bscan.png")
    add_figure(slide, obs, Inches(0.2), Inches(1.4), width=Inches(6.5))
    add_figure(slide, rec, Inches(6.8), Inches(1.4), width=Inches(6.5))
    add_caption(slide, Inches(0.2), Inches(5.9), Inches(6.5),
                "Observed B-scan (synthetic truth, 5 sources, 1.5 GHz Ricker).")
    add_caption(slide, Inches(6.8), Inches(5.9), Inches(6.5),
                "Recovered B-scan from inverted geometry (run 015).")
    add_text(slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.7),
             "Difference is below the floating-point floor: NRMS data = 0. The recovered model is "
             "not just a close fit; on the 1 mm grid it is the same model.",
             size=14, italic=True, color=CHARCOAL)
    set_notes(slide,
              "These are the two B-scans side by side: observed on the "
              "left, recovered from the inverted geometry on the right. "
              "Both show the same direct-wave band at the top, the same "
              "rebar diffraction hyperbola near 2 nanoseconds, and the "
              "same trailing reverberations. The two panels are not just "
              "visually similar — the trace-by-trace difference is below "
              "the floating-point floor, so the normalised data misfit is "
              "exactly zero. This is the strongest possible synthetic "
              "recovery signal: on the 1 mm grid, the recovered model is "
              "the same model as the truth, not just a close fit. The "
              "value of showing both panels is mostly to make the "
              "zero-misfit statement visible to the audience.")


def slide_noise(prs, fig_path):
    slide = new_slide(prs)
    add_title(slide,
              "Noise robustness",
              "Recovery stays correct under 1–10% additive observed-data noise")
    add_figure(slide, fig_path, Inches(0.3), Inches(1.4), width=Inches(13.0))
    add_text(slide, Inches(0.7), Inches(6.0), Inches(12), Inches(1.0),
             "Across two seeds and four noise levels, the polished radius is exactly 6.0 mm and "
             "the data NRMS matches the injected noise. Model NRMS stays at zero, meaning the "
             "rasterized geometry is unchanged.",
             size=14, color=CHARCOAL)
    set_notes(slide,
              "The point of this slide is not that the data misfit is "
              "small — under noise the misfit cannot be smaller than the "
              "noise floor. The point is that the recovered geometry does "
              "not move. Across two noise seeds and four noise levels — "
              "exact, one percent, five percent, and ten percent additive "
              "trace-RMS noise — the polished radius is exactly 6.0 mm in "
              "every case. The data-misfit bars on the right rise linearly "
              "with the injected noise level, exactly as they should: the "
              "residual is now fully explained by the noise we added, not "
              "by a bad model. The model-misfit row stays at zero "
              "throughout because the rasterised geometry the polish "
              "selects is unchanged. This is what makes the polish stage "
              "worth defending under realistic conditions.")


def slide_margin(prs, fig_path):
    slide = new_slide(prs)
    add_title(slide,
              "Margin at the 10% stress level",
              "Where the radius would start being ambiguous")
    add_figure(slide, fig_path, Inches(0.3), Inches(1.3), width=Inches(13.0))
    add_text(slide, Inches(0.7), Inches(6.1), Inches(12), Inches(0.9),
             "At 10% noise, r = 6.0 mm still wins, but by only ≈ 5.6e-4 over r = 6.2 mm. Beyond "
             "this point we should treat radius as a distribution rather than a single value, "
             "and use the top-k candidates list directly.",
             size=14, color=CHARCOAL)
    set_notes(slide,
              "This slide calibrates where the pipeline starts to lose "
              "radius discrimination. The eight bars are the top eight "
              "polish candidates at ten percent additive noise, sorted by "
              "their objective values. The true 6.0 mm radius still wins "
              "in two equivalent depth cells, but the margin against the "
              "next-distinct radius — 6.2 mm — is only about 5.6e-4. "
              "That is a real margin, but it is small enough that beyond "
              "this noise level we should not trust a single point "
              "estimate. The pipeline's recommended behaviour at this "
              "regime is to report the top-k candidates plus the distinct-"
              "radius margin and let the user decide how to interpret. It "
              "also tells us not to use the zero-misfit early-stop option "
              "on noisy data, because the true model cannot reach zero "
              "misfit when noise is present.")


def slide_landscape(prs):
    slide = new_slide(prs)
    add_title(slide,
              "Objective landscape diagnostics",
              "The optimizers are working on a well-behaved local landscape")

    xz_fig = os.path.join(OUTPUTS_DIR, "single_rebar_landscape_9src_run01",
                         "figures", "xz_landscape.png")
    zr_fig = os.path.join(OUTPUTS_DIR, "single_rebar_landscape_9src_run01",
                         "figures", "z_radius_landscape.png")
    add_figure(slide, xz_fig, Inches(0.3), Inches(1.4), width=Inches(6.5))
    add_figure(slide, zr_fig, Inches(6.8), Inches(1.4), width=Inches(6.5))
    add_caption(slide, Inches(0.3), Inches(5.6), Inches(6.5),
                "x–z slice at the true radius: unimodal, true position at the dark minimum.")
    add_caption(slide, Inches(6.8), Inches(5.6), Inches(6.5),
                "z–radius slice at the true x: the coupled valley discussed in the previous slide.")
    add_text(slide, Inches(0.7), Inches(6.2), Inches(12), Inches(1.0),
             "The local landscape has the expected shape: position is convex around truth, "
             "but the z–radius slice shows the elongated valley that causes continuous optimizers "
             "to stop slightly off-radius. That valley is exactly what the grid polish steps over.",
             size=14, color=CHARCOAL)
    set_notes(slide,
              "These two panels are direct diagnostics of the inversion "
              "objective landscape — not optimisation runs, just the "
              "objective evaluated on a grid of candidates. The left panel "
              "fixes the radius at the true value and sweeps x and z; the "
              "dark minimum sits exactly at the true position and the "
              "landscape is smoothly convex around it. The right panel "
              "fixes x at the true value and sweeps z and radius; here we "
              "see the elongated diagonal valley where deeper-and-larger "
              "and shallower-and-smaller candidates all fit the data "
              "almost equally well. That valley is the geometric origin of "
              "the radius bias we have been discussing. The polish stage "
              "works because it enumerates discrete cells across this "
              "valley rather than following a gradient down it.")


def slide_next(prs):
    slide = new_slide(prs)
    add_title(slide,
              "What comes next",
              "Lift the twin-experiment assumptions one at a time")

    cards = [
        ("Multi-rebar",
         "Re-introduce 2-3 rebars with shared global x search and per-rebar polish. The polish stage and identifiability hierarchy stay; the global stage gets harder."),
        ("Nuisance parameters",
         "Let concrete εr, conductivity, time-zero, and source wavelet float. This is the Jazayeri-style realism gap — radius identifiability depends on calibrating them."),
        ("Field-style validation",
         "Add a hyperbola/migration baseline (RGPR-style) to quantify what FWI buys over geometry-only methods, and cross-check the forward against gprMax."),
    ]
    card_w = Inches(4.0)
    card_h = Inches(3.7)
    for i, (title, body) in enumerate(cards):
        left = Inches(0.7 + i * 4.2)
        top = Inches(1.4)
        add_panel(slide, left, top, card_w, card_h, LIGHT_BG, STEEL)
        add_text(slide, left, top + Inches(0.2), card_w, Inches(0.55),
                 title, size=18, bold=True, color=NAVY,
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(0.85), card_w, Inches(2.7),
                 body, size=13, color=CHARCOAL, alignment=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5),
             "Each direction is supported by the existing pipeline: the forward model, "
             "numbered-run convention, and grid-polish stage transfer directly. The work is "
             "in the search strategy and in mapping which parameters the data actually supports.",
             size=14, color=CHARCOAL)
    set_notes(slide,
              "The three forward directions all lift one assumption of the "
              "current twin experiment. Multi-rebar re-introduces two or "
              "three rebars; the polish stage and identifiability hierarchy "
              "carry over, but the global search stage gets harder because "
              "the basin landscape becomes multi-modal. Nuisance parameters "
              "means letting concrete permittivity, conductivity, time-zero, "
              "and source wavelet float — this is the realism gap that "
              "shows up in the FWI literature, and radius identifiability "
              "depends on calibrating it. Field-style validation adds a "
              "hyperbola-fitting and migration baseline so we can quantify "
              "what full-waveform inversion buys over simpler methods, and "
              "lets us cross-check the forward model against an "
              "independent solver. The pipeline pieces that transfer "
              "directly are the forward model, the numbered-run convention, "
              "and the grid-polish stage. The new work is in the search "
              "strategy and in mapping which parameters the observed data "
              "actually constrains.")


# -----------------------------------------------------------------------------
# Driver
# -----------------------------------------------------------------------------


def build():
    print("Generating custom figures...")
    pipeline_fig = fig_pipeline_architecture()
    sweep_pair = fig_radius_sweep_pair()
    noise_fig = fig_noise_robustness()
    margin_fig = fig_top_k_margin()

    print("Assembling deck...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_problem(prs)
    slide_pipeline_arch(prs, pipeline_fig)
    slide_forward(prs)
    slide_radius_hard(prs, sweep_pair)
    slide_coupling(prs)
    slide_grid_polish(prs)
    slide_final_model(prs)
    slide_bscan_compare(prs)
    slide_landscape(prs)
    slide_noise(prs, noise_fig)
    slide_margin(prs, margin_fig)
    slide_next(prs)

    prs.save(PPTX_PATH)
    print(f"\nSaved: {PPTX_PATH}")


if __name__ == "__main__":
    build()
